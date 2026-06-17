"""Phase1 最小実証: 複数KW → 各TikTok1位サムネ → FMTマトリクス(slide36)の複数空枠に並べて注入。

MVP `demo_thumb_to_fmt_mvp.py`(1KW1枚) の拡張版。本番 src 無改修・無AWS・全ローカル。
「DR→出てきたKW群→TikTok並列検索→実物サムネをマトリクスに並べる」を物証(.pptx)で示す。
画像注入ロジックは MVP のローカル関数を再利用（本番 renderer/contract は未改修＝Phase2/3）。

KWの渡し方:
  (a) 直接   : python scripts/demo_campaign_thumbs_mvp.py 集中 作業用BGM 勉強 --template ... --out ...
  (b) DR JSON: --from-dr <GeminiDRJSON.json>  （D_publicity[].trend_word / E_community[].tiktok_tags / C_tiktok[].tag を抽出）

前提 env: SSL_CERT_FILE=~/.hermes/ca_bundle.pem / node+Chromium / ffmpeg
"""

from __future__ import annotations

import argparse
import json
import sys
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT / "src"))
sys.path.insert(0, str(_ROOT / "scripts"))  # 既存MVPのヘルパを再利用

from demo_thumb_to_fmt_mvp import (  # noqa: E402
    _SLOT_SIZE_EMU,
    _SLOT_TOP_EMU,
    _add_picture_fit,
    _dummy_composer_output,
    _iter_shapes,
    _normalize_to_jpeg,
)
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

from teamagent.skills.proposal_deck.renderer import render_deck  # noqa: E402

_MATRIX_SLIDE_INDEX = 36
_MAX_KW_DEFAULT = 6  # slide36 マトリクス最下段の空枠数
_COORD_TOL_EMU = 80000


def extract_keywords_from_dr(dr: dict) -> list[str]:
    """GeminiDRJSON から TikTok検索KWを抽出（trend_word / tiktok_tags / tag）。順序保持で重複排除。

    02_fmt_placeholder_map.md のマッピング: D_publicity[].trend_word→{58-71}、
    E_community[].tiktok_tags→{72-92}、C_tiktok[].tag→界隈言語の根拠。
    """
    kws: list[str] = []
    for item in dr.get("D_publicity", []) or []:
        w = (item or {}).get("trend_word")
        if w:
            kws.append(str(w))
    for item in dr.get("E_community", []) or []:
        for t in (item or {}).get("tiktok_tags", []) or []:
            if t:
                kws.append(str(t).lstrip("#"))
    for item in dr.get("C_tiktok", []) or []:
        t = (item or {}).get("tag")
        if t:
            kws.append(str(t).lstrip("#"))
    seen: set[str] = set()
    uniq: list[str] = []
    for k in kws:
        k = k.strip()
        if k and k not in seen:
            seen.add(k)
            uniq.append(k)
    return uniq


def _thumb_for_kw(kw: str, fallback_image: bytes | None) -> tuple[str, bytes | None, str]:
    """1KW → (kw, jpeg_bytes|None, source_label)。失敗時 fallback_image があればそれを使う。"""
    from teamagent.adapters.tiktok_scraper import search_tiktok
    from teamagent.skills.video_algorithm.thumbnails import fetch_cover

    try:
        result = search_tiktok(kw, max_videos=1, request_id=f"mvp-{kw[:8]}")
        if result.videos:
            top = result.videos[0]  # rank 1
            raw = fetch_cover(top.cover_url, request_id="mvp")
            if raw:
                return kw, _normalize_to_jpeg(raw), f"1位: {top.url}"
    except Exception as e:
        # TikTokScrapeError(captcha/0件/timeout) 等は握りつぶし、フォールバックへ
        sys.stderr.write(f"   [warn] '{kw}' 検索失敗: {type(e).__name__}\n")
    if fallback_image is not None:
        return kw, _normalize_to_jpeg(fallback_image), "fallback(sample)"
    return kw, None, "取得失敗"


def _find_thumb_slots(slide) -> list:
    """マトリクス最下段の空枠（0.67"角・y≈5.98"）を全て・左→右の順で返す。"""
    slots = []
    for shape in _iter_shapes(slide.shapes):
        try:
            w, t = int(shape.width), int(shape.top)
        except (TypeError, ValueError):
            continue
        if abs(w - _SLOT_SIZE_EMU) <= _COORD_TOL_EMU and abs(t - _SLOT_TOP_EMU) <= _COORD_TOL_EMU:
            slots.append(shape)
    return sorted(slots, key=lambda s: int(s.left))


def main() -> int:
    ap = argparse.ArgumentParser(description="Phase1: 複数KW→各1位サムネ→FMTマトリクスに並べる")
    ap.add_argument("kws", nargs="*", help="検索KW（複数可）。--from-dr 指定時は無視")
    ap.add_argument("--template", required=True, help="template_v2.pptx の絶対パス")
    ap.add_argument("--out", default="/tmp/mvp_campaign_thumbs.pptx")
    ap.add_argument("--from-dr", default=None, help="GeminiDRJSON(.json) から KW を抽出")
    ap.add_argument("--sample-image", default=None, help="全KWのフォールバック画像")
    ap.add_argument("--max-kw", type=int, default=_MAX_KW_DEFAULT)
    ap.add_argument("--slide-index", type=int, default=_MATRIX_SLIDE_INDEX)
    args = ap.parse_args()

    template = Path(args.template)
    if not template.exists():
        raise SystemExit(f"テンプレが見つかりません: {template}")

    # KW解決（DR JSON抽出 or 直接指定）
    if args.from_dr:
        dr = json.loads(Path(args.from_dr).read_text())
        kws = extract_keywords_from_dr(dr)
        print(f"📥 DR JSONから抽出したKW({len(kws)}): {kws}")
    else:
        kws = list(args.kws)
    if not kws:
        raise SystemExit("KWがありません。位置引数で渡すか --from-dr を指定してください。")
    kws = kws[: args.max_kw]

    fallback = Path(args.sample_image).read_bytes() if args.sample_image else None

    # 並列でサムネ取得（KW毎 search_tiktok→fetch_cover→JPEG正規化・graceful）
    print(f"🔎 {len(kws)}KW を並列検索（node+Chromium・max_workers=3）…")
    with ThreadPoolExecutor(max_workers=3) as ex:
        results = list(ex.map(lambda k: _thumb_for_kw(k, fallback), kws))
    got = [(k, b, s) for (k, b, s) in results if b is not None]
    for k, _b, s in results:
        print(f"   - {k}: {s}")
    if not got:
        raise SystemExit("全KWでサムネ取得に失敗。--sample-image でフォールバックしてください。")

    # FMTテキスト生成（ダミー95枠・本番 render_deck 無改修）
    print("   ⏳ FMTテキスト流し込み（render_deck・89MBテンプレ）…")
    rendered = render_deck(
        _dummy_composer_output(), template, Path(args.out), fail_if_missing=False
    )

    # マトリクス空枠に左→右で並べて注入
    prs = Presentation(str(rendered))
    slide = prs.slides[args.slide_index]
    slots = _find_thumb_slots(slide)
    if not slots:
        raise SystemExit("マトリクス空枠が見つかりません。slide_index/テンプレを確認。")
    n = min(len(got), len(slots))
    for i in range(n):
        _kw, img, _s = got[i]
        slot = slots[i]
        _add_picture_fit(slide, img, int(slot.left), int(slot.top), int(slot.height))
    prs.save(str(args.out))

    # 自己検証: 対象スライドの PICTURE 実数（テンプレのマトリクス枠は元0枚）
    prs2 = Presentation(str(args.out))
    pics = [
        sh
        for sh in _iter_shapes(prs2.slides[args.slide_index].shapes)
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    ok = len(pics) >= n
    print("\n" + ("✅ 完了・検証OK" if ok else "⚠️ 検証で枚数不一致"))
    print(f"   出力        : {args.out}")
    print(f"   注入        : {n}枚 / 空枠{len(slots)} / KW{len(kws)}（取得成功{len(got)}）")
    print(f"   PICTURE実数 : slide{args.slide_index} に {len(pics)} 枚")
    print(f"   確認        : open '{args.out}'")
    return 0 if ok else 2


if __name__ == "__main__":
    raise SystemExit(main())
