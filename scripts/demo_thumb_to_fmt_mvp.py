"""最小実証: KW→TikTok検索→1位→実物サムネを最終FMT(template_v2.pptx)の1枠に「画像として」貼る。

  feeder(検索) → 真実源(ComposerOutput / テンプレ) → renderer(pptx) + 画像注入

の「画像経路」を、ローカル完結・本番 src 無改修で1本通すデモ。本番 AWS には一切触れない。
本番の renderer.py / contract.py は変更せず、画像注入ロジックはこのデモ内のローカル関数に閉じる
（95枠 validator / 段落跨ぎ置換 / 後方互換への波及を避けるため）。正式実装は次フェーズ。

前提 env:
  - 会社プロキシ: SSL_CERT_FILE=~/.hermes/ca_bundle.pem  （cover_url 取得に必須）
  - tiktok検索:  node + Chromium（tools/tiktok_scraper/ は npm install 済）
  - ffmpeg（PATH）: cover を JPEG/縮小に正規化（webp 等でも安全に）
  - GEMINI_API_KEY は不要（軽量経路＝検索メタ+サムネのみ。動画DL/分析はしない）

使い方:
  export SSL_CERT_FILE=~/.hermes/ca_bundle.pem
  python scripts/demo_thumb_to_fmt_mvp.py 集中 \
      --template "/Users/s-komata/Claude 社内横断/teamagent_consulting/teamagent/assets/template_v2.pptx" \
      --out /tmp/mvp_fmt_with_thumb.pptx

TikTok 不安定時のフォールバック（画像経路だけ先に検証）:
  ... --cover-url https://....jpg
  ... --sample-image /tmp/sample_thumb.jpg
"""

from __future__ import annotations

import argparse
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT / "src"))

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.util import Emu  # noqa: E402

from teamagent.skills.proposal_deck.contract import (  # noqa: E402
    LENGTH_RULES,
    VALID_IDS,
    ComposerOutput,
)
from teamagent.skills.proposal_deck.renderer import render_deck  # noqa: E402

# template_v2.pptx の {58-92} マトリクス・スライド（0始まり）と、最下段の空サムネ枠（EMU）。
# 914400 EMU = 1 inch。空枠は 0.67"角（=609905 EMU）・y=5.98"（=5454720 EMU）・最左 x=1.0"（=914400 EMU）。
_MATRIX_SLIDE_INDEX = 36
_SLOT_TOP_EMU = 5454720
_SLOT_SIZE_EMU = 609905
_SLOT_LEFT_EMU = 914400
_COORD_TOL_EMU = 80000  # 座標一致の許容誤差


def _dummy_composer_output() -> ComposerOutput:
    """全 95 枠を字数規則どおり埋めたダミー（テキストは検証対象外）。

    tests/skills/proposal_deck/test_proposal_deck_skill.py の _full_composer_json と同じ埋め方。
    """
    placeholders: dict[int, str] = {}
    for pid in sorted(VALID_IDS):
        if pid in LENGTH_RULES:
            lo, hi = LENGTH_RULES[pid]
            placeholders[pid] = "サ" * ((lo + hi) // 2)
        else:
            placeholders[pid] = f"値-{pid}"
    return ComposerOutput(placeholders=placeholders)


def _normalize_to_jpeg(data: bytes) -> bytes:
    """画像 bytes を JPEG に正規化（webp/png でも安全に・幅480pxへ縮小）。

    python-pptx は webp を読めないため、ffmpeg で必ず JPEG 化する。ffmpeg 不在なら素通し。
    """
    if not data or not shutil.which("ffmpeg"):
        return data
    with tempfile.TemporaryDirectory() as tmp:
        src = os.path.join(tmp, "src")
        dst = os.path.join(tmp, "out.jpg")
        with open(src, "wb") as f:
            f.write(data)
        try:
            subprocess.run(
                ["ffmpeg", "-y", "-i", src, "-vframes", "1", "-vf", "scale=480:-1", dst],
                capture_output=True,
                timeout=30,
                check=False,
            )
            if os.path.exists(dst) and os.path.getsize(dst) > 0:
                with open(dst, "rb") as f:
                    return f.read()
        except Exception as e:
            print(f"   ⚠️ ffmpeg 正規化に失敗（{type(e).__name__}）→ 生バイトのまま試行")
    return data


def _get_rank1_cover(kw: str) -> tuple[bytes, str]:
    """KW で TikTok 検索 → 1位の (サムネ生バイト, 動画URL) を返す（軽量・Gemini不要）。"""
    from teamagent.adapters.tiktok_scraper import search_tiktok
    from teamagent.skills.video_algorithm.thumbnails import fetch_cover

    result = search_tiktok(kw, max_videos=1, request_id="mvp")
    if not result.videos:
        raise SystemExit(
            "TikTok 検索が 0 件。--cover-url か --sample-image でフォールバックしてください。"
        )
    top = result.videos[0]  # rank 1（上位から）
    data = fetch_cover(top.cover_url, request_id="mvp")
    if not data:
        raise SystemExit(
            f"サムネ取得に失敗（cover_url={top.cover_url[:60]}…）。SSL_CERT_FILE 未設定の可能性。"
            " --cover-url か --sample-image でフォールバックしてください。"
        )
    return data, top.url


def _iter_shapes(shapes):
    """グループ内も再帰して全 shape を yield（slides[36] の SEE 枠は通常トップレベル）。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _find_thumb_slot(slide):
    """最下段の空サムネ枠（0.67"角・y≈5.98"）の最左を座標で探す。無ければ None。"""
    candidates = []
    for shape in _iter_shapes(slide.shapes):
        try:
            w = int(shape.width)
            t = int(shape.top)
        except (TypeError, ValueError):
            continue
        if abs(w - _SLOT_SIZE_EMU) <= _COORD_TOL_EMU and abs(t - _SLOT_TOP_EMU) <= _COORD_TOL_EMU:
            candidates.append(shape)
    if not candidates:
        return None
    return min(candidates, key=lambda s: int(s.left))


def _add_picture_fit(slide, img_bytes: bytes, left: int, top: int, box_h: int):
    """枠の高さに等比で合わせて画像を貼る（歪み無し）。"""
    import io

    pic = slide.shapes.add_picture(io.BytesIO(img_bytes), Emu(int(left)), Emu(int(top)))
    if pic.width and pic.height:
        scale = box_h / pic.height
        pic.height = Emu(int(box_h))
        pic.width = Emu(int(pic.width * scale))
        pic.left = Emu(int(left))
        pic.top = Emu(int(top))
    return pic


def _verify_picture(out_path: Path, slide_index: int) -> bool:
    """保存後の pptx を開き直し、対象スライドに目標座標付近の PICTURE があるか検証。"""
    prs = Presentation(str(out_path))
    slide = prs.slides[slide_index]
    for shape in _iter_shapes(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        try:
            if abs(int(shape.top) - _SLOT_TOP_EMU) <= _COORD_TOL_EMU:
                return True
        except (TypeError, ValueError):
            continue
    return False


def main() -> int:
    ap = argparse.ArgumentParser(description="最小実証: KW→1位サムネ→FMT1枠に画像注入")
    ap.add_argument("kw", nargs="?", default="集中", help="検索 KW（例: 集中）")
    ap.add_argument("--template", required=True, help="template_v2.pptx の絶対パス")
    ap.add_argument("--out", default="/tmp/mvp_fmt_with_thumb.pptx", help="出力 pptx パス")
    ap.add_argument("--cover-url", default=None, help="フォールバック: この cover_url から取得")
    ap.add_argument("--sample-image", default=None, help="フォールバック: ローカル画像ファイル")
    ap.add_argument(
        "--slide-index", type=int, default=_MATRIX_SLIDE_INDEX, help="貼り先スライド（0始まり）"
    )
    args = ap.parse_args()

    template = Path(args.template)
    if not template.exists():
        raise SystemExit(f"テンプレが見つかりません: {template}")
    out = Path(args.out)

    # 1) 画像バイト取得（実物 1位 or フォールバック）→ JPEG 正規化
    print(f"🔎 最小実証: KW='{args.kw}' → 1位サムネ → FMT[{args.slide_index}] に画像注入")
    if args.sample_image:
        raw = Path(args.sample_image).read_bytes()
        source_label = f"sample-image={args.sample_image}"
    elif args.cover_url:
        from teamagent.skills.video_algorithm.thumbnails import fetch_cover

        raw = fetch_cover(args.cover_url, request_id="mvp")
        if not raw:
            raise SystemExit("--cover-url からの取得に失敗（SSL_CERT_FILE を確認）。")
        source_label = f"cover-url={args.cover_url[:60]}"
    else:
        print("   ⏳ TikTok 検索（node+Chromium・~10-60s）…")
        raw, video_url = _get_rank1_cover(args.kw)
        source_label = f"TikTok 1位: {video_url}"
    img_bytes = _normalize_to_jpeg(raw)
    print(f"   画像ソース: {source_label}（{len(img_bytes)} bytes・JPEG正規化済）")

    # 2) 全95枠ダミー → 既存 render_deck でテキスト pptx 生成（本番 src 無改修）
    print("   ⏳ FMT テキスト流し込み（render_deck・89MB テンプレ）…")
    composer_out = _dummy_composer_output()
    rendered = render_deck(composer_out, template, out, fail_if_missing=False)

    # 3) 生成 pptx を開き直し、slides[36] の空枠へ add_picture
    prs = Presentation(str(rendered))
    if args.slide_index >= len(prs.slides):
        raise SystemExit(f"slide_index={args.slide_index} は範囲外（総 {len(prs.slides)} 枚）")
    slide = prs.slides[args.slide_index]
    slot = _find_thumb_slot(slide)
    if slot is None:
        print("   ⚠️ 空枠を座標で発見できず → 既知座標に直接配置")
        _add_picture_fit(slide, img_bytes, _SLOT_LEFT_EMU, _SLOT_TOP_EMU, _SLOT_SIZE_EMU)
    else:
        print(f"   空枠を発見: left={int(slot.left)} top={int(slot.top)} EMU")
        _add_picture_fit(slide, img_bytes, int(slot.left), int(slot.top), int(slot.height))
    prs.save(str(out))

    # 4) 自己検証: 対象スライドに PICTURE が入ったか
    ok = _verify_picture(out, args.slide_index)
    print("\n" + ("✅ 完了・検証 OK" if ok else "⚠️ 完了したが検証で PICTURE を確認できず"))
    print(f"   出力      : {out}")
    print(
        f"   貼り先    : slide_index={args.slide_index}（{args.slide_index + 1}枚目）最下段の空枠"
    )
    print(f"   画像ソース: {source_label}")
    print(f"   確認      : open '{out}'")
    return 0 if ok else 2


if __name__ == "__main__":
    raise SystemExit(main())
