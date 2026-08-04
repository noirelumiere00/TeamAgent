"""Isolated, network-denied renderer entrypoint for the media worker."""

from __future__ import annotations

import io
import json
import os
import re
import sys
from datetime import date, timedelta
from pathlib import Path
from typing import Any

from teamagent.media.operations import (
    _EXTERNAL_HTML_REF,
    _PLACEHOLDER,
    MediaOperationError,
    _iter_proposal_image_slots,
    _iter_text_frames,
    _replace_placeholders,
)

_AUXILIARY = re.compile(r"\{\{(PB-[A-Z0-9_-]{1,60})\}\}")
_DATE = re.compile(r"\{\{PB-DATE:([+-]?\d{1,3}):(%Y/%m/%d|%m/%d|%Y年%m月%d日)\}\}")
_TEMPLATE_VERSION = re.compile(r"\{\{PB-TEMPLATE:([a-z0-9-]{1,40})\}\}")
_PB_TOKEN = re.compile(r"\{\{PB-[^{}]+\}\}")
_LEGACY_INSTRUCTION = re.compile(r"自動入力|貼り付けてください|はめ込|転記|差し替え")
_BRACE_CHARACTER = re.compile(r"[{}｛｝]")
_PROPOSAL_BUILDER_PROFILE = "proposal-builder-v1"
_PROPOSAL_BUILDER_REQUIRED_AUXILIARY = frozenset(
    {
        "PB-ACCOUNTS",
        "PB-CASES",
        "PB-CLIENT-NAME",
        "PB-DATETIME",
        "PB-EXPERIENCE",
        "PB-KEY-MESSAGE",
        "PB-MONTH",
        "PB-PRODUCT-NAME",
    }
)
_PROPOSAL_BUILDER_REQUIRED_DATE_OFFSETS = frozenset(range(-56, 22, 7))
_PROPOSAL_BUILDER_EXPECTED_SLIDE_COUNT = 83


def _exact(value: Any, keys: set[str], name: str) -> dict[str, Any]:
    if not isinstance(value, dict) or set(value) != keys:
        raise MediaOperationError("MEDIA_RENDER_MANIFEST_INVALID", f"{name} is invalid")
    return value


def _path(root: Path, value: Any, *, output: bool = False) -> Path:
    if not isinstance(value, str) or not value or Path(value).is_absolute():
        raise MediaOperationError("MEDIA_RENDER_PATH_INVALID", "renderer path is invalid")
    path = (root / value).resolve()
    if not path.is_relative_to(root):
        raise MediaOperationError("MEDIA_RENDER_PATH_INVALID", "renderer path escaped workdir")
    if not output and (path.is_symlink() or not path.is_file()):
        raise MediaOperationError("MEDIA_RENDER_PATH_INVALID", "renderer input is invalid")
    if output and path.exists():
        raise MediaOperationError("MEDIA_RENDER_PATH_INVALID", "renderer output already exists")
    return path


def _slides(root: Path, manifest: dict[str, Any]) -> dict[str, Any]:
    spec = _exact(
        manifest,
        {"kind", "html", "output", "selector", "width", "height", "scale"},
        "slides renderer manifest",
    )
    html_path = _path(root, spec["html"])
    destination = _path(root, spec["output"], output=True)
    html = html_path.read_text(encoding="utf-8")
    if _EXTERNAL_HTML_REF.search(html):
        raise MediaOperationError(
            "MEDIA_HTML_NETWORK_REFERENCE",
            "slides HTML references network",
        )
    from playwright.sync_api import sync_playwright
    from pptx import Presentation
    from pptx.util import Emu, Inches

    chromium = os.environ.get("CHROMIUM_PATH", "/usr/lib/chromium/chromium")
    images: list[bytes] = []
    with sync_playwright() as playwright:
        browser = playwright.chromium.launch(
            executable_path=chromium,
            headless=True,
            args=[
                "--disable-gpu",
                "--disable-dev-shm-usage",
                "--disable-quic",
                "--disable-setuid-sandbox",
                "--host-resolver-rules=MAP * ~NOTFOUND",
                "--no-proxy-server",
            ],
            chromium_sandbox=True,
        )
        try:
            page = browser.new_page(
                viewport={"width": int(spec["width"]), "height": int(spec["height"])},
                device_scale_factor=int(spec["scale"]),
            )
            page.route("**/*", lambda route: route.abort())
            page.set_content(html, wait_until="domcontentloaded")
            slides = page.locator(str(spec["selector"]))
            count = slides.count()
            if count < 1 or count > 20:
                raise MediaOperationError(
                    "MEDIA_SLIDE_COUNT_INVALID",
                    "slide count is out of range",
                )
            for index in range(count):
                images.append(slides.nth(index).screenshot(type="png"))
        finally:
            browser.close()

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for image in images:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(image),
            Emu(0),
            Emu(0),
            width=presentation.slide_width,
            height=presentation.slide_height,
        )
    presentation.save(str(destination))
    return {"slides": len(images), "network_requests_allowed": 0}


def _replace_proposal_special_tokens(
    text_frame: Any,
    auxiliary: dict[str, str],
    posting_start_date: date | None,
) -> None:
    paragraphs = list(text_frame.paragraphs)
    combined = "".join(paragraph.text for paragraph in paragraphs)
    if not combined:
        return
    replaced = _AUXILIARY.sub(
        lambda match: auxiliary.get(match.group(1), match.group(0)),
        combined,
    )
    if posting_start_date is not None:
        replaced = _DATE.sub(
            lambda match: (posting_start_date + timedelta(days=int(match.group(1)))).strftime(
                match.group(2)
            ),
            replaced,
        )
    replaced = _TEMPLATE_VERSION.sub(
        lambda match: "" if match.group(1) == _PROPOSAL_BUILDER_PROFILE else match.group(0),
        replaced,
    )
    if replaced == combined:
        return
    paragraphs[0].text = replaced
    for paragraph in paragraphs[1:]:
        paragraph.text = ""


def _legacy_artifacts(text: str, *, allow_template_tokens: bool) -> list[str]:
    """Return obsolete brace tokens/manual instructions at a render boundary."""

    scrubbed = text
    if allow_template_tokens:
        for pattern in (_PLACEHOLDER, _DATE, _TEMPLATE_VERSION, _AUXILIARY):
            scrubbed = pattern.sub("", scrubbed)
    findings: list[str] = []
    if _BRACE_CHARACTER.search(scrubbed):
        findings.append("brace")
    findings.extend(match.group(0) for match in _LEGACY_INSTRUCTION.finditer(scrubbed))
    return findings


def _proposal(root: Path, manifest: dict[str, Any]) -> dict[str, Any]:
    spec = _exact(
        manifest,
        {"kind", "template", "composer", "evidence", "output", "fail_if_missing"},
        "proposal renderer manifest",
    )
    template = _path(root, spec["template"])
    composer = _path(root, spec["composer"])
    destination = _path(root, spec["output"], output=True)
    evidence = spec["evidence"]
    if not isinstance(evidence, list) or len(evidence) > 20:
        raise MediaOperationError("MEDIA_RENDER_MANIFEST_INVALID", "evidence is invalid")
    try:
        raw = json.loads(composer.read_text(encoding="utf-8"))
        placeholders = {int(key): str(value) for key, value in raw["placeholders"].items()}
        skipped = {int(item["id"]) for item in raw.get("skipped_placeholders", [])}
        auxiliary = {
            str(key): str(value) for key, value in raw.get("auxiliary_placeholders", {}).items()
        }
        profile = str(raw.get("template_profile", "base"))
        raw_start_date = raw.get("posting_start_date")
        posting_start_date = (
            date.fromisoformat(raw_start_date) if isinstance(raw_start_date, str) else None
        )
    except (KeyError, TypeError, ValueError, json.JSONDecodeError) as exc:
        raise MediaOperationError("MEDIA_COMPOSER_INVALID", "composer JSON is invalid") from exc
    valid_ids = set(range(1, 104)) - set(range(48, 56))
    if set(placeholders) - valid_ids or skipped - valid_ids:
        raise MediaOperationError("MEDIA_COMPOSER_IDS_INVALID", "composer IDs are invalid")
    for placeholder_id in valid_ids:
        placeholders.setdefault(placeholder_id, "要確認（データ未検出）")

    from pptx import Presentation
    from pptx.util import Emu

    presentation = Presentation(str(template))
    template_ids: set[int] = set()
    template_auxiliary: set[str] = set()
    template_date_offsets: set[int] = set()
    template_versions: set[str] = set()
    template_legacy_artifacts: list[str] = []
    for text_frame in _iter_text_frames(presentation):
        combined = "".join(paragraph.text for paragraph in text_frame.paragraphs)
        template_ids.update(int(match.group(1)) for match in _PLACEHOLDER.finditer(combined))
        template_auxiliary.update(match.group(1) for match in _AUXILIARY.finditer(combined))
        template_date_offsets.update(int(match.group(1)) for match in _DATE.finditer(combined))
        template_versions.update(match.group(1) for match in _TEMPLATE_VERSION.finditer(combined))
        template_legacy_artifacts.extend(_legacy_artifacts(combined, allow_template_tokens=True))
    if profile == _PROPOSAL_BUILDER_PROFILE:
        invalid_numeric_inventory = template_ids != valid_ids
        invalid_values = set(auxiliary) != _PROPOSAL_BUILDER_REQUIRED_AUXILIARY
        invalid_tokens = template_auxiliary != _PROPOSAL_BUILDER_REQUIRED_AUXILIARY
        if (
            invalid_numeric_inventory
            or invalid_values
            or invalid_tokens
            or template_legacy_artifacts
        ):
            raise MediaOperationError(
                "MEDIA_PPTX_TEMPLATE_PROFILE_INVALID",
                "integrated proposal template inventory is incomplete",
            )
        if (
            posting_start_date is None
            or len(presentation.slides) != _PROPOSAL_BUILDER_EXPECTED_SLIDE_COUNT
            or template_versions != {_PROPOSAL_BUILDER_PROFILE}
            or not _PROPOSAL_BUILDER_REQUIRED_DATE_OFFSETS.issubset(template_date_offsets)
        ):
            raise MediaOperationError(
                "MEDIA_PPTX_TEMPLATE_PROFILE_INVALID",
                "integrated proposal template schedule markers are incomplete",
            )
    elif profile != "base":
        raise MediaOperationError(
            "MEDIA_COMPOSER_INVALID",
            "composer template profile is invalid",
        )

    for text_frame in _iter_text_frames(presentation):
        _replace_placeholders(text_frame, placeholders)
        _replace_proposal_special_tokens(text_frame, auxiliary, posting_start_date)
    remaining: list[int] = []
    remaining_special: list[str] = []
    remaining_legacy: list[str] = []
    for text_frame in _iter_text_frames(presentation):
        combined = "".join(paragraph.text for paragraph in text_frame.paragraphs)
        remaining.extend(int(match.group(1)) for match in _PLACEHOLDER.finditer(combined))
        remaining_special.extend(match.group(0) for match in _PB_TOKEN.finditer(combined))
        if profile == _PROPOSAL_BUILDER_PROFILE:
            remaining_legacy.extend(_legacy_artifacts(combined, allow_template_tokens=False))
    if (remaining or remaining_special or remaining_legacy) and spec["fail_if_missing"] is True:
        raise MediaOperationError("MEDIA_PPTX_UNFILLED", "unfilled placeholders remain")

    slots = _iter_proposal_image_slots(presentation)
    injected = 0
    ordered: list[tuple[int, int, Path]] = []
    for value in evidence:
        entry = _exact(value, {"placeholder_id", "rank", "path"}, "evidence entry")
        ordered.append(
            (
                int(entry["placeholder_id"]),
                int(entry["rank"]),
                _path(root, entry["path"]),
            )
        )
    for _placeholder_id, _rank, image in sorted(ordered):
        try:
            slide, shape = next(slots)
        except StopIteration:
            break
        try:
            picture = slide.shapes.add_picture(
                io.BytesIO(image.read_bytes()),
                Emu(int(shape.left)),
                Emu(int(shape.top)),
            )
        except Exception:
            continue
        if picture.width and picture.height:
            scale = int(shape.height) / picture.height
            picture.height = Emu(int(shape.height))
            picture.width = Emu(int(picture.width * scale))
        injected += 1
    presentation.save(str(destination))
    return {
        "filled": len(placeholders),
        "skipped": len(skipped),
        "evidence_images": injected,
        "remaining_placeholders": len(remaining),
        "remaining_special_placeholders": len(remaining_special),
        "remaining_legacy_artifacts": len(remaining_legacy),
        "template_numeric_ids": len(template_ids),
        "template_slides": len(presentation.slides),
    }


def _pdf(root: Path, manifest: dict[str, Any]) -> dict[str, Any]:
    spec = _exact(manifest, {"kind", "html", "output"}, "PDF renderer manifest")
    html_path = _path(root, spec["html"])
    destination = _path(root, spec["output"], output=True)
    html = html_path.read_text(encoding="utf-8")
    if _EXTERNAL_HTML_REF.search(html):
        raise MediaOperationError(
            "MEDIA_HTML_NETWORK_REFERENCE",
            "PDF HTML references network",
        )
    from weasyprint import HTML

    def block_url_fetcher(_url: str, *_args: Any, **_kwargs: Any) -> dict[str, Any]:
        raise MediaOperationError(
            "MEDIA_HTML_NETWORK_REFERENCE",
            "PDF rendering network access is blocked",
        )

    HTML(string=html, url_fetcher=block_url_fetcher).write_pdf(str(destination))
    if not destination.is_file() or destination.stat().st_size < 5:
        raise MediaOperationError("MEDIA_PDF_EMPTY", "weasyprint produced an empty PDF")
    return {"network_requests_allowed": 0}


def main() -> int:
    try:
        if len(sys.argv) != 2:
            raise MediaOperationError(
                "MEDIA_RENDER_MANIFEST_INVALID",
                "renderer manifest argument is required",
            )
        root = Path.cwd().resolve()
        manifest_path = _path(root, sys.argv[1])
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        if not isinstance(manifest, dict):
            raise MediaOperationError(
                "MEDIA_RENDER_MANIFEST_INVALID",
                "renderer manifest is invalid",
            )
        kind = manifest.get("kind")
        if kind == "slides":
            metadata = _slides(root, manifest)
        elif kind == "proposal_pptx":
            metadata = _proposal(root, manifest)
        elif kind == "pdf":
            metadata = _pdf(root, manifest)
        else:
            raise MediaOperationError(
                "MEDIA_RENDER_MANIFEST_INVALID",
                "renderer kind is invalid",
            )
        print(json.dumps({"ok": True, "metadata": metadata}, sort_keys=True))
        return 0
    except MediaOperationError as exc:
        print(json.dumps({"ok": False, "code": exc.code}, sort_keys=True))
        return 2
    except Exception:
        print(json.dumps({"ok": False, "code": "MEDIA_RENDER_FAILED"}, sort_keys=True))
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
