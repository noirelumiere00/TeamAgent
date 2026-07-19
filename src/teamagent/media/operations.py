"""media worker 内で実行する最小 transform/render/download 実装。"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import shutil
import signal
import subprocess
import sys
import threading
import urllib.request
from collections.abc import Callable, Iterator
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from teamagent.media.contracts import (
    AcquireOperation,
    FrameOperation,
    MediaOperation,
    PdfOperation,
    ProposalPptxOperation,
    ProxyOperation,
    S3ObjectRef,
    SlidesOperation,
    ThumbnailOperation,
    TikTokAcquireOperation,
)
from teamagent.media.deadline import DeadlineBudget, MediaDeadlineExceededError
from teamagent.media.security import (
    ALLOWED_YTDLP_EXTRACTORS,
    PublicHttpsRedirectHandler,
    public_dns_only,
    validate_acquire_url,
    validate_public_https_url,
)

logger = logging.getLogger(__name__)

ObjectLoader = Callable[[S3ObjectRef, Path], Path]
_FFMPEG_TIMEOUT_S = 180
_FFMPEG_PROTOCOL_WHITELIST = "file,pipe"
_FFMPEG_PROTOCOL_BLACKLIST = "http,https,tcp,tls,udp,rtp,ftp,gopher,sftp,ssh"
_MEDIA_FILENAME = re.compile(r"^[A-Za-z0-9_.-]+$")
_EXTERNAL_HTML_REF = re.compile(
    r"""(?is)(?:src|href)\s*=\s*["']\s*(?:https?:)?//|url\(\s*["']?\s*(?:https?:)?//"""
)
_PLACEHOLDER = re.compile(r"[｛{]\s*(\d+)\s*[:：]?[^｝}]*[｝}]")
_ACTIVE_PROCESS_GROUPS: set[int] = set()
_ACTIVE_PROCESS_GROUPS_LOCK = threading.Lock()


class MediaOperationError(RuntimeError):
    """workerが安全なerror_codeへ変換できる操作失敗。"""

    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code


def _remaining(budget: DeadlineBudget, cap_s: float | None = None) -> float:
    try:
        return budget.remaining(cap_s=cap_s)
    except MediaDeadlineExceededError as exc:
        raise MediaOperationError(
            "MEDIA_JOB_DEADLINE_EXCEEDED",
            "media job deadline exceeded",
        ) from exc


@dataclass(frozen=True, slots=True)
class ProducedArtifact:
    name: str
    path: Path
    content_type: str
    relative_key: str | None = None


@dataclass(frozen=True, slots=True)
class OperationOutput:
    artifacts: tuple[ProducedArtifact, ...]
    metadata: dict[str, Any]


def _terminate_process_group(
    process: subprocess.Popen[bytes],
    *,
    grace_s: float = 2.0,
) -> None:
    if process.poll() is not None:
        return
    try:
        os.killpg(process.pid, signal.SIGTERM)
    except ProcessLookupError:
        return
    try:
        process.communicate(timeout=grace_s)
        return
    except subprocess.TimeoutExpired:
        pass
    try:
        os.killpg(process.pid, signal.SIGKILL)
    except ProcessLookupError:
        return
    process.communicate()


def terminate_active_process_groups() -> None:
    """Terminate every renderer/tool process group owned by this worker."""

    with _ACTIVE_PROCESS_GROUPS_LOCK:
        process_groups = tuple(_ACTIVE_PROCESS_GROUPS)
    for process_group in process_groups:
        try:
            os.killpg(process_group, signal.SIGTERM)
        except ProcessLookupError:
            continue


def _run_process(
    command: list[str],
    *,
    budget: DeadlineBudget,
    timeout_s: float,
    timeout_code: str,
    timeout_message: str,
    cwd: Path | None = None,
    env: dict[str, str] | None = None,
) -> subprocess.CompletedProcess[bytes]:
    timeout = _remaining(budget, timeout_s)
    process = subprocess.Popen(
        command,
        cwd=cwd,
        env=env,
        stdin=subprocess.DEVNULL,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        start_new_session=True,
    )
    with _ACTIVE_PROCESS_GROUPS_LOCK:
        _ACTIVE_PROCESS_GROUPS.add(process.pid)
    try:
        try:
            stdout, stderr = process.communicate(timeout=timeout)
        except subprocess.TimeoutExpired as exc:
            _terminate_process_group(process)
            try:
                _remaining(budget)
            except MediaOperationError:
                raise
            raise MediaOperationError(timeout_code, timeout_message) from exc
    finally:
        with _ACTIVE_PROCESS_GROUPS_LOCK:
            _ACTIVE_PROCESS_GROUPS.discard(process.pid)
    return subprocess.CompletedProcess(command, process.returncode, stdout, stderr)


def _run(
    command: list[str],
    *,
    budget: DeadlineBudget,
    timeout_s: int = _FFMPEG_TIMEOUT_S,
) -> None:
    completed = _run_process(
        command,
        budget=budget,
        timeout_s=float(timeout_s),
        timeout_code="MEDIA_PROCESS_TIMEOUT",
        timeout_message="media subprocess timed out",
    )
    if completed.returncode != 0:
        logger.warning(
            "media subprocess failed: command=%s stderr=%s",
            command[0],
            completed.stderr.decode("utf-8", "replace")[-300:],
        )
        raise MediaOperationError("MEDIA_PROCESS_FAILED", "media subprocess failed")
    _remaining(budget)


def _ffmpeg_input(source: Path) -> list[str]:
    """Return fail-closed input flags for already-staged local media only."""

    return [
        "-protocol_whitelist",
        _FFMPEG_PROTOCOL_WHITELIST,
        "-protocol_blacklist",
        _FFMPEG_PROTOCOL_BLACKLIST,
        "-i",
        str(source),
    ]


def _safe_file(path: Path, root: Path) -> Path:
    if path.is_symlink():
        raise MediaOperationError("MEDIA_OUTPUT_PATH_INVALID", "output is a symlink")
    resolved = path.resolve(strict=True)
    if not resolved.is_relative_to(root.resolve()) or not resolved.is_file():
        raise MediaOperationError("MEDIA_OUTPUT_PATH_INVALID", "output escaped request directory")
    return resolved


def _acquire(
    operation: AcquireOperation,
    workdir: Path,
    budget: DeadlineBudget,
) -> OperationOutput:
    _remaining(budget)
    validate_acquire_url(operation.url)
    try:
        import yt_dlp
    except ImportError as exc:
        raise MediaOperationError("MEDIA_YTDLP_UNAVAILABLE", "yt-dlp is not installed") from exc

    download_dir = workdir / "download"
    download_dir.mkdir(mode=0o700)
    options: dict[str, Any] = {
        "allowed_extractors": list(ALLOWED_YTDLP_EXTRACTORS),
        "cachedir": False,
        "continuedl": False,
        "format": (
            f"best[filesize<={operation.max_bytes}][vcodec!=none][acodec!=none]/"
            f"best[filesize_approx<={operation.max_bytes}][vcodec!=none][acodec!=none]/worst"
        ),
        "fragment_retries": 2,
        # The worker's native ffmpeg is local-file-only. Prevent yt-dlp from
        # discovering or spawning it for remote manifests/merges; formats that
        # require an external downloader or postprocessor fail closed instead.
        "external_downloader": None,
        "ffmpeg_location": "/nonexistent/teamagent-ffmpeg-disabled",
        "fixup": "never",
        "hls_prefer_native": True,
        "max_filesize": operation.max_bytes,
        "noplaylist": True,
        "nopart": True,
        "outtmpl": str(download_dir / "media.%(ext)s"),
        "overwrites": False,
        "paths": {"home": str(download_dir), "temp": str(download_dir)},
        "postprocessors": [],
        "quiet": True,
        "retries": 2,
        "socket_timeout": max(1.0, _remaining(budget, 20.0)),
        "progress_hooks": [lambda _status: _remaining(budget)],
        "verbose": False,
        "writeinfojson": False,
        "writethumbnail": False,
    }
    try:
        with public_dns_only(), yt_dlp.YoutubeDL(options) as ydl:
            info = ydl.extract_info(operation.url, download=True)
    except MediaOperationError:
        raise
    except Exception as exc:
        _remaining(budget)
        raise MediaOperationError("MEDIA_ACQUIRE_FAILED", "allowed media acquire failed") from exc
    _remaining(budget)
    if not isinstance(info, dict):
        raise MediaOperationError("MEDIA_ACQUIRE_EMPTY", "yt-dlp returned no media metadata")
    extractor = str(info.get("extractor_key") or info.get("extractor") or "").lower()
    if not extractor.startswith(("youtube", "tiktok", "instagram")):
        raise MediaOperationError("MEDIA_EXTRACTOR_BLOCKED", "yt-dlp selected a blocked extractor")
    candidates = sorted(
        path
        for path in download_dir.iterdir()
        if path.is_file() and not path.name.endswith((".part", ".ytdl", ".json"))
    )
    if len(candidates) != 1:
        raise MediaOperationError("MEDIA_ACQUIRE_OUTPUT_INVALID", "expected exactly one media file")
    source = _safe_file(candidates[0], workdir)
    size = source.stat().st_size
    if size < 1 or size > operation.max_bytes:
        raise MediaOperationError("MEDIA_ACQUIRE_SIZE_EXCEEDED", "acquired media exceeded limit")
    suffix = source.suffix.lower() if _MEDIA_FILENAME.fullmatch(source.name) else ".bin"
    destination = workdir / f"acquired{suffix}"
    source.replace(destination)
    mime = {
        ".mp4": "video/mp4",
        ".webm": "video/webm",
        ".mov": "video/quicktime",
        ".mkv": "video/x-matroska",
    }.get(suffix, "application/octet-stream")
    return OperationOutput(
        artifacts=(ProducedArtifact("media", destination, mime),),
        metadata={"extractor": extractor, "size": size},
    )


def _node_json(
    command: list[str],
    *,
    workdir: Path,
    budget: DeadlineBudget,
    timeout_s: int,
) -> dict[str, Any]:
    environment = {
        **os.environ,
        "HOME": str(workdir / "home"),
        "TMPDIR": str(workdir / "tmp"),
        "XDG_CACHE_HOME": str(workdir / "cache"),
        "XDG_CONFIG_HOME": str(workdir / "config"),
        "XDG_DATA_HOME": str(workdir / "data"),
        "XDG_STATE_HOME": str(workdir / "state"),
    }
    for directory in ("home", "tmp", "cache", "config", "data", "state"):
        (workdir / directory).mkdir(mode=0o700, exist_ok=True)
    completed = _run_process(
        command,
        cwd=workdir,
        env=environment,
        budget=budget,
        timeout_s=float(timeout_s),
        timeout_code="MEDIA_TIKTOK_TIMEOUT",
        timeout_message="TikTok browser timed out",
    )
    try:
        result = json.loads(completed.stdout.decode("utf-8"))
    except (UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise MediaOperationError(
            "MEDIA_TIKTOK_OUTPUT_INVALID",
            "TikTok output is invalid",
        ) from exc
    if completed.returncode != 0 or not isinstance(result, dict) or not result.get("ok"):
        raise MediaOperationError("MEDIA_TIKTOK_FAILED", "TikTok browser job failed")
    _remaining(budget)
    return result


def _renderer_json(
    manifest: dict[str, Any],
    *,
    workdir: Path,
    budget: DeadlineBudget,
) -> dict[str, Any]:
    manifest_path = workdir / "render-request.json"
    manifest_path.write_text(
        json.dumps(
            manifest,
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
        ),
        encoding="utf-8",
    )
    manifest_path.chmod(0o600)
    completed = _run_process(
        [
            sys.executable,
            "-m",
            "teamagent.media.render_child",
            manifest_path.name,
        ],
        cwd=workdir,
        env={**os.environ, "PYTHONUNBUFFERED": "1"},
        budget=budget,
        timeout_s=_remaining(budget),
        timeout_code="MEDIA_RENDER_TIMEOUT",
        timeout_message="media renderer timed out",
    )
    try:
        result = json.loads(completed.stdout.decode("utf-8"))
    except (UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise MediaOperationError(
            "MEDIA_RENDER_OUTPUT_INVALID",
            "renderer output is invalid",
        ) from exc
    if completed.returncode != 0 or not isinstance(result, dict) or result.get("ok") is not True:
        code = str(result.get("code", "MEDIA_RENDER_FAILED"))
        if not re.fullmatch(r"MEDIA_[A-Z0-9_]{1,80}", code):
            code = "MEDIA_RENDER_FAILED"
        raise MediaOperationError(code, "media renderer failed")
    metadata = result.get("metadata")
    if not isinstance(metadata, dict):
        raise MediaOperationError(
            "MEDIA_RENDER_OUTPUT_INVALID",
            "renderer metadata is invalid",
        )
    _remaining(budget)
    return metadata


def _post_row(video: dict[str, Any], keyword: str, rank: int, pid: str) -> dict[str, Any]:
    stats_value = video.get("stats")
    author_value = video.get("author")
    stats: dict[str, Any] = stats_value if isinstance(stats_value, dict) else {}
    author: dict[str, Any] = author_value if isinstance(author_value, dict) else {}
    plays = int(stats.get("playCount") or 0)
    likes = int(stats.get("diggCount") or 0)
    comments = int(stats.get("commentCount") or 0)
    shares = int(stats.get("shareCount") or 0)
    saves = int(stats.get("collectCount") or 0)
    followers = int(author.get("followerCount") or 0)
    return {
        "pid": pid,
        "kw": keyword,
        "rank_display": rank,
        "url": str(video.get("url") or ""),
        "title": str(video.get("desc") or ""),
        "account_id": str(author.get("uniqueId") or ""),
        "account_name": str(author.get("nickname") or ""),
        "followers": followers,
        "plays": plays,
        "likes": likes,
        "comments": comments,
        "shares": shares,
        "saves": saves,
        "eg_rate": round((likes + comments + shares + saves) / max(plays, 1) * 100, 6),
        "save_rate": round(saves / max(plays, 1) * 100, 6),
        "create_time": int(video.get("createTime") or 0),
        "cover_url": str(video.get("coverUrl") or ""),
    }


def _fetch_public_image(
    url: str,
    destination: Path,
    *,
    budget: DeadlineBudget,
    width: int = 480,
) -> bool:
    if not url:
        return False
    try:
        validate_public_https_url(url)
        request = urllib.request.Request(
            url,
            headers={
                "User-Agent": (
                    "Mozilla/5.0 (X11; Linux aarch64) AppleWebKit/537.36 Chrome/150 Safari/537.36"
                )
            },
        )
        opener = urllib.request.build_opener(PublicHttpsRedirectHandler())
        with (
            public_dns_only(),
            opener.open(
                request,
                timeout=_remaining(budget, 20.0),
            ) as response,
        ):
            length = int(response.headers.get("Content-Length") or 0)
            if length > 8 * 1024 * 1024:
                return False
            body = response.read(8 * 1024 * 1024 + 1)
        if not body or len(body) > 8 * 1024 * 1024:
            return False
        raw = destination.with_suffix(".raw")
        raw.write_bytes(body)
        _run(
            [
                "ffmpeg",
                "-nostdin",
                "-y",
                *_ffmpeg_input(raw),
                "-frames:v",
                "1",
                "-vf",
                f"scale=w={width}:h=-2:flags=lanczos",
                "-q:v",
                "5",
                str(destination),
            ],
            budget=budget,
            timeout_s=30,
        )
        raw.unlink(missing_ok=True)
        return destination.exists() and destination.stat().st_size > 0
    except MediaOperationError:
        raise
    except Exception:
        _remaining(budget)
        return False


def _tiktok_acquire(
    operation: TikTokAcquireOperation,
    workdir: Path,
    budget: DeadlineBudget,
) -> OperationOutput:
    node = os.environ.get("TIKTOK_NODE_BIN", "/usr/bin/node")
    scraper = os.environ.get(
        "TIKTOK_SCRAPER_PATH",
        "/app/tools/tiktok_scraper/search.mjs",
    )
    if not Path(node).exists() or not Path(scraper).exists():
        raise MediaOperationError("MEDIA_TIKTOK_RUNTIME_MISSING", "TikTok runtime is missing")

    posts: list[dict[str, Any]] = []
    for keyword_index, keyword in enumerate(operation.keywords):
        _remaining(budget)
        result = _node_json(
            [
                node,
                scraper,
                "--query",
                keyword,
                "--type",
                "keyword",
                "--max",
                str(operation.n_per_kw),
            ],
            workdir=workdir,
            budget=budget,
            timeout_s=120,
        )
        videos = result.get("videos")
        if not isinstance(videos, list):
            raise MediaOperationError("MEDIA_TIKTOK_OUTPUT_INVALID", "TikTok videos are invalid")
        for rank, raw_video in enumerate(videos[: operation.n_per_kw], 1):
            if not isinstance(raw_video, dict):
                continue
            pid = f"p{keyword_index + 1:02d}{rank:03d}"
            row = _post_row(raw_video, keyword, rank, pid)
            if not row["url"]:
                continue
            try:
                validate_acquire_url(str(row["url"]))
            except ValueError:
                continue
            posts.append(row)

    artifacts: list[ProducedArtifact] = []
    manifest_items: list[dict[str, Any]] = []
    posts_by_keyword = {
        keyword: [post for post in posts if post["kw"] == keyword] for keyword in operation.keywords
    }
    selected: set[str] = set()
    for keyword in operation.keywords:
        _remaining(budget)
        candidates = list(posts_by_keyword[keyword])
        if operation.sort == "save_rate":
            candidates.sort(key=lambda post: (-float(post["save_rate"]), int(post["rank_display"])))
        elif operation.sort == "recent":
            candidates.sort(key=lambda post: (-int(post["create_time"]), int(post["rank_display"])))
        selected.update(str(post["pid"]) for post in candidates[: operation.videos_per_kw])

    for post in posts:
        _remaining(budget)
        pid = str(post["pid"])
        thumbnail_relative = f"thumbs/{pid}.jpg"
        thumbnail = workdir / f"{pid}.jpg"
        if _fetch_public_image(
            str(post["cover_url"]),
            thumbnail,
            budget=budget,
        ):
            artifacts.append(
                ProducedArtifact(
                    f"thumb-{pid}",
                    thumbnail,
                    "image/jpeg",
                    relative_key=thumbnail_relative,
                )
            )
        downloaded = False
        video_relative = f"videos/{pid}.mp4"
        video_path = workdir / f"{pid}.mp4"
        if pid in selected:
            try:
                _node_json(
                    [
                        node,
                        scraper,
                        "--mode",
                        "download",
                        "--url",
                        str(post["url"]),
                        "--out",
                        str(video_path),
                        "--max-bytes",
                        str(operation.max_video_bytes),
                    ],
                    workdir=workdir,
                    budget=budget,
                    timeout_s=120,
                )
                downloaded = (
                    video_path.exists()
                    and 0 < video_path.stat().st_size <= operation.max_video_bytes
                )
            except MediaOperationError as exc:
                if exc.code == "MEDIA_JOB_DEADLINE_EXCEEDED":
                    raise
                downloaded = False
        if downloaded:
            artifacts.append(
                ProducedArtifact(
                    f"video-{pid}",
                    video_path,
                    "video/mp4",
                    relative_key=video_relative,
                )
            )
        else:
            video_path.unlink(missing_ok=True)
        manifest_items.append(
            {
                "pid": pid,
                "kw": post["kw"],
                "downloaded": downloaded,
                "video_path": video_relative,
                "thumb_path": thumbnail_relative,
                "tiktok_url": post["url"],
            }
        )

    posts_path = workdir / "posts.normalized.json"
    posts_path.write_text(
        json.dumps({"posts": posts}, ensure_ascii=False, separators=(",", ":")),
        encoding="utf-8",
    )
    config_path = workdir / "config.json"
    config_path.write_text(
        json.dumps(
            operation.client.model_dump(mode="json", exclude_none=True),
            ensure_ascii=False,
            separators=(",", ":"),
        ),
        encoding="utf-8",
    )
    manifest_path = workdir / "manifest.json"
    manifest_path.write_text(
        json.dumps({"items": manifest_items}, ensure_ascii=False, separators=(",", ":")),
        encoding="utf-8",
    )
    artifacts.extend(
        (
            ProducedArtifact(
                "posts.json",
                posts_path,
                "application/json",
                relative_key="posts.normalized.json",
            ),
            ProducedArtifact(
                "config.json",
                config_path,
                "application/json",
                relative_key="config.json",
            ),
            ProducedArtifact(
                "manifest.json",
                manifest_path,
                "application/json",
                relative_key="videos/manifest.json",
            ),
        )
    )
    downloaded_count = sum(bool(item["downloaded"]) for item in manifest_items)
    return OperationOutput(
        tuple(artifacts),
        {
            "counts": {
                "kw": len(operation.keywords),
                "posts": len(posts),
                "videos": downloaded_count,
            },
            "s3_prefix": "",
        },
    )


def _proxy(
    operation: ProxyOperation,
    workdir: Path,
    load: ObjectLoader,
    budget: DeadlineBudget,
) -> OperationOutput:
    _remaining(budget)
    source = load(operation.source, workdir / "source.bin")
    if source.stat().st_size <= operation.limit_bytes and not operation.preview:
        destination = workdir / "proxy.bin"
        shutil.copyfile(source, destination)
        return OperationOutput(
            artifacts=(ProducedArtifact("proxy", destination, operation.source.content_type),),
            metadata={"transcoded": False, "size": destination.stat().st_size},
        )

    destination = workdir / "proxy.mp4"
    starting_crf = 30 if operation.preview else 27
    ladder = (
        (starting_crf, operation.long_edge),
        (starting_crf + 3, max(360, operation.long_edge * 3 // 4)),
        (starting_crf + 6, max(240, operation.long_edge // 2)),
    )
    for crf, edge in ladder:
        destination.unlink(missing_ok=True)
        _run(
            [
                "ffmpeg",
                "-nostdin",
                "-y",
                *_ffmpeg_input(source),
                "-vf",
                (
                    f"scale=w={edge}:h={edge}:"
                    "force_original_aspect_ratio=decrease:force_divisible_by=2"
                ),
                "-c:v",
                "libx264",
                "-preset",
                "veryfast",
                "-crf",
                str(crf),
                "-c:a",
                "aac",
                "-b:a",
                "96k",
                "-movflags",
                "+faststart",
                str(destination),
            ],
            budget=budget,
        )
        if destination.stat().st_size <= operation.limit_bytes:
            break
    if not destination.exists() or destination.stat().st_size > operation.limit_bytes:
        raise MediaOperationError("MEDIA_PROXY_SIZE_EXCEEDED", "proxy did not fit size bound")
    return OperationOutput(
        artifacts=(ProducedArtifact("proxy", destination, "video/mp4"),),
        metadata={"transcoded": True, "size": destination.stat().st_size},
    )


def _frames(
    operation: FrameOperation,
    workdir: Path,
    load: ObjectLoader,
    budget: DeadlineBudget,
) -> OperationOutput:
    source = load(operation.source, workdir / "source.bin")
    artifacts: list[ProducedArtifact] = []
    frame_rows: list[dict[str, Any]] = []
    for index, second in enumerate(operation.timecodes):
        _remaining(budget)
        destination = workdir / f"frame-{index:02d}.jpg"
        _run(
            [
                "ffmpeg",
                "-nostdin",
                "-y",
                "-ss",
                f"{second:.3f}",
                *_ffmpeg_input(source),
                "-frames:v",
                "1",
                "-vf",
                f"scale=w={operation.width}:h=-2:flags=lanczos",
                "-q:v",
                "6",
                str(destination),
            ],
            budget=budget,
            timeout_s=60,
        )
        if not destination.exists() or destination.stat().st_size == 0:
            raise MediaOperationError("MEDIA_FRAME_EMPTY", "ffmpeg produced an empty frame")
        name = f"frame-{index:02d}"
        artifacts.append(ProducedArtifact(name, destination, "image/jpeg"))
        frame_rows.append({"name": name, "sec": second})
    manifest = workdir / "frames.json"
    manifest.write_text(
        json.dumps({"frames": frame_rows}, ensure_ascii=False, separators=(",", ":")),
        encoding="utf-8",
    )
    artifacts.append(ProducedArtifact("frames.json", manifest, "application/json"))
    return OperationOutput(tuple(artifacts), {"count": len(frame_rows)})


def _palette(rgb: bytes) -> tuple[list[str], float, float]:
    pixels = [(rgb[i], rgb[i + 1], rgb[i + 2]) for i in range(0, len(rgb) - 2, 3)]
    if not pixels:
        return [], 0.0, 0.0
    bins: dict[tuple[int, int, int], list[tuple[int, int, int]]] = {}
    for red, green, blue in pixels:
        bins.setdefault((red >> 5, green >> 5, blue >> 5), []).append((red, green, blue))
    colors: list[str] = []
    for group in sorted(bins.values(), key=len, reverse=True)[:3]:
        count = len(group)
        red = sum(pixel[0] for pixel in group) // count
        green = sum(pixel[1] for pixel in group) // count
        blue = sum(pixel[2] for pixel in group) // count
        colors.append(f"#{red:02x}{green:02x}{blue:02x}")
    count = len(pixels)
    brightness = sum(0.299 * r + 0.587 * g + 0.114 * b for r, g, b in pixels) / count
    warmth = sum(r - b for r, _g, b in pixels) / count
    return colors, round(brightness / 255.0, 3), round(max(-1.0, min(1.0, warmth / 255)), 3)


def _thumbnail(
    operation: ThumbnailOperation,
    workdir: Path,
    load: ObjectLoader,
    budget: DeadlineBudget,
) -> OperationOutput:
    if operation.source is not None:
        source = load(operation.source, workdir / "source.bin")
    else:
        source = workdir / "url-source.jpg"
        if operation.url is None or not _fetch_public_image(
            operation.url,
            source,
            budget=budget,
            width=operation.width,
        ):
            raise MediaOperationError(
                "MEDIA_THUMBNAIL_FETCH_FAILED",
                "thumbnail URL fetch failed",
            )
    image = workdir / "thumbnail.jpg"
    rgb = workdir / "thumbnail.rgb"
    _run(
        [
            "ffmpeg",
            "-nostdin",
            "-y",
            *_ffmpeg_input(source),
            "-frames:v",
            "1",
            "-vf",
            f"scale=w={operation.width}:h=-2:flags=lanczos",
            "-q:v",
            "5",
            str(image),
        ],
        budget=budget,
        timeout_s=60,
    )
    _run(
        [
            "ffmpeg",
            "-nostdin",
            "-y",
            *_ffmpeg_input(source),
            "-frames:v",
            "1",
            "-vf",
            "scale=8:8:flags=area",
            "-f",
            "rawvideo",
            "-pix_fmt",
            "rgb24",
            str(rgb),
        ],
        budget=budget,
        timeout_s=60,
    )
    colors, brightness, warmth = _palette(rgb.read_bytes())
    rgb.unlink(missing_ok=True)
    metadata = {
        "swatches": colors,
        "brightness01": brightness,
        "warmth": warmth,
    }
    manifest = workdir / "thumbnail.json"
    manifest.write_text(
        json.dumps(metadata, ensure_ascii=False, separators=(",", ":")),
        encoding="utf-8",
    )
    return OperationOutput(
        (
            ProducedArtifact("thumbnail", image, "image/jpeg"),
            ProducedArtifact("thumbnail.json", manifest, "application/json"),
        ),
        metadata,
    )


def _slides(
    operation: SlidesOperation,
    workdir: Path,
    load: ObjectLoader,
    budget: DeadlineBudget,
) -> OperationOutput:
    html_path = load(operation.html, workdir / "slides.html")
    if html_path.stat().st_size > 2 * 1024 * 1024:
        raise MediaOperationError("MEDIA_HTML_SIZE_EXCEEDED", "slides HTML exceeds limit")
    html = html_path.read_text(encoding="utf-8")
    if _EXTERNAL_HTML_REF.search(html):
        raise MediaOperationError("MEDIA_HTML_NETWORK_REFERENCE", "slides HTML references network")
    destination = workdir / "slides.pptx"
    metadata = _renderer_json(
        {
            "kind": "slides",
            "html": html_path.name,
            "output": destination.name,
            "selector": operation.selector,
            "width": operation.width,
            "height": operation.height,
            "scale": operation.device_scale_factor,
        },
        workdir=workdir,
        budget=budget,
    )
    if not destination.is_file() or destination.stat().st_size < 5:
        raise MediaOperationError("MEDIA_PPTX_EMPTY", "slides renderer produced no PPTX")
    return OperationOutput(
        (
            ProducedArtifact(
                "slides.pptx",
                destination,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ),
        ),
        metadata,
    )


def _iter_text_frames(presentation: Any) -> Iterator[Any]:
    def walk(shapes: Any) -> Iterator[Any]:
        for shape in shapes:
            if getattr(shape, "shape_type", None) == 6:
                yield from walk(shape.shapes)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame
            elif getattr(shape, "has_text_frame", False):
                yield shape.text_frame

    for slide in presentation.slides:
        yield from walk(slide.shapes)


def _replace_placeholders(text_frame: Any, placeholders: dict[int, str]) -> None:
    paragraphs = list(text_frame.paragraphs)
    combined = "".join(paragraph.text for paragraph in paragraphs)
    if not combined:
        return
    replaced = _PLACEHOLDER.sub(
        lambda match: placeholders.get(int(match.group(1)), match.group(0)),
        combined,
    )
    if replaced == combined:
        return
    paragraphs[0].text = replaced
    for paragraph in paragraphs[1:]:
        paragraph.text = ""


def _iter_proposal_image_slots(presentation: Any) -> Iterator[tuple[Any, Any]]:
    for slide in presentation.slides:
        candidates: list[Any] = []
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) in (6, 13):
                continue
            try:
                width = int(shape.width)
                top = int(shape.top)
            except (TypeError, ValueError):
                continue
            if abs(width - 609905) <= 80000 and abs(top - 5454720) <= 80000:
                candidates.append(shape)
        for shape in sorted(candidates, key=lambda item: int(item.left)):
            yield slide, shape


def _inject_proposal_evidence(
    presentation: Any,
    operation: ProposalPptxOperation,
    workdir: Path,
    load: ObjectLoader,
    budget: DeadlineBudget,
) -> int:
    from pptx.util import Emu

    slots = _iter_proposal_image_slots(presentation)
    injected = 0
    for index, evidence in enumerate(
        sorted(operation.evidence, key=lambda item: (item.placeholder_id, item.rank))
    ):
        _remaining(budget)
        image = load(evidence.source, workdir / f"evidence-{index:02d}.bin")
        try:
            slide, shape = next(slots)
        except StopIteration:
            break
        try:
            picture = slide.shapes.add_picture(
                io.BytesIO(image.read_bytes()),
                Emu(int(shape.left)),
                Emu(int(shape.top)),
            )
        except Exception:
            continue
        if picture.width and picture.height:
            scale = int(shape.height) / picture.height
            picture.height = Emu(int(shape.height))
            picture.width = Emu(int(picture.width * scale))
        injected += 1
    return injected


def _proposal_pptx(
    operation: ProposalPptxOperation,
    workdir: Path,
    load: ObjectLoader,
    budget: DeadlineBudget,
) -> OperationOutput:
    template = load(operation.template, workdir / "template.pptx")
    composer_path = load(operation.composer_json, workdir / "composer.json")
    evidence_manifest: list[dict[str, Any]] = []
    for index, evidence in enumerate(
        sorted(operation.evidence, key=lambda item: (item.placeholder_id, item.rank))
    ):
        image = load(evidence.source, workdir / f"evidence-{index:02d}.bin")
        evidence_manifest.append(
            {
                "placeholder_id": evidence.placeholder_id,
                "rank": evidence.rank,
                "path": image.name,
            }
        )
    destination = workdir / "proposal.pptx"
    metadata = _renderer_json(
        {
            "kind": "proposal_pptx",
            "template": template.name,
            "composer": composer_path.name,
            "evidence": evidence_manifest,
            "output": destination.name,
            "fail_if_missing": operation.fail_if_missing,
        },
        workdir=workdir,
        budget=budget,
    )
    if not destination.is_file() or destination.stat().st_size < 5:
        raise MediaOperationError("MEDIA_PPTX_EMPTY", "proposal renderer produced no PPTX")
    return OperationOutput(
        (
            ProducedArtifact(
                "proposal.pptx",
                destination,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ),
        ),
        metadata,
    )


def _pdf(
    operation: PdfOperation,
    workdir: Path,
    load: ObjectLoader,
    budget: DeadlineBudget,
) -> OperationOutput:
    html_path = load(operation.html, workdir / "document.html")
    if html_path.stat().st_size > 2 * 1024 * 1024:
        raise MediaOperationError("MEDIA_HTML_SIZE_EXCEEDED", "PDF HTML exceeds limit")
    html = html_path.read_text(encoding="utf-8")
    if _EXTERNAL_HTML_REF.search(html):
        raise MediaOperationError("MEDIA_HTML_NETWORK_REFERENCE", "PDF HTML references network")
    destination = workdir / "document.pdf"
    metadata = _renderer_json(
        {
            "kind": "pdf",
            "html": html_path.name,
            "output": destination.name,
        },
        workdir=workdir,
        budget=budget,
    )
    if not destination.exists() or destination.stat().st_size < 5:
        raise MediaOperationError("MEDIA_PDF_EMPTY", "weasyprint produced an empty PDF")
    return OperationOutput(
        (ProducedArtifact("document.pdf", destination, "application/pdf"),),
        metadata,
    )


def execute_operation(
    operation: MediaOperation,
    *,
    workdir: Path,
    load_object: ObjectLoader,
    budget: DeadlineBudget,
) -> OperationOutput:
    """strict operationを1つだけ実行する。"""

    _remaining(budget)
    if isinstance(operation, AcquireOperation):
        output = _acquire(operation, workdir, budget)
    elif isinstance(operation, TikTokAcquireOperation):
        output = _tiktok_acquire(operation, workdir, budget)
    elif isinstance(operation, ProxyOperation):
        output = _proxy(operation, workdir, load_object, budget)
    elif isinstance(operation, FrameOperation):
        output = _frames(operation, workdir, load_object, budget)
    elif isinstance(operation, ThumbnailOperation):
        output = _thumbnail(operation, workdir, load_object, budget)
    elif isinstance(operation, SlidesOperation):
        output = _slides(operation, workdir, load_object, budget)
    elif isinstance(operation, ProposalPptxOperation):
        output = _proposal_pptx(operation, workdir, load_object, budget)
    elif isinstance(operation, PdfOperation):
        output = _pdf(operation, workdir, load_object, budget)
    else:
        raise MediaOperationError("MEDIA_OPERATION_UNSUPPORTED", "unsupported media operation")
    _remaining(budget)
    return output


__all__ = [
    "MediaOperationError",
    "OperationOutput",
    "ProducedArtifact",
    "execute_operation",
    "terminate_active_process_groups",
]
