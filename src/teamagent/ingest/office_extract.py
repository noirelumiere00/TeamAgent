"""docx / pptx / xlsx と Google native gdoc のテキスト抽出 + ページ化ユーティリティ。

ingest パイプライン用の純粋関数群。pdf_extract.py と同じ I/F を提供する
（pipeline 側が ``chunk_pages`` を再利用できるように ``[(page_num, text), ...]`` で返す）。

依存（python-docx / python-pptx / openpyxl）は遅延 import — CI で未インストールでも
モジュール import は通る。pyproject に依存記載済み（>=1.1.0 / >=1.0.0 / >=3.1.0）。

Usage:
    from teamagent.ingest.office_extract import extract_office_pages, OFFICE_BINARY_MIMES

    pages = extract_office_pages(data, mime_type="...pptx")
    # docx → [(1, "doc text...")]
    # pptx → [(1, "slide1..."), (2, "slide2..."), ...]
    # xlsx → [(1, "sheet1 text..."), (2, "sheet2 text..."), ...]
"""

from __future__ import annotations

import hashlib
import re
import zipfile
from collections.abc import Callable
from io import BytesIO
from xml.etree import ElementTree
from xml.parsers import expat

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
GDOC_NATIVE_MIME = "application/vnd.google-apps.document"

# pipeline 側の mime_type 判定用にひとまとめで export。
OFFICE_BINARY_MIMES = frozenset({DOCX_MIME, PPTX_MIME, XLSX_MIME})

# known-invalid fingerprint の validator 世代。検証規則・上限・必須partの解釈を変えたら
# 必ず値を更新し、同じ Drive payload でも新規規則で再検証させる。
OFFICE_VALIDATOR_SCHEMA_VERSION = "ooxml-safe-v3"

# ZIP/OOXML の有界検証。Drive download の hard cap と同じ compressed input 上限を持ち、
# metadata 上限を通った member だけを chunk 単位で読み、CRC を確認する。
MAX_OFFICE_COMPRESSED_BYTES = 256 * 1024 * 1024
MAX_OFFICE_ZIP_MEMBERS = 20_000
MAX_OFFICE_TOTAL_UNCOMPRESSED_BYTES = 512 * 1024 * 1024
MAX_OFFICE_MEMBER_UNCOMPRESSED_BYTES = 128 * 1024 * 1024
MAX_OFFICE_COMPRESSION_RATIO = 500.0
MAX_OFFICE_REQUIRED_XML_BYTES = 16 * 1024 * 1024
MAX_OFFICE_CONTENT_TYPES_XML_BYTES = 2 * 1024 * 1024
MAX_OFFICE_XML_MEMBER_BYTES = 32 * 1024 * 1024
MAX_OFFICE_TOTAL_XML_BYTES = 128 * 1024 * 1024
MAX_OFFICE_EXTRACTED_CHARACTERS = 2_000_000
_RATIO_CHECK_MIN_BYTES = 1024 * 1024
_ZIP_READ_CHUNK_BYTES = 1024 * 1024

_OOXML_REQUIRED_PART = {
    DOCX_MIME: "word/document.xml",
    PPTX_MIME: "ppt/presentation.xml",
    XLSX_MIME: "xl/workbook.xml",
}
_OOXML_REQUIRED_ROOT = {
    DOCX_MIME: "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}document",
    PPTX_MIME: ("{http://schemas.openxmlformats.org/presentationml/2006/main}presentation"),
    XLSX_MIME: "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}workbook",
}
_OOXML_REQUIRED_CONTENT_TYPE = {
    DOCX_MIME: ("application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"),
    PPTX_MIME: (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
    ),
    XLSX_MIME: ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"),
}
_CONTENT_TYPES_PART = "[Content_Types].xml"
_CONTENT_TYPES_ROOT = "{http://schemas.openxmlformats.org/package/2006/content-types}Types"
_CONTENT_TYPES_OVERRIDE = "{http://schemas.openxmlformats.org/package/2006/content-types}Override"
_ZIP_PREFIXES = (b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")
_COMPOUND_FILE_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
_ENCRYPTED_COMPOUND_MARKERS = (
    "EncryptedPackage".encode("utf-16le"),
    "EncryptionInfo".encode("utf-16le"),
)
_XML_MEMBER_SUFFIXES = (".xml", ".rels")


class _OfficeProgressCallbackError(Exception):
    """progress callback由来の例外をpayload破損へ誤分類しないための内部wrapper。"""

    def __init__(self, cause: Exception) -> None:
        self.cause = cause
        super().__init__(str(cause))


def _report_progress(progress_callback: Callable[[], None] | None) -> None:
    if progress_callback is None:
        return
    try:
        progress_callback()
    except Exception as exc:
        raise _OfficeProgressCallbackError(exc) from exc


class OfficePayloadError(zipfile.BadZipFile):
    """抽出前に判別できる Office payload 不正。

    ``category`` は運用ログで機械集計する安定値:

    - ``html_response``: Drive 本体ではなく HTML 応答
    - ``truncated_download`` / ``size_mismatch``: Drive metadata の size と不一致
    - ``checksum_mismatch``: Drive metadata の MD5 と実バイト列が不一致
    - ``corrupt_zip``: ZIP シグネチャはあるが中央ディレクトリ等が壊れている
    - ``format_mismatch``: Office MIME と実体の OOXML package が一致しない
    - ``unsafe_archive``: 展開量・member数・圧縮率・XML量が安全上限を超える
    - ``unsafe_content_volume``: 抽出本文がファイル単位の文字数上限を超える
    - ``encrypted_office``: password 保護された OOXML / encrypted compound file

    ``BadZipFile`` の subclass にして、既存の fail-open 呼び出し側との互換性を保つ。
    """

    def __init__(
        self,
        category: str,
        *,
        mime_type: str,
        actual_bytes: int,
        expected_bytes: int | None = None,
    ) -> None:
        self.category = category
        self.mime_type = mime_type
        self.actual_bytes = actual_bytes
        self.expected_bytes = expected_bytes
        super().__init__(f"invalid Office payload: {category}")


def _looks_like_html(data: bytes) -> bool:
    """BOM/空白付きも含め、確認画面・エラーページの HTML 応答を判定する。"""
    head = data[:512].lstrip()
    if head.startswith(b"\xef\xbb\xbf"):
        head = head[3:].lstrip()
    lowered = head.lower()
    return (
        lowered.startswith(b"<!doctype html")
        or lowered.startswith(b"<html")
        or (lowered.startswith(b"<?xml") and b"<html" in lowered)
    )


def _looks_like_encrypted_compound_office(data: bytes) -> bool:
    """OOXML password保護で使われるOLE compound containerを識別する。"""
    return data.startswith(_COMPOUND_FILE_MAGIC) and all(
        marker in data for marker in _ENCRYPTED_COMPOUND_MARKERS
    )


def _read_zip_member_bounded(
    package: zipfile.ZipFile,
    info: zipfile.ZipInfo,
    *,
    capture_limit: int | None = None,
    progress_callback: Callable[[], None] | None = None,
) -> bytes | None:
    """1 memberを有界chunkで最後まで読み、CRCを検証する。

    ``ZipExtFile`` は EOF まで読むとCRCを検査する。``testzip``のような無制限展開はせず、
    metadataで保証した上限内だけを読み、必要XML以外はメモリへ保持しない。
    """
    captured = bytearray() if capture_limit is not None else None
    read_bytes = 0
    with package.open(info, "r") as member:
        while True:
            _report_progress(progress_callback)
            chunk = member.read(_ZIP_READ_CHUNK_BYTES)
            if not chunk:
                break
            read_bytes += len(chunk)
            if read_bytes > info.file_size or read_bytes > MAX_OFFICE_MEMBER_UNCOMPRESSED_BYTES:
                raise ValueError("zip member exceeded declared safety bound")
            if captured is not None and capture_limit is not None:
                if len(captured) + len(chunk) > capture_limit:
                    raise ValueError("required XML exceeded safety bound")
                captured.extend(chunk)
    if read_bytes != info.file_size:
        raise zipfile.BadZipFile("zip member size did not match central directory")
    return bytes(captured) if captured is not None else None


def _validate_bounded_xml_syntax(
    data: bytes,
    *,
    progress_callback: Callable[[], None] | None = None,
) -> None:
    """DTD/entity展開を許さず、上限確認済みXMLをstreaming pre-parseする。"""
    parser = expat.ParserCreate()

    def _reject_declaration(*_args: object) -> None:
        raise ValueError("DTD/entity declarations are not permitted in OOXML validation")

    def _reject_external_entity(*_args: object) -> int:
        raise ValueError("external entities are not permitted in OOXML validation")

    parser.StartDoctypeDeclHandler = _reject_declaration
    parser.EntityDeclHandler = _reject_declaration
    parser.UnparsedEntityDeclHandler = _reject_declaration
    parser.NotationDeclHandler = _reject_declaration
    parser.ExternalEntityRefHandler = _reject_external_entity
    parser.SetParamEntityParsing(expat.XML_PARAM_ENTITY_PARSING_NEVER)
    try:
        for start in range(0, len(data), _ZIP_READ_CHUNK_BYTES):
            _report_progress(progress_callback)
            end = min(start + _ZIP_READ_CHUNK_BYTES, len(data))
            parser.Parse(data[start:end], end == len(data))
        if not data:
            parser.Parse(b"", True)
    except ValueError:
        raise
    except expat.ExpatError as exc:
        raise ElementTree.ParseError(str(exc)) from exc


def _parse_bounded_xml(
    data: bytes,
    *,
    progress_callback: Callable[[], None] | None = None,
) -> ElementTree.Element:
    """secure pre-parse済みの上限内XMLだけをElementTree化する。"""
    _validate_bounded_xml_syntax(data, progress_callback=progress_callback)
    return ElementTree.fromstring(data)  # nosec B314  # bounded preparse blocks DTD/entity/external


def _validate_office_payload(
    data: bytes,
    *,
    mime_type: str,
    expected_size: int | None,
    expected_md5: str | None,
    progress_callback: Callable[[], None] | None = None,
) -> None:
    """OOXMLを開く前に、payload同一性・ZIP integrity・必須構造を有界検証する。"""
    actual_size = len(data)

    if _looks_like_html(data):
        raise OfficePayloadError(
            "html_response",
            mime_type=mime_type,
            actual_bytes=actual_size,
            expected_bytes=expected_size,
        )

    if actual_size > MAX_OFFICE_COMPRESSED_BYTES:
        raise OfficePayloadError(
            "unsafe_archive",
            mime_type=mime_type,
            actual_bytes=actual_size,
            expected_bytes=expected_size,
        )

    if expected_size is not None and actual_size != expected_size:
        category = "truncated_download" if actual_size < expected_size else "size_mismatch"
        raise OfficePayloadError(
            category,
            mime_type=mime_type,
            actual_bytes=actual_size,
            expected_bytes=expected_size,
        )

    if expected_md5:
        actual_md5 = hashlib.md5(data, usedforsecurity=False).hexdigest()
        if actual_md5 != expected_md5.lower():
            raise OfficePayloadError(
                "checksum_mismatch",
                mime_type=mime_type,
                actual_bytes=actual_size,
                expected_bytes=expected_size,
            )

    if _looks_like_encrypted_compound_office(data):
        raise OfficePayloadError(
            "encrypted_office",
            mime_type=mime_type,
            actual_bytes=actual_size,
            expected_bytes=expected_size,
        )

    if not zipfile.is_zipfile(BytesIO(data)):
        category = "corrupt_zip" if data.startswith(_ZIP_PREFIXES) else "format_mismatch"
        raise OfficePayloadError(
            category,
            mime_type=mime_type,
            actual_bytes=actual_size,
            expected_bytes=expected_size,
        )

    required_part = _OOXML_REQUIRED_PART[mime_type]
    try:
        _report_progress(progress_callback)
        with zipfile.ZipFile(BytesIO(data)) as package:
            infos = package.infolist()
            if len(infos) > MAX_OFFICE_ZIP_MEMBERS:
                raise OfficePayloadError(
                    "unsafe_archive",
                    mime_type=mime_type,
                    actual_bytes=actual_size,
                    expected_bytes=expected_size,
                )

            names = [info.filename for info in infos]
            if len(set(names)) != len(names):
                raise OfficePayloadError(
                    "unsafe_archive",
                    mime_type=mime_type,
                    actual_bytes=actual_size,
                    expected_bytes=expected_size,
                )
            info_by_name = {info.filename: info for info in infos}

            total_uncompressed = 0
            total_compressed = 0
            total_xml_uncompressed = 0
            for info in infos:
                _report_progress(progress_callback)
                if info.flag_bits & 0x1:
                    raise OfficePayloadError(
                        "encrypted_office",
                        mime_type=mime_type,
                        actual_bytes=actual_size,
                        expected_bytes=expected_size,
                    )
                if info.file_size < 0 or info.compress_size < 0:
                    raise OfficePayloadError(
                        "corrupt_zip",
                        mime_type=mime_type,
                        actual_bytes=actual_size,
                        expected_bytes=expected_size,
                    )
                if info.file_size > MAX_OFFICE_MEMBER_UNCOMPRESSED_BYTES:
                    raise OfficePayloadError(
                        "unsafe_archive",
                        mime_type=mime_type,
                        actual_bytes=actual_size,
                        expected_bytes=expected_size,
                    )
                if info.filename.lower().endswith(_XML_MEMBER_SUFFIXES):
                    if info.file_size > MAX_OFFICE_XML_MEMBER_BYTES:
                        raise OfficePayloadError(
                            "unsafe_archive",
                            mime_type=mime_type,
                            actual_bytes=actual_size,
                            expected_bytes=expected_size,
                        )
                    total_xml_uncompressed += info.file_size
                    if total_xml_uncompressed > MAX_OFFICE_TOTAL_XML_BYTES:
                        raise OfficePayloadError(
                            "unsafe_archive",
                            mime_type=mime_type,
                            actual_bytes=actual_size,
                            expected_bytes=expected_size,
                        )
                if info.file_size >= _RATIO_CHECK_MIN_BYTES:
                    ratio = info.file_size / max(1, info.compress_size)
                    if ratio > MAX_OFFICE_COMPRESSION_RATIO:
                        raise OfficePayloadError(
                            "unsafe_archive",
                            mime_type=mime_type,
                            actual_bytes=actual_size,
                            expected_bytes=expected_size,
                        )
                total_uncompressed += info.file_size
                total_compressed += info.compress_size

            if total_uncompressed > MAX_OFFICE_TOTAL_UNCOMPRESSED_BYTES:
                raise OfficePayloadError(
                    "unsafe_archive",
                    mime_type=mime_type,
                    actual_bytes=actual_size,
                    expected_bytes=expected_size,
                )
            if (
                total_uncompressed >= _RATIO_CHECK_MIN_BYTES
                and total_uncompressed / max(1, total_compressed) > MAX_OFFICE_COMPRESSION_RATIO
            ):
                raise OfficePayloadError(
                    "unsafe_archive",
                    mime_type=mime_type,
                    actual_bytes=actual_size,
                    expected_bytes=expected_size,
                )
            if required_part not in info_by_name or _CONTENT_TYPES_PART not in info_by_name:
                # CRCは先に全memberで確認済み。構造だけが不足しているpackageを区別する。
                for info in infos:
                    _read_zip_member_bounded(
                        package,
                        info,
                        progress_callback=progress_callback,
                    )
                raise OfficePayloadError(
                    "format_mismatch",
                    mime_type=mime_type,
                    actual_bytes=actual_size,
                    expected_bytes=expected_size,
                )
            if info_by_name[required_part].file_size > MAX_OFFICE_REQUIRED_XML_BYTES:
                raise OfficePayloadError(
                    "unsafe_archive",
                    mime_type=mime_type,
                    actual_bytes=actual_size,
                    expected_bytes=expected_size,
                )
            if info_by_name[_CONTENT_TYPES_PART].file_size > MAX_OFFICE_CONTENT_TYPES_XML_BYTES:
                raise OfficePayloadError(
                    "unsafe_archive",
                    mime_type=mime_type,
                    actual_bytes=actual_size,
                    expected_bytes=expected_size,
                )

            required_root: ElementTree.Element | None = None
            content_types_root: ElementTree.Element | None = None
            for info in infos:
                capture_limit: int | None = None
                if info.filename == required_part:
                    capture_limit = MAX_OFFICE_REQUIRED_XML_BYTES
                elif info.filename == _CONTENT_TYPES_PART:
                    capture_limit = MAX_OFFICE_CONTENT_TYPES_XML_BYTES
                elif info.filename.lower().endswith(_XML_MEMBER_SUFFIXES):
                    capture_limit = MAX_OFFICE_XML_MEMBER_BYTES
                captured = _read_zip_member_bounded(
                    package,
                    info,
                    capture_limit=capture_limit,
                    progress_callback=progress_callback,
                )
                if captured is not None:
                    try:
                        if info.filename == required_part:
                            required_root = _parse_bounded_xml(
                                captured,
                                progress_callback=progress_callback,
                            )
                        elif info.filename == _CONTENT_TYPES_PART:
                            content_types_root = _parse_bounded_xml(
                                captured,
                                progress_callback=progress_callback,
                            )
                        else:
                            _validate_bounded_xml_syntax(
                                captured,
                                progress_callback=progress_callback,
                            )
                    except ElementTree.ParseError as exc:
                        raise OfficePayloadError(
                            "corrupt_zip",
                            mime_type=mime_type,
                            actual_bytes=actual_size,
                            expected_bytes=expected_size,
                        ) from exc
                    except ValueError as exc:
                        raise OfficePayloadError(
                            "unsafe_archive",
                            mime_type=mime_type,
                            actual_bytes=actual_size,
                            expected_bytes=expected_size,
                        ) from exc

            if required_root is None or content_types_root is None:
                raise OfficePayloadError(
                    "format_mismatch",
                    mime_type=mime_type,
                    actual_bytes=actual_size,
                    expected_bytes=expected_size,
                )

            if (
                required_root.tag != _OOXML_REQUIRED_ROOT[mime_type]
                or content_types_root.tag != _CONTENT_TYPES_ROOT
            ):
                raise OfficePayloadError(
                    "format_mismatch",
                    mime_type=mime_type,
                    actual_bytes=actual_size,
                    expected_bytes=expected_size,
                )
            expected_part_name = f"/{required_part}"
            expected_content_type = _OOXML_REQUIRED_CONTENT_TYPE[mime_type]
            if not any(
                element.tag == _CONTENT_TYPES_OVERRIDE
                and element.attrib.get("PartName") == expected_part_name
                and element.attrib.get("ContentType") == expected_content_type
                for element in content_types_root
            ):
                raise OfficePayloadError(
                    "format_mismatch",
                    mime_type=mime_type,
                    actual_bytes=actual_size,
                    expected_bytes=expected_size,
                )
    except OfficePayloadError:
        raise
    except _OfficeProgressCallbackError as exc:
        raise exc.cause from exc
    except (zipfile.BadZipFile, zipfile.LargeZipFile, OSError, RuntimeError, ValueError) as exc:
        raise OfficePayloadError(
            "corrupt_zip",
            mime_type=mime_type,
            actual_bytes=actual_size,
            expected_bytes=expected_size,
        ) from exc


# xlsx 抽出時に拾うセル数の上限ガード（極端に大きい sheet で OOM 防止）。
_XLSX_MAX_CELLS_PER_SHEET = 10000

# python-pptx の GROUP shape type 値（MSO_SHAPE_TYPE.GROUP == 6）。
# 列挙体を import せずに数値で比較する（renderer.py:27-28 と同じ流儀）。
_GROUP_SHAPE_TYPE = 6


class _OfficeTextLimitExceededError(ValueError):
    """Office本文抽出がファイル単位の文字数上限を超えた。"""


class _TextBudget:
    """抽出器をまたいで共有する、本文文字数の単純なhard cap。"""

    def __init__(
        self,
        max_chars: int,
        progress_callback: Callable[[], None] | None = None,
    ) -> None:
        if max_chars < 1:
            raise ValueError("max_chars must be positive")
        self.max_chars = max_chars
        self.used_chars = 0
        self.progress_callback = progress_callback

    def append(self, parts: list[str], text: str) -> None:
        _report_progress(self.progress_callback)
        if not text:
            return
        separator_chars = 1 if parts else 0
        projected = self.used_chars + separator_chars + len(text)
        if projected > self.max_chars:
            raise _OfficeTextLimitExceededError(
                f"Office extracted text exceeded {self.max_chars} characters"
            )
        parts.append(text)
        self.used_chars = projected


def _normalize_text(text: str) -> str:
    """NUL バイト除去 + 連続空白を 1 スペースに圧縮（pdf_extract.py と同じ流儀）."""
    if not text:
        return ""
    # NUL は postgres TEXT で禁止
    text = text.replace("\x00", "")
    text = re.sub(r"\s+", " ", text).strip()
    return text


def extract_docx_text(
    data: bytes,
    *,
    max_chars: int = MAX_OFFICE_EXTRACTED_CHARACTERS,
    progress_callback: Callable[[], None] | None = None,
) -> str:
    """docx → 単一テキスト。paragraph + table セルを順に連結."""
    from docx import Document  # python-docx（lazy import）

    _report_progress(progress_callback)
    document = Document(BytesIO(data))
    _report_progress(progress_callback)
    parts: list[str] = []
    budget = _TextBudget(max_chars, progress_callback)
    for para in document.paragraphs:
        _report_progress(progress_callback)
        if para.text:
            budget.append(parts, para.text)
    # 表セルも本文として拾う（書式は捨てる）
    for table in document.tables:
        _report_progress(progress_callback)
        for row in table.rows:
            _report_progress(progress_callback)
            for cell in row.cells:
                _report_progress(progress_callback)
                if cell.text:
                    budget.append(parts, cell.text)
    return _normalize_text("\n".join(parts))


def _collect_shape_text(
    shapes: object,
    parts: list[str],
    *,
    include_tables: bool,
    recurse_groups: bool,
    budget: _TextBudget,
) -> None:
    """shape 群を走査して本文テキストを ``parts`` へ追記.

    ``include_tables=True`` のとき ``shape.has_table`` のセル文字列も拾う。
    ``recurse_groups=True`` のとき group shape（``shape_type ==
    _GROUP_SHAPE_TYPE``）の子 shape を再帰する。両 flag False のときは
    従来の非再帰・本文 run のみの挙動と完全一致する。
    renderer.py:56-73 の walk 実装を参考にした。
    """
    for shape in shapes:  # type: ignore[attr-defined]
        _report_progress(budget.progress_callback)
        # group shape は（要求時のみ）子 shape を再帰
        if recurse_groups and getattr(shape, "shape_type", None) == _GROUP_SHAPE_TYPE:
            _collect_shape_text(
                shape.shapes,
                parts,
                include_tables=include_tables,
                recurse_groups=recurse_groups,
                budget=budget,
            )
            continue
        # 表セル（include_tables のときだけ）
        if include_tables and getattr(shape, "has_table", False):
            for row in shape.table.rows:
                _report_progress(budget.progress_callback)
                for cell in row.cells:
                    _report_progress(budget.progress_callback)
                    if cell.text:
                        budget.append(parts, cell.text)
            continue
        if getattr(shape, "has_text_frame", False):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        budget.append(parts, run.text)
        elif hasattr(shape, "text") and getattr(shape, "text", None):
            budget.append(parts, shape.text)


def extract_pptx_pages(
    data: bytes,
    *,
    include_notes: bool = False,
    include_tables: bool = False,
    max_chars: int = MAX_OFFICE_EXTRACTED_CHARACTERS,
    progress_callback: Callable[[], None] | None = None,
) -> list[tuple[int, str]]:
    """pptx → slide ごとのテキスト。空 slide は除外.

    既定（両 flag False）では従来挙動（本文 run のみ・非再帰）と完全一致。
    ``include_notes=True`` で各 slide のノート（``slide.notes_slide``）本文を、
    ``include_tables=True`` で表セル文字列を追加で拾う。
    いずれかの flag が True のときだけ group shape を再帰し、子 shape の
    テキスト/表も拾う（従来は group 配下を取りこぼしていた）。
    """
    from pptx import Presentation  # python-pptx（lazy import）

    # 拡張要求があるときだけ group を再帰（既定は従来の非再帰挙動を維持）
    recurse_groups = include_notes or include_tables

    _report_progress(progress_callback)
    prs = Presentation(BytesIO(data))
    _report_progress(progress_callback)
    budget = _TextBudget(max_chars, progress_callback)
    out: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        _collect_shape_text(
            slide.shapes,
            parts,
            include_tables=include_tables,
            recurse_groups=recurse_groups,
            budget=budget,
        )
        if include_notes and slide.has_notes_slide:
            notes = slide.notes_slide
            notes_tf = getattr(notes, "notes_text_frame", None)
            if notes_tf is not None and notes_tf.text:
                budget.append(parts, notes_tf.text)
        text = _normalize_text("\n".join(parts))
        if text:
            out.append((i, text))
    return out


def _collect_sheet_cells(
    ws: object,
    sheet_index: int,
    *,
    budget: _TextBudget,
) -> list[str]:
    """1 worksheet を走査してセル文字列を返す（OOM ガード付き）.

    1 sheet あたり最大 ``_XLSX_MAX_CELLS_PER_SHEET`` セルまで。
    上限超過はpartial本文を返さず、dispatcherで``unsafe_content_volume``に分類する。
    """
    parts: list[str] = []
    cell_count = 0
    for row in ws.iter_rows(values_only=True):  # type: ignore[attr-defined]
        _report_progress(budget.progress_callback)
        for v in row:
            if v is None:
                continue
            s = str(v).strip()
            if not s:
                continue
            if cell_count >= _XLSX_MAX_CELLS_PER_SHEET:
                raise _OfficeTextLimitExceededError(
                    f"xlsx sheet exceeded {_XLSX_MAX_CELLS_PER_SHEET} non-empty cells"
                )
            budget.append(parts, s)
            cell_count += 1
    return parts


def extract_xlsx_pages(
    data: bytes,
    *,
    formula_fallback: bool = False,
    max_chars: int = MAX_OFFICE_EXTRACTED_CHARACTERS,
    progress_callback: Callable[[], None] | None = None,
) -> list[tuple[int, str]]:
    """xlsx → sheet ごとのテキスト。空 sheet・空セルは除外.

    1 sheet あたり最大 ``_XLSX_MAX_CELLS_PER_SHEET`` セルまで（OOM ガード）。
    ``data_only=True`` で式の評価結果を取得（最後に Excel が保存した値）。
    ``read_only=True`` でストリーミング読込（メモリ削減）。

    ``formula_fallback=True`` のとき、``data_only=True`` で全セル None だった
    （＝キャッシュ値を持たない＝openpyxl 等で書かれた未評価の式 sheet）について
    のみ、``data_only=False`` で book を再ロードし式文字列（``=SUM(...)`` 等）を
    拾う。既定 False では従来挙動（キャッシュ値のみ）と完全一致。
    """
    from openpyxl import load_workbook  # lazy import

    _report_progress(progress_callback)
    wb = load_workbook(BytesIO(data), data_only=True, read_only=True)
    _report_progress(progress_callback)
    budget = _TextBudget(max_chars, progress_callback)
    # data_only で空だった sheet を index→title で記録（fallback 用）
    empty_titles: dict[int, str] = {}
    out: list[tuple[int, str]] = []
    for i, ws in enumerate(wb.worksheets, start=1):
        parts = _collect_sheet_cells(ws, i, budget=budget)
        text = _normalize_text(" ".join(parts))
        if text:
            out.append((i, text))
        elif formula_fallback:
            empty_titles[i] = ws.title

    if formula_fallback and empty_titles:
        # data_only=False で式文字列を再取得（None だった sheet のみ）
        wb_raw = load_workbook(BytesIO(data), data_only=False, read_only=True)
        raw_sheets = wb_raw.worksheets
        for i, _title in empty_titles.items():
            ws_raw = raw_sheets[i - 1]  # worksheets 並びは安定
            parts = _collect_sheet_cells(ws_raw, i, budget=budget)
            text = _normalize_text(" ".join(parts))
            if text:
                out.append((i, text))
        # sheet index 昇順に整列（fallback 分を本来の位置へ）
        out.sort(key=lambda pair: pair[0])

    return out


def extract_office_pages(
    data: bytes,
    mime_type: str,
    *,
    expected_size: int | None = None,
    expected_md5: str | None = None,
    include_notes: bool = False,
    include_tables: bool = False,
    formula_fallback: bool = False,
    min_chars: int = 0,
    max_extracted_chars: int = MAX_OFFICE_EXTRACTED_CHARACTERS,
    progress_callback: Callable[[], None] | None = None,
) -> list[tuple[int, str]]:
    """mime_type に応じて office 抽出器を dispatch.

    docx → [(1, text)]（page 概念なし→ 1 ページ扱い）
    pptx → [(slide_num, text), ...]
    xlsx → [(sheet_idx, text), ...]

    pdf_extract.extract_pdf_pages と同じ I/F なので chunk_pages を再利用できる。

    すべての追加 kwarg は後方互換（既定値は現行挙動と完全一致）:

    * ``expected_size`` — Drive metadata の size。実バイト数との差を途中切れとして分類する。
    * ``expected_md5`` — Drive metadata の MD5。本文が metadata と同一 payload か検証する。
    * ``include_notes`` / ``include_tables`` — pptx のノート・表セルも拾う
      （mime が pptx 以外なら無視）。
    * ``formula_fallback`` — xlsx でキャッシュ値が無い式 sheet の式文字列を拾う
      （mime が xlsx 以外なら無視）。
    * ``min_chars`` — ``>0`` のとき、結合テキスト長が ``min_chars`` 未満の
      ページを空扱いで落とす（極小テキストの偽「中身あり」防止）。全 mime 共通。
    * ``max_extracted_chars`` — ファイル全体で保持する抽出本文文字数のhard cap。
    * ``progress_callback`` — ZIP検証・本文走査中に定期的に呼ぶheartbeat hook。

    例外は現行同様にそのまま上位へ伝播（pipeline 側で fail-open）。
    """
    if mime_type not in OFFICE_BINARY_MIMES:
        raise ValueError(f"Unsupported office mime type: {mime_type}")

    _validate_office_payload(
        data,
        mime_type=mime_type,
        expected_size=expected_size,
        expected_md5=expected_md5,
        progress_callback=progress_callback,
    )

    try:
        if mime_type == DOCX_MIME:
            text = extract_docx_text(
                data,
                max_chars=max_extracted_chars,
                progress_callback=progress_callback,
            )
            pages = [(1, text)] if text else []
        elif mime_type == PPTX_MIME:
            pages = extract_pptx_pages(
                data,
                include_notes=include_notes,
                include_tables=include_tables,
                max_chars=max_extracted_chars,
                progress_callback=progress_callback,
            )
        elif mime_type == XLSX_MIME:
            pages = extract_xlsx_pages(
                data,
                formula_fallback=formula_fallback,
                max_chars=max_extracted_chars,
                progress_callback=progress_callback,
            )
    except _OfficeProgressCallbackError as exc:
        raise exc.cause from exc
    except _OfficeTextLimitExceededError as exc:
        raise OfficePayloadError(
            "unsafe_content_volume",
            mime_type=mime_type,
            actual_bytes=len(data),
            expected_bytes=expected_size,
        ) from exc
    if min_chars > 0:
        pages = [(num, text) for num, text in pages if len(text) >= min_chars]
    return pages


__all__ = [
    "DOCX_MIME",
    "GDOC_NATIVE_MIME",
    "MAX_OFFICE_COMPRESSED_BYTES",
    "MAX_OFFICE_EXTRACTED_CHARACTERS",
    "OFFICE_BINARY_MIMES",
    "OFFICE_VALIDATOR_SCHEMA_VERSION",
    "PPTX_MIME",
    "XLSX_MIME",
    "OfficePayloadError",
    "extract_docx_text",
    "extract_office_pages",
    "extract_pptx_pages",
    "extract_xlsx_pages",
]
