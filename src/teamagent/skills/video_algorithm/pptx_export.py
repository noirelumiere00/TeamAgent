"""要点スライドHTML（slides.render_slides）→ PPTX（16:9・1スライド=1画像）。

方式A（MVP）: 各 <section class="slide"> を playwright の要素スクショ（1280x720・2x）で PNG 化し、
python-pptx の add_picture で 16:9 スライドにフルブリード配置する。見た目が slides.py の CSS と
完全一致する（CSS再現の悩みゼロ）。PPTX 上では画像のため文字編集は不可（=Phase2 でネイティブ化）。

Fargate(非root mcp/uid10001)・拡張版イメージ前提:
  - chromium は apt の system 版（CHROMIUM_PATH=/usr/bin/chromium）。executable_path を明示する。
  - 非root + 小さい /dev/shm のため --no-sandbox / --disable-dev-shm-usage / --disable-gpu。
  - 日本語は fonts-noto-cjk（Dockerfile 拡張版）で豆腐回避。

shoot_sections は playwright 注入でテスト差し替え可＝chromium 無しでも build_pptx を検証できる。
"""

from __future__ import annotations

import io
import os
from collections.abc import Callable

import structlog

from teamagent.skills.video_algorithm.schema import VideoAlgorithmOutput
from teamagent.skills.video_algorithm.slides import SLIDE_H, SLIDE_W, render_slides

logger = structlog.get_logger(__name__)

# 各 .slide section を撮った PNG bytes のリストを返す関数の型（テスト差し替え用）。
Shooter = Callable[[str], list[bytes]]

_DEVICE_SCALE = 2  # 2x で 2560x1440 PNG（鮮明・PPTX は数MBに収まる。3x は容量過大で避ける）
_LAUNCH_ARGS = ["--no-sandbox", "--disable-dev-shm-usage", "--disable-gpu"]


def shoot_sections(html: str) -> list[bytes]:
    """slides HTML を描画し、各 .slide section を要素単位で PNG 化（full_page にしない）。

    full_page だと全スライド連結の縦長1枚になるため、locator('.slide').nth(i).screenshot() で
    1枚ずつ撮る。.slide は 1280x720 固定なのでアスペクトずれが出ない。
    """
    from playwright.sync_api import sync_playwright

    exe = os.environ.get("CHROMIUM_PATH")  # 拡張版イメージの system chromium
    pngs: list[bytes] = []
    with sync_playwright() as p:
        browser = p.chromium.launch(executable_path=exe or None, args=_LAUNCH_ARGS)
        try:
            page = browser.new_page(
                viewport={"width": SLIDE_W, "height": SLIDE_H},
                device_scale_factor=_DEVICE_SCALE,
            )
            page.set_content(html, wait_until="networkidle")
            slides = page.locator(".slide")
            for i in range(slides.count()):
                pngs.append(slides.nth(i).screenshot(type="png"))
        finally:
            browser.close()
    return pngs


def build_pptx(pngs: list[bytes], out_path: str) -> str:
    """PNG bytes のリスト → 16:9 PPTX（1スライド=1画像・フルブリード）。失敗時は例外。"""
    from pptx import Presentation
    from pptx.util import Emu, Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9（python-pptx が EMU へ自動変換）
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # 完全な空レイアウト
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            io.BytesIO(png), Emu(0), Emu(0), width=prs.slide_width, height=prs.slide_height
        )
    prs.save(out_path)
    return out_path


def render_pptx(
    out: VideoAlgorithmOutput,
    out_path: str,
    *,
    generated_at: str = "",
    shooter: Shooter | None = None,
) -> str | None:
    """VideoAlgorithmOutput → 要点スライドHTML → PPTX を out_path に書く。失敗で None（graceful）。

    shooter 注入時はそれを使う（テスト/将来の weasyprint 代替）。無ければ playwright を使う。
    """
    try:
        html = render_slides(out, generated_at=generated_at)
        shoot = shooter or shoot_sections
        pngs = shoot(html)
        if not pngs:
            logger.warning("vseo_pptx_no_sections")
            return None
        return build_pptx(pngs, out_path)
    except Exception as e:
        logger.warning("vseo_pptx_failed", error=type(e).__name__)
        return None
