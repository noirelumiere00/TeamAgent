"""python-pptx 経由のテンプレート流し込みレンダラ（proposal_deck Skill 専用）。

teamagent_consulting の rendering/pptx_renderer.py + agents/slidewriter.py を移植。

placeholder 表記: 全角「｛N：ラベル｝」「｛N｝」/ 半角「{N:ラベル}」「{N}」（N は十進整数）。

**段落跨ぎ対応**: 実テンプレ（界隈マトリクス）では 1 placeholder が複数段落 (<a:p>) に
分割される。テキストフレーム単位で全段落を連結して置換・監査し、書き戻し時は段落構造を保つ
（マッチ開始段落に値を入れ、消費された段落は空にする）。書式損失は許容（自動生成テキスト主役）。
"""

from __future__ import annotations

import io
import re
from collections.abc import Iterator
from dataclasses import dataclass, field
from datetime import date, timedelta
from pathlib import Path
from typing import Any

from teamagent.skills.proposal_deck.contract import (
    PROPOSAL_BUILDER_REQUIRED_AUXILIARY,
    PROPOSAL_BUILDER_TEMPLATE_PROFILE,
    VALID_IDS,
    ComposerOutput,
    EvidenceImage,
)

# 全角・半角両対応。{N:ラベル} の `N` だけをキャプチャ。`[^｝}]*` は改行も含むため
# テキストフレーム単位で連結すれば段落跨ぎトークンも 1 マッチになる。
PLACEHOLDER_PATTERN = re.compile(r"[｛{]\s*(\d+)\s*[:：]?[^｝}]*[｝}]")
AUXILIARY_PATTERN = re.compile(r"\{\{(PB-[A-Z0-9_-]{1,60})\}\}")
DATE_PATTERN = re.compile(r"\{\{PB-DATE:([+-]?\d{1,3}):(%Y/%m/%d|%m/%d|%Y年%m月%d日)\}\}")
TEMPLATE_VERSION_PATTERN = re.compile(r"\{\{PB-TEMPLATE:([a-z0-9-]{1,40})\}\}")
PB_TOKEN_PATTERN = re.compile(r"\{\{PB-[^{}]+\}\}")
LEGACY_INSTRUCTION_PATTERN = re.compile(r"自動入力|貼り付けてください|はめ込|転記|差し替え")
BRACE_CHARACTER_PATTERN = re.compile(r"[{}｛｝]")
_PROPOSAL_BUILDER_REQUIRED_DATE_OFFSETS = frozenset(range(-56, 22, 7))
_PROPOSAL_BUILDER_EXPECTED_SLIDE_COUNT = 83

# python-pptx の GROUP shape type 値（MSO_SHAPE_TYPE.GROUP == 6）。
_GROUP_SHAPE_TYPE = 6

# template_v2.pptx の {58-92} マトリクス最下段サムネ空枠（EMU・1inch=914400）。
# 0.67"角(=609905)・y≈5.98"(=5454720)。座標一致の許容誤差 ±80000 EMU。
# NOTE: テンプレ固有値。別テンプレ採用時はスロット仕様の拡張が必要（座標未一致なら静かに skip）。
_PICTURE_SHAPE_TYPE = 13  # MSO_SHAPE_TYPE.PICTURE（既存スロット除外用）
_SLOT_TOP_EMU = 5454720
_SLOT_SIZE_EMU = 609905
_COORD_TOL_EMU = 80000


class UnfilledPlaceholderError(RuntimeError):
    """fail_if_missing=True で render_pptx を呼んで残存 placeholder があった。"""


@dataclass(slots=True)
class AuditResult:
    """audit_presentation の戻り値。"""

    filled_ids: frozenset[int]
    unfilled_ids: list[int]
    extra_braces: list[str]
    missing_template_ids: list[int] = field(default_factory=list)
    unfilled_auxiliary: list[str] = field(default_factory=list)
    unfilled_dates: list[str] = field(default_factory=list)
    legacy_artifacts: list[str] = field(default_factory=list)

    @property
    def is_clean(self) -> bool:
        return not (
            self.unfilled_ids
            or self.extra_braces
            or self.missing_template_ids
            or self.unfilled_auxiliary
            or self.unfilled_dates
            or self.legacy_artifacts
        )


@dataclass(slots=True, frozen=True)
class TemplateInventory:
    """置換前テンプレートに実在した動的トークンの一覧。"""

    numeric_ids: frozenset[int]
    auxiliary_keys: frozenset[str]
    date_token_count: int
    slide_count: int
    date_offsets: frozenset[int]
    template_versions: frozenset[str]
    legacy_artifacts: tuple[str, ...]


def _iter_text_frames(prs: Any) -> Iterator[Any]:
    """全 slide の全 shape（group / table 内も再帰）からテキストフレームを yield。"""

    def walk(shapes: Any) -> Iterator[Any]:
        for shape in shapes:
            if getattr(shape, "shape_type", None) == _GROUP_SHAPE_TYPE:
                yield from walk(shape.shapes)
                continue
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        yield cell.text_frame
                continue
            if getattr(shape, "has_text_frame", False):
                yield shape.text_frame

    for slide in prs.slides:
        yield from walk(slide.shapes)


def _emit_range(
    cat: str, a: int, b: int, bounds: list[tuple[int, int]], pieces: list[list[str]]
) -> None:
    """連結文字列 cat[a:b] を、出現元の段落ごとに pieces へ振り分ける。"""
    if a >= b:
        return
    for i, (s, e) in enumerate(bounds):
        lo = max(a, s)
        hi = min(b, e)
        if lo < hi:
            pieces[i].append(cat[lo:hi])


def _set_paragraph_text(paragraph: Any, text: str) -> None:
    """段落の text を再設定。書式は先頭 run のものに統一。"""
    runs = list(paragraph.runs)
    if runs:
        runs[0].text = text
        for r in runs[1:]:
            r.text = ""
    elif text:
        paragraph.add_run().text = text


def _replace_in_text_frame(text_frame: Any, placeholders: dict[int, str]) -> set[int]:
    """1 text frame の全段落を連結 → placeholder 置換 → 段落構造を保って書き戻し。"""
    paragraphs = list(text_frame.paragraphs)
    texts = [p.text for p in paragraphs]
    bounds: list[tuple[int, int]] = []
    pos = 0
    for t in texts:
        bounds.append((pos, pos + len(t)))
        pos += len(t)
    cat = "".join(texts)
    if not cat:
        return set()

    matches = [m for m in PLACEHOLDER_PATTERN.finditer(cat) if int(m.group(1)) in placeholders]
    if not matches:
        return set()

    def para_of(p: int) -> int:
        for i, (s, e) in enumerate(bounds):
            if s <= p < e:
                return i
        return len(bounds) - 1

    filled: set[int] = set()
    pieces: list[list[str]] = [[] for _ in texts]
    cursor = 0
    for m in matches:
        _emit_range(cat, cursor, m.start(), bounds, pieces)
        pid = int(m.group(1))
        pieces[para_of(m.start())].append(placeholders[pid])
        filled.add(pid)
        cursor = m.end()
    _emit_range(cat, cursor, len(cat), bounds, pieces)

    new_texts = ["".join(pp) for pp in pieces]
    for i, paragraph in enumerate(paragraphs):
        if new_texts[i] != texts[i]:
            _set_paragraph_text(paragraph, new_texts[i])
    return filled


def _replace_pattern_in_text_frame(
    text_frame: Any,
    pattern: re.Pattern[str],
    resolve: Any,
) -> set[str]:
    """段落跨ぎトークンを ``resolve(match)`` の値へ置換する。"""

    paragraphs = list(text_frame.paragraphs)
    texts = [p.text for p in paragraphs]
    bounds: list[tuple[int, int]] = []
    pos = 0
    for text in texts:
        bounds.append((pos, pos + len(text)))
        pos += len(text)
    combined = "".join(texts)
    if not combined:
        return set()

    resolved: list[tuple[re.Match[str], str, str]] = []
    for match in pattern.finditer(combined):
        replacement = resolve(match)
        if replacement is not None:
            resolved.append((match, str(replacement), match.group(0)))
    if not resolved:
        return set()

    def paragraph_of(position: int) -> int:
        for index, (start, end) in enumerate(bounds):
            if start <= position < end:
                return index
        return len(bounds) - 1

    replaced: set[str] = set()
    pieces: list[list[str]] = [[] for _ in texts]
    cursor = 0
    for match, replacement, token in resolved:
        _emit_range(combined, cursor, match.start(), bounds, pieces)
        pieces[paragraph_of(match.start())].append(replacement)
        replaced.add(token)
        cursor = match.end()
    _emit_range(combined, cursor, len(combined), bounds, pieces)

    new_texts = ["".join(piece) for piece in pieces]
    for index, paragraph in enumerate(paragraphs):
        if new_texts[index] != texts[index]:
            _set_paragraph_text(paragraph, new_texts[index])
    return replaced


def inspect_template(prs: Any) -> TemplateInventory:
    """置換前の数値・補助・日付トークンとslide数を収集する。"""

    numeric_ids: set[int] = set()
    auxiliary_keys: set[str] = set()
    date_token_count = 0
    date_offsets: set[int] = set()
    template_versions: set[str] = set()
    legacy_artifacts: list[str] = []
    for text_frame in _iter_text_frames(prs):
        combined = "".join(paragraph.text for paragraph in text_frame.paragraphs)
        numeric_ids.update(int(match.group(1)) for match in PLACEHOLDER_PATTERN.finditer(combined))
        auxiliary_keys.update(match.group(1) for match in AUXILIARY_PATTERN.finditer(combined))
        date_token_count += sum(1 for _ in DATE_PATTERN.finditer(combined))
        date_offsets.update(int(match.group(1)) for match in DATE_PATTERN.finditer(combined))
        template_versions.update(
            match.group(1) for match in TEMPLATE_VERSION_PATTERN.finditer(combined)
        )
        legacy_artifacts.extend(_find_legacy_artifacts(combined, allow_template_tokens=True))
    return TemplateInventory(
        numeric_ids=frozenset(numeric_ids),
        auxiliary_keys=frozenset(auxiliary_keys),
        date_token_count=date_token_count,
        slide_count=len(prs.slides),
        date_offsets=frozenset(date_offsets),
        template_versions=frozenset(template_versions),
        legacy_artifacts=tuple(dict.fromkeys(legacy_artifacts))[:20],
    )


def _find_legacy_artifacts(
    text: str,
    *,
    allow_template_tokens: bool,
) -> list[str]:
    """Find obsolete brace tokens and manual operator instructions.

    The proposal-builder profile accepts only the explicit numeric/PB/date/version
    token grammar.  Everything else from the Windows FMT (for example
    ``｛自動入力：…｝`` or ``｛A-1｝``) is rejected before and after rendering.
    """

    scrubbed = text
    if allow_template_tokens:
        for pattern in (
            PLACEHOLDER_PATTERN,
            DATE_PATTERN,
            TEMPLATE_VERSION_PATTERN,
            AUXILIARY_PATTERN,
        ):
            scrubbed = pattern.sub("", scrubbed)
    findings: list[str] = []
    brace = BRACE_CHARACTER_PATTERN.search(scrubbed)
    if brace is not None:
        start = max(0, brace.start() - 40)
        findings.append(scrubbed[start : brace.start() + 120])
    findings.extend(match.group(0) for match in LEGACY_INSTRUCTION_PATTERN.finditer(scrubbed))
    return findings


def _replace_auxiliary(
    prs: Any,
    auxiliary_placeholders: dict[str, str],
) -> None:
    for text_frame in _iter_text_frames(prs):
        _replace_pattern_in_text_frame(
            text_frame,
            AUXILIARY_PATTERN,
            lambda match: auxiliary_placeholders.get(match.group(1)),
        )


def _replace_schedule_dates(prs: Any, posting_start_date: date | None) -> None:
    if posting_start_date is None:
        return

    def resolve(match: re.Match[str]) -> str:
        offset = int(match.group(1))
        return (posting_start_date + timedelta(days=offset)).strftime(match.group(2))

    for text_frame in _iter_text_frames(prs):
        _replace_pattern_in_text_frame(text_frame, DATE_PATTERN, resolve)


def _replace_template_version(prs: Any, template_profile: str) -> None:
    for text_frame in _iter_text_frames(prs):
        _replace_pattern_in_text_frame(
            text_frame,
            TEMPLATE_VERSION_PATTERN,
            lambda match: "" if match.group(1) == template_profile else None,
        )


def _resolve_evidence_image_bytes(img: EvidenceImage) -> bytes | None:
    """EvidenceImage を画像バイトに解決。image_path のみ参照（httpx/fetch 非依存）。

    image_path が無い／読めない／source_url のみ → None（graceful skip）。
    フィーダ（Phase4）が正規化済み JPEG を image_path に用意する前提。
    """
    if not img.image_path:
        return None
    try:
        data = Path(img.image_path).read_bytes()
    except OSError:
        return None
    return data or None


def _iter_image_slots(prs: Any) -> Iterator[tuple[Any, Any]]:
    """全 slide から「空サムネ枠」候補を (slide, shape) で yield（group 再帰）。

    座標(top≈_SLOT_TOP_EMU)・サイズ(width≈_SLOT_SIZE_EMU)が許容誤差内の shape を
    候補とし、PICTURE は除外。slide ごとに left 昇順で返す。
    """

    def walk(shapes: Any) -> Iterator[Any]:
        for shape in shapes:
            if getattr(shape, "shape_type", None) == _GROUP_SHAPE_TYPE:
                yield from walk(shape.shapes)
                continue
            yield shape

    for slide in prs.slides:
        candidates: list[Any] = []
        for shape in walk(slide.shapes):
            if getattr(shape, "shape_type", None) == _PICTURE_SHAPE_TYPE:
                continue
            try:
                w = int(shape.width)
                t = int(shape.top)
            except (TypeError, ValueError):
                continue
            if (
                abs(w - _SLOT_SIZE_EMU) <= _COORD_TOL_EMU
                and abs(t - _SLOT_TOP_EMU) <= _COORD_TOL_EMU
            ):
                candidates.append(shape)
        for shape in sorted(candidates, key=lambda s: int(s.left)):
            yield slide, shape


def _add_picture_fit(slide: Any, img_bytes: bytes, shape: Any) -> None:
    """空枠の高さに等比で合わせて add_picture（歪み無し）。左上は枠に合わせる。"""
    from pptx.util import Emu

    left, top, box_h = int(shape.left), int(shape.top), int(shape.height)
    pic = slide.shapes.add_picture(io.BytesIO(img_bytes), Emu(left), Emu(top))
    if pic.width and pic.height:
        scale = box_h / pic.height
        pic.height = Emu(int(box_h))
        pic.width = Emu(int(pic.width * scale))
        pic.left = Emu(left)
        pic.top = Emu(top)


def _inject_evidence_images(prs: Any, evidence_images: dict[int, list[EvidenceImage]]) -> int:
    """evidence_images を placeholder_id 昇順→rank 昇順で空枠へ順に add_picture。

    image_path を解決できない画像はスキップ。空枠が尽きたら以降は捨てる
    （pptx はテキスト主役・graceful by design）。注入した枚数を返す。
    add_picture が画像形式を認識できない（壊れたファイル等）場合もスキップする。
    """
    from PIL import UnidentifiedImageError  # python-pptx は Pillow backend

    ordered: list[EvidenceImage] = []
    for pid in sorted(evidence_images):
        for img in sorted(evidence_images[pid], key=lambda e: e.rank):
            ordered.append(img)
    if not ordered:
        return 0

    slots = _iter_image_slots(prs)
    injected = 0
    for img in ordered:
        data = _resolve_evidence_image_bytes(img)
        if data is None:
            continue
        try:
            slide, shape = next(slots)
        except StopIteration:
            break
        try:
            _add_picture_fit(slide, data, shape)
        except UnidentifiedImageError:
            continue
        injected += 1
    return injected


def render_pptx(
    template_path: Path,
    placeholders: dict[int, str],
    out_path: Path,
    *,
    fail_if_missing: bool = True,
    evidence_images: dict[int, list[EvidenceImage]] | None = None,
    auxiliary_placeholders: dict[str, str] | None = None,
    posting_start_date: date | None = None,
    template_profile: str = "base",
) -> Path:
    """テンプレ pptx を読み、placeholder を全置換して out_path に保存。

    evidence_images（任意）が渡されたら、テキスト置換・audit 後・保存前に空枠へ
    画像を add_picture する（image_path のみ解決・httpx 非依存）。
    """
    from pptx import Presentation

    prs = Presentation(str(template_path))
    inventory = inspect_template(prs)
    required_template_ids: set[int] = set()
    if template_profile == PROPOSAL_BUILDER_TEMPLATE_PROFILE:
        required_template_ids = set(VALID_IDS)
        missing_auxiliary_values = PROPOSAL_BUILDER_REQUIRED_AUXILIARY - set(
            auxiliary_placeholders or {}
        )
        invalid_auxiliary_values = (
            set(auxiliary_placeholders or {}) != PROPOSAL_BUILDER_REQUIRED_AUXILIARY
        )
        invalid_auxiliary_tokens = (
            set(inventory.auxiliary_keys) != PROPOSAL_BUILDER_REQUIRED_AUXILIARY
        )
        if missing_auxiliary_values or invalid_auxiliary_values:
            raise UnfilledPlaceholderError(
                "proposal-builder auxiliary values must exactly match the integrated "
                f"template contract; missing={sorted(missing_auxiliary_values)}"
            )
        if invalid_auxiliary_tokens:
            raise UnfilledPlaceholderError(
                "integrated template auxiliary token inventory must exactly match "
                "the proposal-builder contract"
            )
        if inventory.legacy_artifacts:
            raise UnfilledPlaceholderError(
                "integrated template contains legacy placeholders or operator "
                f"instructions: {list(inventory.legacy_artifacts[:5])}"
            )
        if inventory.numeric_ids != VALID_IDS:
            raise UnfilledPlaceholderError(
                "integrated template numeric inventory must exactly match the 95-ID contract"
            )
        if inventory.slide_count != _PROPOSAL_BUILDER_EXPECTED_SLIDE_COUNT:
            raise UnfilledPlaceholderError("integrated template must contain exactly 83 slides")
        if inventory.template_versions != {PROPOSAL_BUILDER_TEMPLATE_PROFILE}:
            raise UnfilledPlaceholderError(
                "integrated template version marker is missing or ambiguous"
            )
        if posting_start_date is None:
            raise UnfilledPlaceholderError("proposal-builder posting_start_date is required")
        if not _PROPOSAL_BUILDER_REQUIRED_DATE_OFFSETS.issubset(inventory.date_offsets):
            raise UnfilledPlaceholderError(
                "integrated template D-56 through D+21 schedule markers are incomplete"
            )

    filled: set[int] = set()
    for tf in _iter_text_frames(prs):
        filled |= _replace_in_text_frame(tf, placeholders)
    _replace_auxiliary(prs, auxiliary_placeholders or {})
    _replace_schedule_dates(prs, posting_start_date)
    _replace_template_version(prs, template_profile)

    audit = audit_presentation(
        prs,
        expected_filled_ids=set(placeholders),
        observed_template_ids=set(inventory.numeric_ids),
        required_template_ids=required_template_ids,
        reject_legacy_artifacts=(template_profile == PROPOSAL_BUILDER_TEMPLATE_PROFILE),
    )
    if fail_if_missing and not audit.is_clean:
        raise UnfilledPlaceholderError(
            f"unfilled placeholders remain: {audit.unfilled_ids[:20]} "
            f"missing template ids: {audit.missing_template_ids[:20]} "
            f"auxiliary: {audit.unfilled_auxiliary[:5]} "
            f"dates: {audit.unfilled_dates[:5]} "
            f"legacy: {audit.legacy_artifacts[:5]} "
            f"(samples: {audit.extra_braces[:5]})"
        )

    if evidence_images:
        _inject_evidence_images(prs, evidence_images)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def audit_presentation(
    prs: Any,
    *,
    expected_filled_ids: set[int] | None = None,
    observed_template_ids: set[int] | None = None,
    required_template_ids: set[int] | None = None,
    reject_legacy_artifacts: bool = False,
) -> AuditResult:
    """全 slide を走査し、残存 `{N}` を検出する（テキストフレーム単位 = 段落跨ぎ対応）。"""
    unfilled: set[int] = set()
    extra: list[str] = []
    unfilled_auxiliary: set[str] = set()
    unfilled_dates: set[str] = set()
    legacy_artifacts: list[str] = []
    for tf in _iter_text_frames(prs):
        cat = "".join(p.text for p in tf.paragraphs)
        for m in PLACEHOLDER_PATTERN.finditer(cat):
            unfilled.add(int(m.group(1)))
            if len(extra) < 20:
                extra.append(m.group(0))
        unfilled_auxiliary.update(match.group(0) for match in PB_TOKEN_PATTERN.finditer(cat))
        unfilled_dates.update(match.group(0) for match in DATE_PATTERN.finditer(cat))
        if reject_legacy_artifacts:
            legacy_artifacts.extend(_find_legacy_artifacts(cat, allow_template_tokens=False))
    filled = (expected_filled_ids or set()) - unfilled
    missing_template_ids = (required_template_ids or set()) - (observed_template_ids or set())
    return AuditResult(
        filled_ids=frozenset(filled),
        unfilled_ids=sorted(unfilled),
        extra_braces=extra,
        missing_template_ids=sorted(missing_template_ids),
        unfilled_auxiliary=sorted(unfilled_auxiliary),
        unfilled_dates=sorted(unfilled_dates),
        legacy_artifacts=list(dict.fromkeys(legacy_artifacts))[:20],
    )


def materialize_placeholders(out: ComposerOutput) -> dict[int, str]:
    """skipped も『要確認（データ未検出）』で実体化し、全 VALID_IDS をカバーする。

    renderer 段階では unfilled を残さないため、未検出 placeholder にも必ず文字列を割り当てる。
    """
    materialized = dict(out.placeholders)
    for skip in out.skipped_placeholders:
        materialized.setdefault(skip.id, "要確認（データ未検出）")
    for pid in VALID_IDS:
        materialized.setdefault(pid, "要確認（データ未検出）")
    return materialized


def render_deck(
    composer_out: ComposerOutput,
    template_path: Path,
    out_path: Path,
    *,
    fail_if_missing: bool = True,
    enable_images: bool = False,
) -> Path:
    """ComposerOutput を pptx に流し込み、ファイルパスを返す。

    enable_images=True のとき composer_out.evidence_images を空枠へ画像注入する
    （既定 False＝従来どおりテキストのみ・後方互換）。
    """
    placeholders = materialize_placeholders(composer_out)
    return render_pptx(
        template_path,
        placeholders,
        out_path,
        fail_if_missing=fail_if_missing,
        evidence_images=composer_out.evidence_images if enable_images else None,
        auxiliary_placeholders=composer_out.auxiliary_placeholders,
        posting_start_date=composer_out.posting_start_date,
        template_profile=composer_out.template_profile,
    )


__all__ = [
    "AUXILIARY_PATTERN",
    "BRACE_CHARACTER_PATTERN",
    "DATE_PATTERN",
    "LEGACY_INSTRUCTION_PATTERN",
    "PB_TOKEN_PATTERN",
    "PLACEHOLDER_PATTERN",
    "TEMPLATE_VERSION_PATTERN",
    "AuditResult",
    "TemplateInventory",
    "UnfilledPlaceholderError",
    "audit_presentation",
    "inspect_template",
    "materialize_placeholders",
    "render_deck",
    "render_pptx",
]
