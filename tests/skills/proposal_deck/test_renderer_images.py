"""renderer の evidence_images 画像注入（Phase3）の単体テスト。実テンプレ(89MB)不使用。

合成方針:
- python-pptx で「既知スロット座標(top=_SLOT_TOP_EMU, 幅/高=_SLOT_SIZE_EMU)に空 AUTO_SHAPE を
  置いた最小テンプレ」を作る → _iter_image_slots が同ロジックで発見できる。
- テスト画像は **埋め込み最小 JPEG バイト**（PIL/ffmpeg を使わない＝CI の ci.yml 手動列挙で
  Pillow 未導入でも import で落ちない／既知の罠回避）。
- 検証は保存後 pptx を開き直し PICTURE 形状を計数（MVP _verify_picture と同パターン）。
"""

from __future__ import annotations

import base64
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Emu, Inches

from teamagent.skills.proposal_deck.contract import (
    LENGTH_RULES,
    VALID_IDS,
    ComposerOutput,
    EvidenceImage,
)
from teamagent.skills.proposal_deck.renderer import (
    _SLOT_SIZE_EMU,
    _SLOT_TOP_EMU,
    render_deck,
    render_pptx,
)

_SLOT_LEFT_EMU = 914400

# 1x1 の最小有効 JPEG（base64）。python-pptx の add_picture が寸法を読める。
_TINY_JPEG = base64.b64decode(
    "/9j/4AAQSkZJRgABAQEAYABgAAD/2wBDAAMCAgICAgMCAgIDAwMDBAYEBAQEBAgGBgUGCQgKCgkICQkKDA8MCgsO"
    "CwkJDRENDg8QEBEQCgwSExIQEw8QEBD/wAALCAABAAEBAREA/8QAFAABAAAAAAAAAAAAAAAAAAAACP/EABQQAQAA"
    "AAAAAAAAAAAAAAAAAAD/2gAIAQEAAD8AfwD/2Q=="
)


def _full_placeholders() -> dict[int, str]:
    ph: dict[int, str] = {}
    for pid in sorted(VALID_IDS):
        if pid in LENGTH_RULES:
            lo, hi = LENGTH_RULES[pid]
            ph[pid] = "サ" * ((lo + hi) // 2)
        else:
            ph[pid] = f"値-{pid}"
    return ph


def _composer(evidence: dict[int, list[EvidenceImage]] | None = None) -> ComposerOutput:
    return ComposerOutput(placeholders=_full_placeholders(), evidence_images=evidence or {})


def _template_with_slots(path: Path, n_slots: int = 2) -> Path:
    """全95枠テキスト + 既知スロット座標に空 AUTO_SHAPE を n 個置いた合成テンプレ。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    ids = sorted(VALID_IDS)
    slide = None
    for i in range(0, len(ids), 30):
        slide = prs.slides.add_slide(blank)
        tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.0)).text_frame
        tf.word_wrap = True
        for pid in ids[i : i + 30]:
            tf.add_paragraph().text = f"｛{pid}：ラベル{pid}｝"
    assert slide is not None
    for k in range(n_slots):
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(_SLOT_LEFT_EMU + k * (_SLOT_SIZE_EMU + 50000)),
            Emu(_SLOT_TOP_EMU),
            Emu(_SLOT_SIZE_EMU),
            Emu(_SLOT_SIZE_EMU),
        )
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def _make_jpeg(path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(_TINY_JPEG)
    return path


def _count_pictures(pptx_path: Path) -> int:
    prs = Presentation(str(pptx_path))
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                n += 1
    return n


def test_enable_images_false_is_default(tmp_path: Path) -> None:
    """enable_images 省略=False。evidence があっても PICTURE 0（後方互換）。"""
    tpl = _template_with_slots(tmp_path / "t.pptx")
    jpg = _make_jpeg(tmp_path / "a.jpg")
    ev = {92: [EvidenceImage(placeholder_id=92, rank=1, keyword="k", image_path=str(jpg))]}
    out = render_deck(_composer(ev), tpl, tmp_path / "out.pptx", fail_if_missing=False)
    assert Path(out).exists()
    assert _count_pictures(Path(out)) == 0


def test_enable_images_true_empty_images(tmp_path: Path) -> None:
    """enable_images=True でも evidence_images={} なら PICTURE 0。"""
    tpl = _template_with_slots(tmp_path / "t.pptx")
    out = render_deck(
        _composer({}), tpl, tmp_path / "out.pptx", fail_if_missing=False, enable_images=True
    )
    assert _count_pictures(Path(out)) == 0


def test_enable_images_true_with_image_path(tmp_path: Path) -> None:
    """image_path 付き 1枚 → PICTURE 1枚注入される。"""
    tpl = _template_with_slots(tmp_path / "t.pptx", n_slots=2)
    jpg = _make_jpeg(tmp_path / "a.jpg")
    ev = {92: [EvidenceImage(placeholder_id=92, rank=1, keyword="集中", image_path=str(jpg))]}
    out = render_deck(
        _composer(ev), tpl, tmp_path / "out.pptx", fail_if_missing=False, enable_images=True
    )
    assert _count_pictures(Path(out)) == 1


def test_graceful_skip_source_url_only(tmp_path: Path) -> None:
    """source_url のみ（image_path=None）→ fetch せず skip・PICTURE 0・no crash。"""
    tpl = _template_with_slots(tmp_path / "t.pptx")
    ev = {
        92: [
            EvidenceImage(
                placeholder_id=92, rank=1, keyword="k", source_url="https://cdn.example/x.jpg"
            )
        ]
    }
    out = render_deck(
        _composer(ev), tpl, tmp_path / "out.pptx", fail_if_missing=False, enable_images=True
    )
    assert _count_pictures(Path(out)) == 0


def test_graceful_skip_nonexistent_path(tmp_path: Path) -> None:
    """image_path が存在しないパス → 読めず skip・PICTURE 0・no crash。"""
    tpl = _template_with_slots(tmp_path / "t.pptx")
    ev = {92: [EvidenceImage(placeholder_id=92, rank=1, keyword="k", image_path="/no/such.jpg")]}
    out = render_deck(
        _composer(ev), tpl, tmp_path / "out.pptx", fail_if_missing=False, enable_images=True
    )
    assert _count_pictures(Path(out)) == 0


def test_multiple_images_by_rank(tmp_path: Path) -> None:
    """同一 placeholder に rank=1,2 → 2枠あれば両方 add_picture。"""
    tpl = _template_with_slots(tmp_path / "t.pptx", n_slots=2)
    j1 = _make_jpeg(tmp_path / "1.jpg")
    j2 = _make_jpeg(tmp_path / "2.jpg")
    ev = {
        92: [
            EvidenceImage(placeholder_id=92, rank=1, keyword="k", image_path=str(j1)),
            EvidenceImage(placeholder_id=92, rank=2, keyword="k", image_path=str(j2)),
        ]
    }
    out = render_deck(
        _composer(ev), tpl, tmp_path / "out.pptx", fail_if_missing=False, enable_images=True
    )
    assert _count_pictures(Path(out)) == 2


def test_render_pptx_evidence_images_param(tmp_path: Path) -> None:
    """低レベル render_pptx に evidence_images を直接渡しても注入される。"""
    tpl = _template_with_slots(tmp_path / "t.pptx")
    jpg = _make_jpeg(tmp_path / "a.jpg")
    ph = _full_placeholders()
    ev = {92: [EvidenceImage(placeholder_id=92, rank=1, keyword="k", image_path=str(jpg))]}
    out = render_pptx(tpl, ph, tmp_path / "out.pptx", fail_if_missing=False, evidence_images=ev)
    assert _count_pictures(Path(out)) == 1


def test_graceful_skip_malformed_image(tmp_path: Path) -> None:
    """image_path が画像でない（壊れたバイト）→ add_picture skip・PICTURE 0・no crash。"""
    tpl = _template_with_slots(tmp_path / "t.pptx")
    bad = tmp_path / "bad.jpg"
    bad.write_bytes(b"not a real image at all")
    ev = {92: [EvidenceImage(placeholder_id=92, rank=1, keyword="k", image_path=str(bad))]}
    out = render_deck(
        _composer(ev), tpl, tmp_path / "out.pptx", fail_if_missing=False, enable_images=True
    )
    assert _count_pictures(Path(out)) == 0
