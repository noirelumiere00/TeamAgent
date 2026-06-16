"""Batch C3: proposal_deck の PDF コンパニオン + version_id 配線を検証する。

build_proposal_html は純関数として直接検証。weasyprint(_html_to_pdf) は monkeypatch して
CI に重い C 依存を持ち込まない。emit_pdf 既定 OFF＝従来挙動（PDF なし）も固定する。
"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock

import pytest

from teamagent.adapters.bedrock_client import ConverseResponse, TokenUsage
from teamagent.skills.base import SkillContext
from teamagent.skills.proposal_deck.contract import (
    LENGTH_RULES,
    VALID_IDS,
    ComposerOutput,
    SkippedPlaceholder,
)
from teamagent.skills.proposal_deck.pdf_export import build_proposal_html
from teamagent.skills.proposal_deck.schema import ProposalDeckInput
from teamagent.skills.proposal_deck.skill import ProposalDeckSkill


def _full_composer_json() -> str:
    placeholders: dict[int, str] = {}
    for pid in sorted(VALID_IDS):
        if pid in LENGTH_RULES:
            lo, hi = LENGTH_RULES[pid]
            placeholders[pid] = "サ" * ((lo + hi) // 2)
        else:
            placeholders[pid] = f"値-{pid}"
    return ComposerOutput(placeholders=placeholders).model_dump_json()


def _dummy_template(path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]
    ids = sorted(VALID_IDS)
    for i in range(0, len(ids), 30):
        slide = prs.slides.add_slide(blank)
        tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5)).text_frame
        tf.word_wrap = True
        for pid in ids[i : i + 30]:
            tf.add_paragraph().text = f"｛{pid}：ラベル{pid}｝"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def _resp(text: str) -> ConverseResponse:
    return ConverseResponse(
        text=text,
        usage=TokenUsage(
            input_tokens=100,
            output_tokens=50,
            cache_read_input_tokens=0,
            cache_creation_input_tokens=0,
            cost_usd=0.01,
        ),
        model_id="jp.anthropic.claude-sonnet-4-6",
        latency_ms=100,
        stop_reason="end_turn",
    )


def _input(template: Path, out: Path, **kw: object) -> ProposalDeckInput:
    return ProposalDeckInput(
        product_name="ACME 青汁",
        goal="認知獲得",
        target_persona="20代女性",
        template_path=str(template),
        out_dir=str(out),
        **kw,  # type: ignore[arg-type]
    )


# -----------------------------------------------------------
# build_proposal_html（純関数）
# -----------------------------------------------------------
def test_build_proposal_html_includes_content_and_version() -> None:
    co = ComposerOutput.model_construct(
        placeholders={1: "本文サンプル", 47: "PRワード案"},
        citations_per_placeholder={1: ["src.pdf"]},
        skipped_placeholders=[SkippedPlaceholder(id=2, reason="要確認（データ未検出）")],
    )
    html = build_proposal_html(co, product_name="ACME 青汁", version_id="v-abc123")
    assert "ACME 青汁 提案書" in html
    assert "v-abc123" in html
    assert "本文サンプル" in html
    assert "src.pdf" in html
    # skipped ブロックが出る
    assert "要確認（データ未検出）" in html
    assert "{2}" in html


def test_build_proposal_html_escapes_html() -> None:
    co = ComposerOutput.model_construct(
        placeholders={1: "<script>alert(1)</script>"},
        citations_per_placeholder={},
        skipped_placeholders=[],
    )
    html = build_proposal_html(co, product_name="<b>x</b>", version_id="v-1")
    assert "<script>" not in html
    assert "&lt;script&gt;" in html


# -----------------------------------------------------------
# version_id は常に付く
# -----------------------------------------------------------
def test_version_id_is_emitted(tmp_path: Path) -> None:
    template = _dummy_template(tmp_path / "t.pptx")
    bedrock = MagicMock()
    bedrock.converse.return_value = _resp(_full_composer_json())
    skill = ProposalDeckSkill(bedrock=bedrock)
    out = skill.run(_input(template, tmp_path / "out"), ctx=SkillContext())
    assert out.version_id.startswith("v-")
    assert len(out.version_id) > 2
    # emit_pdf 既定 OFF → PDF なし（後方互換）
    assert out.pdf_path is None
    assert out.pdf_url is None


# -----------------------------------------------------------
# emit_pdf=True → PDF を生成（weasyprint は fake）
# -----------------------------------------------------------
def test_emit_pdf_generates_file(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    def _fake_pdf(html_str: str, out_path: Path) -> None:
        Path(out_path).write_bytes(b"%PDF-1.4 fake")

    monkeypatch.setattr("teamagent.skills.proposal_deck.pdf_export._html_to_pdf", _fake_pdf)
    monkeypatch.delenv("USE_PROPOSAL_DECK_PUBLISH", raising=False)
    template = _dummy_template(tmp_path / "t.pptx")
    bedrock = MagicMock()
    bedrock.converse.return_value = _resp(_full_composer_json())
    skill = ProposalDeckSkill(bedrock=bedrock)
    out = skill.run(_input(template, tmp_path / "out", emit_pdf=True), ctx=SkillContext())
    assert out.pdf_path is not None
    assert Path(out.pdf_path).exists()
    assert out.pdf_path.endswith(".pdf")
    assert out.pdf_url is None  # publish 未有効


def test_emit_pdf_failure_is_graceful(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    def _boom(html_str: str, out_path: Path) -> None:
        raise RuntimeError("weasyprint missing")

    monkeypatch.setattr("teamagent.skills.proposal_deck.pdf_export._html_to_pdf", _boom)
    template = _dummy_template(tmp_path / "t.pptx")
    bedrock = MagicMock()
    bedrock.converse.return_value = _resp(_full_composer_json())
    skill = ProposalDeckSkill(bedrock=bedrock)
    out = skill.run(_input(template, tmp_path / "out", emit_pdf=True), ctx=SkillContext())
    # PDF 失敗でも skill は成功（PPTX が正本）
    assert out.pdf_path is None
    assert out.coverage_ratio == 1.0
    assert Path(out.pptx_path).exists()


def test_emit_pdf_publishes_when_enabled(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    from unittest.mock import patch

    def _fake_pdf(html_str: str, out_path: Path) -> None:
        Path(out_path).write_bytes(b"%PDF-1.4 fake")

    monkeypatch.setattr("teamagent.skills.proposal_deck.pdf_export._html_to_pdf", _fake_pdf)
    monkeypatch.setenv("USE_PROPOSAL_DECK_PUBLISH", "1")
    template = _dummy_template(tmp_path / "t.pptx")
    bedrock = MagicMock()
    bedrock.converse.return_value = _resp(_full_composer_json())
    skill = ProposalDeckSkill(bedrock=bedrock)
    fake_pdf_url = "https://s3/vseo-proposals/x.pdf?sig=..."
    fake_pptx_url = "https://s3/vseo-proposals/x.pptx?sig=..."
    with (
        patch(
            "teamagent.adapters.report_publish.publish_pdf_file", return_value=fake_pdf_url
        ) as mock_pdf,
        patch("teamagent.adapters.report_publish.publish_pptx_file", return_value=fake_pptx_url),
    ):
        out = skill.run(_input(template, tmp_path / "out", emit_pdf=True), ctx=SkillContext())
    assert out.pdf_url == fake_pdf_url
    assert out.pptx_url == fake_pptx_url
    mock_pdf.assert_called_once()
