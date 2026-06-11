"""ProposalDeckSkill の単体テスト（bedrock を MagicMock、dummy template で pptx 生成）。"""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock

import pytest

from teamagent.adapters.bedrock_client import ConverseResponse, TokenUsage
from teamagent.skills.base import SkillContext
from teamagent.skills.proposal_deck.contract import LENGTH_RULES, VALID_IDS, ComposerOutput
from teamagent.skills.proposal_deck.schema import ProposalDeckInput
from teamagent.skills.proposal_deck.skill import ProposalDeckSkill


def _full_composer_json() -> str:
    """全 95 placeholder を文字数規則どおり埋めた ComposerOutput の JSON。"""
    placeholders: dict[int, str] = {}
    for pid in sorted(VALID_IDS):
        if pid in LENGTH_RULES:
            lo, hi = LENGTH_RULES[pid]
            placeholders[pid] = "サ" * ((lo + hi) // 2)
        else:
            placeholders[pid] = f"値-{pid}"
    return ComposerOutput(placeholders=placeholders).model_dump_json()


def _dummy_template(path: Path) -> Path:
    """VALID_IDS から ｛N：ラベルN｝ を敷き詰めたダミーテンプレ pptx。"""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]
    ids = sorted(VALID_IDS)
    for i in range(0, len(ids), 30):
        slide = prs.slides.add_slide(blank)
        tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5)).text_frame
        tf.word_wrap = True
        for pid in ids[i : i + 30]:
            tf.add_paragraph().text = f"｛{pid}：ラベル{pid}｝"
    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def _resp(text: str) -> ConverseResponse:
    return ConverseResponse(
        text=text,
        usage=TokenUsage(
            input_tokens=100,
            output_tokens=50,
            cache_read_input_tokens=0,
            cache_creation_input_tokens=0,
            cost_usd=0.01,
        ),
        model_id="jp.anthropic.claude-sonnet-4-6",
        latency_ms=100,
        stop_reason="end_turn",
    )


def _input(template: Path, out: Path, **kw: object) -> ProposalDeckInput:
    return ProposalDeckInput(
        product_name="ACME 青汁",
        goal="認知獲得",
        target_persona="20代女性",
        template_path=str(template),
        out_dir=str(out),
        **kw,  # type: ignore[arg-type]
    )


def test_generates_pptx_full_coverage(tmp_path: Path) -> None:
    template = _dummy_template(tmp_path / "t.pptx")
    bedrock = MagicMock()
    bedrock.converse.return_value = _resp(_full_composer_json())

    skill = ProposalDeckSkill(bedrock=bedrock)
    out = skill.run(_input(template, tmp_path / "out"), ctx=SkillContext())

    assert out.coverage_ratio == 1.0
    assert out.filled_count == 95
    assert out.skipped_count == 0
    assert Path(out.pptx_path).exists()
    assert out.total_cost_usd == pytest.approx(0.01)
    assert bedrock.converse.call_count == 1


def test_self_repair_invalid_then_valid(tmp_path: Path) -> None:
    template = _dummy_template(tmp_path / "t.pptx")
    bedrock = MagicMock()
    bedrock.converse.side_effect = [
        _resp('{"placeholders": {"1": "x"}}'),  # 網羅不足
        _resp(_full_composer_json()),
    ]
    skill = ProposalDeckSkill(bedrock=bedrock)
    out = skill.run(_input(template, tmp_path / "out", max_repair=1), ctx=SkillContext())

    assert out.coverage_ratio == 1.0
    assert bedrock.converse.call_count == 2
    # 累計コスト（2 回分）
    assert out.total_cost_usd == pytest.approx(0.02)


def test_json_code_fence_is_extracted(tmp_path: Path) -> None:
    template = _dummy_template(tmp_path / "t.pptx")
    bedrock = MagicMock()
    bedrock.converse.return_value = _resp("```json\n" + _full_composer_json() + "\n```")
    skill = ProposalDeckSkill(bedrock=bedrock)
    out = skill.run(_input(template, tmp_path / "out"), ctx=SkillContext())
    assert out.coverage_ratio == 1.0


def test_exhausted_repair_raises(tmp_path: Path) -> None:
    template = _dummy_template(tmp_path / "t.pptx")
    bedrock = MagicMock()
    bedrock.converse.return_value = _resp('{"placeholders": {"1": "x"}}')
    skill = ProposalDeckSkill(bedrock=bedrock)
    with pytest.raises(ValueError):
        skill.run(_input(template, tmp_path / "out", max_repair=1), ctx=SkillContext())
    assert bedrock.converse.call_count == 2
