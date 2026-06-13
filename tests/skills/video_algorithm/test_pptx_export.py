"""§Q-HTML→PPTX: スライドHTML→PPTX 変換（pptx_export）を chromium 無しで検証する。

shooter（playwright 要素スクショ）を注入差し替えしてフェイクPNGを流し込み、python-pptx 側の
16:9・1スライド=1画像・スライド枚数を検証する（CIに chromium/playwright browser は不要）。
"""

from __future__ import annotations

import struct
import zlib

from pptx import Presentation
from pptx.util import Inches

from teamagent.skills.video_algorithm.pptx_export import build_pptx, render_pptx
from teamagent.skills.video_algorithm.schema import VideoAlgorithmOutput


def _fake_png(w: int = 8, h: int = 8) -> bytes:
    """python-pptx が読める最小の有効PNG（単色）を生成する。"""

    def chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)  # 8bit truecolor
    raw = b"".join(b"\x00" + b"\xcc\xcc\xcc" * w for _ in range(h))  # 各行 filter0 + RGB
    idat = zlib.compress(raw)
    return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr) + chunk(b"IDAT", idat) + chunk(b"IEND", b"")


def test_build_pptx_is_16x9_one_picture_per_slide(tmp_path) -> None:
    pngs = [_fake_png(), _fake_png(), _fake_png()]
    out_path = str(tmp_path / "deck.pptx")
    build_pptx(pngs, out_path)

    prs = Presentation(out_path)
    assert prs.slide_width == Inches(13.333)  # 16:9
    assert prs.slide_height == Inches(7.5)
    assert len(prs.slides) == 3
    for slide in prs.slides:
        pics = [sh for sh in slide.shapes if sh.shape_type == 13]  # 13 = PICTURE
        assert len(pics) == 1
        assert pics[0].left == 0 and pics[0].top == 0
        assert pics[0].width == prs.slide_width  # フルブリード


def test_render_pptx_uses_injected_shooter(tmp_path) -> None:
    seen: dict[str, str] = {}

    def fake_shooter(html: str) -> list[bytes]:
        seen["html"] = html
        return [_fake_png(), _fake_png()]

    out = VideoAlgorithmOutput(query="新宿 ランチ")
    out_path = str(tmp_path / "p.pptx")
    res = render_pptx(out, out_path, shooter=fake_shooter)

    assert res == out_path
    assert "新宿 ランチ" in seen["html"]  # render_slides の HTML が shooter に渡る
    assert "data:video" not in seen["html"]  # 動画base64は載らない
    assert len(Presentation(out_path).slides) == 2


def test_render_pptx_graceful_on_no_sections(tmp_path) -> None:
    res = render_pptx(
        VideoAlgorithmOutput(query="kw"), str(tmp_path / "x.pptx"), shooter=lambda _h: []
    )
    assert res is None  # セクション0なら None（本体分析は壊さない）


def test_render_pptx_graceful_on_shooter_error(tmp_path) -> None:
    def boom(_html: str) -> list[bytes]:
        raise RuntimeError("chromium down")

    res = render_pptx(VideoAlgorithmOutput(query="kw"), str(tmp_path / "y.pptx"), shooter=boom)
    assert res is None


def test_skill_wires_proposal_urls_and_slack_links(monkeypatch, tmp_path) -> None:
    """skill が outputs 要求時に slides/pptx を生成→URL を out に載せ→Slack要約にリンクを出す。"""
    from teamagent.skills.video_algorithm.schema import VideoAlgorithmInput
    from teamagent.skills.video_algorithm.skill import VideoAlgorithmSkill

    # slides は publisher 注入で、pptx 生成（chromium 必要）は _build_pptx 差し替えで chromium 不要に。
    def fake_pub(path: str, *, request_id: str, query: str) -> str | None:
        return "https://signed.example/slides"

    skill = VideoAlgorithmSkill(publisher=fake_pub, report_dir=str(tmp_path))
    monkeypatch.setattr(skill, "_build_pptx", lambda *a, **k: "https://signed.example/pptx")

    out = VideoAlgorithmOutput(query="kw")
    inp = VideoAlgorithmInput(query="kw", outputs=["slides", "pptx"])
    skill._build_proposal_outputs(out, inp, "req1")

    assert out.slides_url == "https://signed.example/slides"
    assert out.pptx_url == "https://signed.example/pptx"
    summary = skill._slack_summary(out)
    assert "提案用パワポ" in summary and "https://signed.example/pptx" in summary
    assert "編集用スライド" in summary
