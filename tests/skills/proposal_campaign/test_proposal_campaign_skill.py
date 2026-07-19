"""ProposalCampaignSkill のオーケストレーションテスト（DI・ネット非依存）。"""

from __future__ import annotations

import json
from collections.abc import Callable
from pathlib import Path

import pytest

from teamagent.adapters.tiktok_scraper import TikTokScrapeError, TikTokVideo
from teamagent.skills.base import SkillContext
from teamagent.skills.proposal_campaign.schema import ProposalCampaignInput
from teamagent.skills.proposal_campaign.skill import ProposalCampaignSkill


def _ctx(rid: str) -> SkillContext:
    return SkillContext(request_id=rid)


@pytest.fixture(autouse=True)
def _explicit_local_media_runtime(monkeypatch: pytest.MonkeyPatch) -> None:
    """These renderer integration tests intentionally exercise local media deps."""

    monkeypatch.setenv("TEAMAGENT_LOCAL_MEDIA_RUNTIME", "true")


def test_direct_keywords(
    mock_searcher: Callable[..., list[TikTokVideo]],
    mock_fetcher: Callable[..., bytes | None],
    tmp_path: Path,
) -> None:
    skill = ProposalCampaignSkill(
        searcher=mock_searcher, fetcher=mock_fetcher, normalizer=lambda b: b
    )
    inp = ProposalCampaignInput(keywords=["集中", "作業用BGM"], image_cache_dir=str(tmp_path))
    out = skill.run(inp, _ctx("t1"))
    assert out.total_keywords == 2 and out.success_count == 2 and out.error_count == 0
    assert set(out.evidence_images) == {58, 60}
    for pid in (58, 60):
        ev = out.evidence_images[pid][0]
        assert ev.image_path is not None and Path(ev.image_path).read_bytes()


def test_from_dr_json(
    mock_searcher: Callable[..., list[TikTokVideo]],
    mock_fetcher: Callable[..., bytes | None],
    tmp_path: Path,
) -> None:
    dr = {"D_publicity": [{"trend_word": "ご褒美"}], "E_community": [{"tiktok_tags": ["#BGM"]}]}
    p = tmp_path / "dr.json"
    p.write_text(json.dumps(dr), encoding="utf-8")
    skill = ProposalCampaignSkill(
        searcher=mock_searcher, fetcher=mock_fetcher, normalizer=lambda b: b
    )
    out = skill.run(
        ProposalCampaignInput(gemini_dr_json_path=str(p), image_cache_dir=str(tmp_path)), _ctx("t2")
    )
    assert out.success_count == 2 and {kw.keyword for kw in out.results} == {"ご褒美", "BGM"}


def test_partial_failure_isolated(
    make_video: Callable[..., TikTokVideo],
    mock_fetcher: Callable[..., bytes | None],
    tmp_path: Path,
) -> None:
    def searcher(query: str, n: int, rid: str) -> list[TikTokVideo]:
        if query == "落ちる":
            raise TikTokScrapeError("captcha")
        return [make_video(cover_url=f"https://cdn.example/{query}.jpg")]

    skill = ProposalCampaignSkill(searcher=searcher, fetcher=mock_fetcher, normalizer=lambda b: b)
    out = skill.run(
        ProposalCampaignInput(keywords=["集中", "落ちる"], image_cache_dir=str(tmp_path)),
        _ctx("t3"),
    )
    # 1KW失敗でも例外を投げず他KWは成功（KW単位 isolation）
    assert out.success_count == 1 and out.error_count == 1
    assert set(out.evidence_images) == {58}


def test_all_fail_no_fallback(mock_fetcher: Callable[..., bytes | None], tmp_path: Path) -> None:
    def searcher(query: str, n: int, rid: str) -> list[TikTokVideo]:
        raise TikTokScrapeError("captcha")

    skill = ProposalCampaignSkill(
        searcher=searcher, fetcher=lambda url, rid: None, normalizer=lambda b: b
    )
    out = skill.run(
        ProposalCampaignInput(keywords=["a", "b"], image_cache_dir=str(tmp_path)), _ctx("t4")
    )
    assert out.evidence_images == {} and out.error_count == 2 and out.success_count == 0


def test_no_keywords_returns_empty(tmp_path: Path) -> None:
    skill = ProposalCampaignSkill(
        searcher=lambda q, n, r: [], fetcher=lambda u, r: None, normalizer=lambda b: b
    )
    out = skill.run(ProposalCampaignInput(keywords=[]), _ctx("t5"))
    assert out.total_keywords == 0 and out.evidence_images == {}


def test_default_request_directory_is_removed_by_cleanup(
    mock_searcher: Callable[..., list[TikTokVideo]],
    mock_fetcher: Callable[..., bytes | None],
) -> None:
    skill = ProposalCampaignSkill(
        searcher=mock_searcher, fetcher=mock_fetcher, normalizer=lambda b: b
    )

    out = skill.run(ProposalCampaignInput(keywords=["集中"]), _ctx("cleanup-request"))

    image_path = Path(out.evidence_images[58][0].image_path or "")
    request_dir = image_path.parents[1]
    assert image_path.is_file()
    assert request_dir.name.startswith("teamagent-campaign-cleanup-request-")

    skill.cleanup_output(out)

    assert not request_dir.exists()
    assert skill._temporary_output_dirs == {}


def test_request_directory_is_removed_when_render_fails(
    mock_searcher: Callable[..., list[TikTokVideo]],
    mock_fetcher: Callable[..., bytes | None],
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    request_dir = tmp_path / "request-temp"

    def fake_mkdtemp(*, prefix: str) -> str:
        assert prefix.startswith("teamagent-campaign-")
        request_dir.mkdir()
        return str(request_dir)

    monkeypatch.setattr("teamagent.skills.proposal_campaign.skill.tempfile.mkdtemp", fake_mkdtemp)
    skill = ProposalCampaignSkill(
        searcher=mock_searcher, fetcher=mock_fetcher, normalizer=lambda b: b
    )
    monkeypatch.setattr(
        skill,
        "_render",
        lambda **_kwargs: (_ for _ in ()).throw(RuntimeError("render failed")),
    )

    with pytest.raises(RuntimeError, match="render failed"):
        skill.run(
            ProposalCampaignInput(
                keywords=["集中"],
                composer_output_json_path="unused.json",
                enable_pptx_render=True,
            ),
            _ctx("failed-request"),
        )

    assert not request_dir.exists()
    assert skill._temporary_output_dirs == {}


def test_pptx_render_integration(
    mock_searcher: Callable[..., list[TikTokVideo]],
    mock_fetcher: Callable[..., bytes | None],
    tmp_path: Path,
) -> None:
    """合成テンプレ＋95枠JSON で enable_pptx_render → PICTURE が success 数だけ入る。"""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
    from pptx.util import Emu, Inches

    from teamagent.skills.proposal_deck.contract import LENGTH_RULES, VALID_IDS, ComposerOutput
    from teamagent.skills.proposal_deck.renderer import _SLOT_SIZE_EMU, _SLOT_TOP_EMU

    # 合成テンプレ（95枠テキスト + 既知スロット2個）
    tpl = tmp_path / "tpl.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    ids = sorted(VALID_IDS)
    slide = None
    for i in range(0, len(ids), 30):
        slide = prs.slides.add_slide(blank)
        tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.0)).text_frame
        for pid in ids[i : i + 30]:
            tf.add_paragraph().text = f"｛{pid}：l{pid}｝"
    assert slide is not None
    for k in range(2):
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Emu(914400 + k * (_SLOT_SIZE_EMU + 50000)),
            Emu(_SLOT_TOP_EMU),
            Emu(_SLOT_SIZE_EMU),
            Emu(_SLOT_SIZE_EMU),
        )
    prs.save(str(tpl))

    # 95枠 ComposerOutput JSON
    ph: dict[int, str] = {}
    for pid in sorted(VALID_IDS):
        if pid in LENGTH_RULES:
            lo, hi = LENGTH_RULES[pid]
            ph[pid] = "サ" * ((lo + hi) // 2)
        else:
            ph[pid] = f"v{pid}"
    cj = tmp_path / "composer.json"
    cj.write_text(ComposerOutput(placeholders=ph).model_dump_json(), encoding="utf-8")

    skill = ProposalCampaignSkill(
        searcher=mock_searcher, fetcher=mock_fetcher, normalizer=lambda b: b
    )
    out = skill.run(
        ProposalCampaignInput(
            keywords=["集中", "作業用BGM"],
            image_cache_dir=str(tmp_path),
            composer_output_json_path=str(cj),
            template_path=str(tpl),
            out_dir=str(tmp_path),
            enable_pptx_render=True,
        ),
        _ctx("t6"),
    )
    assert out.pptx_path is not None and Path(out.pptx_path).exists()
    rendered = Presentation(out.pptx_path)
    pics = sum(
        1
        for s in rendered.slides
        for sh in s.shapes
        if getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE
    )
    assert pics == out.success_count == 2
