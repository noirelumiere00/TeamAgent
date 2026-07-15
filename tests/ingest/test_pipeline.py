"""ingest/pipeline.py のテスト。

実 adapter / 実 DB は呼ばず、IngestRunner の orchestration logic を検証。
"""

from __future__ import annotations

from typing import Any
from unittest.mock import MagicMock

import pytest

from teamagent.ingest.loader import (
    GDriveFolderSpec,
    GSheetSpec,
    GSheetsTabSpec,
    IngestSources,
    SlackChannelSpec,
)
from teamagent.ingest.pipeline import IngestResult, IngestRunner, IngestStats


@pytest.fixture(autouse=True)
def _disable_freshness_monitor(monkeypatch: pytest.MonkeyPatch) -> None:
    """鮮度監視(#200)を本モジュールでは無効化する。

    _maybe_check_freshness は env ゲート無しで run 毎に走る常時ONの安全網で、
    自前の ``_ops_connection`` を1本開く。本ファイルの各テストは「対象機能（docdedup/
    boilerplate）が余分な接続を開かない」ことを ``ops_conn_calls`` で検証するため、
    無関係な freshness の +1 接続が混じると誤検知する。freshness 自体は専用の
    tests/ingest/test_freshness.py で検証済みなので、ここでは no-op 化して分離する。
    """
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.IngestRunner._maybe_check_freshness",
        lambda self, *, request_id: None,
    )


class _FakeEmbedder:
    def embed(self, text: str) -> list[float]:
        return [0.1] * 1024

    def embed_passage(self, text: str) -> list[float]:
        return self.embed(text)


class _FakeRepository:
    """実 DB なしで upsert を記録する fake。"""

    def __init__(self) -> None:
        self.upsert_calls: list[dict[str, Any]] = []

    def upsert_document_with_chunks(
        self,
        doc: Any,
        chunks: list[Any],
        request_id: str,
        *,
        replace_existing_chunks: bool = True,
    ) -> str:
        self.upsert_calls.append(
            {
                "external_id": doc.external_id,
                "source_type": doc.source_type,
                "acl_groups": list(doc.acl_groups),
                "metadata": dict(doc.metadata),
                "chunk_count": len(chunks),
                "chunks": list(chunks),
                "request_id": request_id,
            }
        )
        return "fake-doc-id"


# -----------------------------------------------------------
# IngestRunner.run() — dry-run / 集計
# -----------------------------------------------------------
def test_runner_with_no_sources_returns_empty_stats() -> None:
    """空 IngestSources を渡しても落ちず、kind ごとに stats=0 を返す。"""
    runner = IngestRunner(
        repository=_FakeRepository(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=True,
    )
    empty = IngestSources(version=1, slack_channels=(), gdrive_folders=(), gsheets=())
    result = runner.run(empty)
    assert isinstance(result, IngestResult)
    assert result.total_documents() == 0
    assert result.total_errors() == 0
    # 既定 kinds = ['slack','gdrive','gsheets','shared_drives'] すべて key として作られる
    assert set(result.by_kind.keys()) == {"slack", "gdrive", "gsheets", "shared_drives"}


def test_runner_filters_by_kinds() -> None:
    """kinds=['slack'] を渡すと gdrive/gsheets は処理されない。"""
    runner = IngestRunner(
        repository=_FakeRepository(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=True,
    )
    empty = IngestSources(version=1, slack_channels=(), gdrive_folders=(), gsheets=())
    result = runner.run(empty, kinds=["slack"])
    assert "slack" in result.by_kind
    assert "gdrive" not in result.by_kind
    assert "gsheets" not in result.by_kind


def test_runner_catches_handler_exception(monkeypatch: pytest.MonkeyPatch) -> None:
    """1 source の取り込みが例外を投げても、他 source は処理される（partial failure 許容）。"""
    repo = _FakeRepository()
    runner = IngestRunner(
        repository=repo,  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=True,
        alerter=None,  # from_env() で no-op alerter（OPS_SLACK_WEBHOOK_URL 未設定想定）
    )

    sources = IngestSources(
        version=1,
        slack_channels=(),
        gdrive_folders=(
            GDriveFolderSpec(folder_id="OK1", folder_name="ok", description=""),
            GDriveFolderSpec(folder_id="FAIL", folder_name="fail", description=""),
            GDriveFolderSpec(folder_id="OK2", folder_name="ok", description=""),
        ),
        gsheets=(),
    )

    # gdrive handler を fake で差し替え（FAIL spec のみ例外）
    def _fake_handler(spec: GDriveFolderSpec, **kwargs: Any) -> tuple[int, int]:
        if spec.folder_id == "FAIL":
            raise RuntimeError("simulated failure")
        return (2, 3)

    monkeypatch.setattr("teamagent.ingest.pipeline._ingest_gdrive_folder", _fake_handler)

    result = runner.run(sources, kinds=["gdrive"])
    stats = result.by_kind["gdrive"]
    assert stats.sources_processed == 2
    assert stats.sources_skipped == 1
    assert len(stats.errors) == 1
    assert stats.documents_upserted == 4  # 2 OK source × 2 docs
    assert stats.chunks_inserted == 6


def test_runner_calls_alerter_on_handler_exception(monkeypatch: pytest.MonkeyPatch) -> None:
    """Wave1-③: handler 例外時に alerter.send_ingest_failure が（dry_run=False で）呼ばれる。"""
    from teamagent.ingest.ops_alert import IngestOpsAlerter

    fake_alerter = MagicMock(spec=IngestOpsAlerter)
    runner = IngestRunner(
        repository=_FakeRepository(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
        alerter=fake_alerter,
    )

    sources = IngestSources(
        version=1,
        slack_channels=(),
        gdrive_folders=(GDriveFolderSpec(folder_id="FAIL", folder_name="fail", description=""),),
        gsheets=(),
    )

    def _fake_handler(spec: GDriveFolderSpec, **kwargs: Any) -> tuple[int, int]:
        raise RuntimeError("boom")

    monkeypatch.setattr("teamagent.ingest.pipeline._ingest_gdrive_folder", _fake_handler)

    runner.run(sources, kinds=["gdrive"])
    # 1 source 1 failure ⇒ 1 alert
    assert fake_alerter.send_ingest_failure.call_count == 1
    kwargs = fake_alerter.send_ingest_failure.call_args.kwargs
    assert kwargs["kind"] == "gdrive"
    assert kwargs["dry_run"] is False
    assert isinstance(kwargs["exc"], RuntimeError)


def test_runner_passes_dry_run_to_alerter(monkeypatch: pytest.MonkeyPatch) -> None:
    """dry_run=True なら alerter は呼ばれるが dry_run=True が伝達され、内部で no-op になる。"""
    from teamagent.ingest.ops_alert import IngestOpsAlerter

    fake_alerter = MagicMock(spec=IngestOpsAlerter)
    runner = IngestRunner(
        repository=_FakeRepository(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=True,
        alerter=fake_alerter,
    )
    sources = IngestSources(
        version=1,
        slack_channels=(),
        gdrive_folders=(GDriveFolderSpec(folder_id="FAIL", folder_name="fail", description=""),),
        gsheets=(),
    )

    def _fake_handler(spec: GDriveFolderSpec, **kwargs: Any) -> tuple[int, int]:
        raise RuntimeError("boom")

    monkeypatch.setattr("teamagent.ingest.pipeline._ingest_gdrive_folder", _fake_handler)

    runner.run(sources, kinds=["gdrive"])
    fake_alerter.send_ingest_failure.assert_called_once()
    assert fake_alerter.send_ingest_failure.call_args.kwargs["dry_run"] is True


# -----------------------------------------------------------
# IngestStats / IngestResult helpers
# -----------------------------------------------------------
def test_ingest_result_aggregates_totals() -> None:
    r = IngestResult()
    r.by_kind["slack"] = IngestStats(source_kind="slack", documents_upserted=10, chunks_inserted=10)
    r.by_kind["gdrive"] = IngestStats(
        source_kind="gdrive", documents_upserted=5, chunks_inserted=5, errors=["e1", "e2"]
    )
    assert r.total_documents() == 15
    assert r.total_errors() == 2


# -----------------------------------------------------------
# Slack ingest handler を単独で確認（adapter を MagicMock 差し替え）
# -----------------------------------------------------------
def test_ingest_slack_channel_handler_calls_repository(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """_ingest_slack_channel が SlackChannelIngestClient → repository.upsert を呼ぶ。"""
    from teamagent.adapters.slack_channel_ingest_client import HistoryBatch, SlackMessage

    parent = SlackMessage(
        ts="1700000001.000001",
        user="U001",
        text="親メッセージ",
        thread_ts="1700000001.000001",
        reply_count=2,
    )
    fake_history = HistoryBatch(messages=(parent,), next_cursor=None, has_more=False)
    fake_replies = HistoryBatch(messages=(parent,), next_cursor=None, has_more=False)

    from teamagent.adapters.slack_channel_ingest_client import SlackChannelMember

    fake_client = MagicMock()
    fake_client.list_channel_history.return_value = fake_history
    fake_client.list_thread_replies.return_value = fake_replies
    # 新規 ACL 解決経路: list_channel_members + get_user_emails
    fake_client.list_channel_members.return_value = (["U001", "U002"], None)
    fake_client.get_user_emails.return_value = [
        SlackChannelMember(user_id="U001", email="taro@x.jp", display_name="Taro"),
        SlackChannelMember(user_id="U002", email="jiro@x.jp", display_name="Jiro"),
    ]

    monkeypatch.setattr(
        "teamagent.adapters.slack_channel_ingest_client.SlackChannelIngestClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    # 共有 cache をテスト間で汚さないようにクリア
    from teamagent.ingest import pipeline as pipeline_mod

    pipeline_mod._USER_EMAIL_CACHE.clear()

    from teamagent.ingest.pipeline import _ingest_slack_channel

    spec = SlackChannelSpec(
        channel_id="C0XYZ",
        channel_name="#test",
        description="test",
        extra_acl_emails=("alice@x.jp",),
    )
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_slack_channel(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bob@x.jp",
        dry_run=False,
        request_id="req-1",
    )
    assert docs_n == 1
    assert chunks_n == 1
    assert len(repo.upsert_calls) == 1
    call = repo.upsert_calls[0]
    assert call["external_id"] == "C0XYZ:1700000001.000001"
    assert call["source_type"] == "slack"
    assert call["acl_groups"] == []  # §G env 未設定なら従来どおり空（後方互換）


def test_company_acl_groups_env(monkeypatch: pytest.MonkeyPatch) -> None:
    """§G: TEAMAGENT_SHARED_COMPANY_DOMAINS 設定時のみ会社ドメインを acl_groups に付与。"""
    from teamagent.ingest.pipeline import _company_acl_groups

    monkeypatch.delenv("TEAMAGENT_SHARED_COMPANY_DOMAINS", raising=False)
    assert _company_acl_groups() == []
    monkeypatch.setenv("TEAMAGENT_SHARED_COMPANY_DOMAINS", "VectorInc.co.jp")
    assert _company_acl_groups() == ["vectorinc.co.jp"]


def test_ingest_slack_channel_dry_run_skips_repository(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """dry_run=True なら repository.upsert は呼ばれない。"""
    from teamagent.adapters.slack_channel_ingest_client import HistoryBatch, SlackMessage

    parent = SlackMessage(ts="1700.000001", user="U1", text="x")
    fake_client = MagicMock()
    fake_client.list_channel_history.return_value = HistoryBatch(
        messages=(parent,), next_cursor=None, has_more=False
    )
    fake_client.list_channel_members.return_value = ([], None)
    fake_client.get_user_emails.return_value = []
    monkeypatch.setattr(
        "teamagent.adapters.slack_channel_ingest_client.SlackChannelIngestClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest import pipeline as pipeline_mod

    pipeline_mod._USER_EMAIL_CACHE.clear()

    from teamagent.ingest.pipeline import _ingest_slack_channel

    spec = SlackChannelSpec(channel_id="C0", channel_name="#t", description="")
    repo = _FakeRepository()
    _ingest_slack_channel(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=True,
        request_id="r",
    )
    assert len(repo.upsert_calls) == 0


# -----------------------------------------------------------
# ACL 解決ヘルパー（_collect_all_member_ids / _resolve_member_emails）
# -----------------------------------------------------------
def test_collect_all_member_ids_paginates() -> None:
    """list_channel_members を cursor 続く限り呼んで全 ID を集める。"""
    from teamagent.ingest.pipeline import _collect_all_member_ids

    fake = MagicMock()
    fake.list_channel_members.side_effect = [
        (["U001", "U002"], "PAGE2"),
        (["U003"], None),
    ]
    ids = _collect_all_member_ids(fake, "C0", "r")
    assert ids == ["U001", "U002", "U003"]
    assert fake.list_channel_members.call_count == 2


def test_ingest_slack_channel_skips_when_no_acl_resolved(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """email 解決 0 件 + extra_acl_emails 空 → channel ごと skip (fail-safe)。"""
    from teamagent.adapters.slack_channel_ingest_client import HistoryBatch, SlackMessage

    parent = SlackMessage(ts="1700.000001", user="U1", text="x")
    fake_client = MagicMock()
    fake_client.list_channel_history.return_value = HistoryBatch(
        messages=(parent,), next_cursor=None, has_more=False
    )
    fake_client.list_channel_members.return_value = (["U001", "U002"], None)
    # email 全部 None → 解決ゼロ
    from teamagent.adapters.slack_channel_ingest_client import SlackChannelMember

    fake_client.get_user_emails.return_value = [
        SlackChannelMember(user_id="U001", email=None, display_name="A"),
        SlackChannelMember(user_id="U002", email=None, display_name="B"),
    ]
    monkeypatch.setattr(
        "teamagent.adapters.slack_channel_ingest_client.SlackChannelIngestClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    from teamagent.ingest import pipeline as pipeline_mod

    pipeline_mod._USER_EMAIL_CACHE.clear()

    from teamagent.ingest.pipeline import _ingest_slack_channel

    spec = SlackChannelSpec(
        channel_id="C0NOACL",
        channel_name="#noacl",
        description="",
        extra_acl_emails=(),  # 空
    )
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_slack_channel(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r",
    )
    # skip された → 0 件、repository も呼ばれない
    assert docs_n == 0
    assert chunks_n == 0
    assert len(repo.upsert_calls) == 0
    # list_channel_history が呼ばれる前に skip するので、呼ばれていてもいいし
    # 呼ばれていなくてもいい（実装は呼ばれない想定だが厳密にしすぎないこと）


def test_ingest_slack_channel_uses_extra_acl_when_email_unresolved(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """extra_acl_emails があれば email 解決ゼロでも取り込み続行。"""
    from teamagent.adapters.slack_channel_ingest_client import HistoryBatch, SlackMessage

    parent = SlackMessage(ts="1700.000001", user="U1", text="hello world")
    fake_client = MagicMock()
    fake_client.list_channel_history.return_value = HistoryBatch(
        messages=(parent,), next_cursor=None, has_more=False
    )
    fake_client.list_channel_members.return_value = (["U001"], None)
    from teamagent.adapters.slack_channel_ingest_client import SlackChannelMember

    fake_client.get_user_emails.return_value = [
        SlackChannelMember(user_id="U001", email=None, display_name="A"),
    ]
    monkeypatch.setattr(
        "teamagent.adapters.slack_channel_ingest_client.SlackChannelIngestClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    from teamagent.ingest import pipeline as pipeline_mod

    pipeline_mod._USER_EMAIL_CACHE.clear()

    from teamagent.ingest.pipeline import _ingest_slack_channel

    spec = SlackChannelSpec(
        channel_id="C0EXTRA",
        channel_name="#extra",
        description="",
        extra_acl_emails=("alice@x.jp",),  # 解決失敗してもこれが ACL
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_slack_channel(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r",
    )
    assert docs_n == 1
    assert len(repo.upsert_calls) == 1


def test_resolve_member_emails_caches_and_excludes_bots() -> None:
    """get_user_emails を呼んで cache、Bot/deleted は除外。"""
    from teamagent.adapters.slack_channel_ingest_client import SlackChannelMember
    from teamagent.ingest import pipeline as pipeline_mod
    from teamagent.ingest.pipeline import _resolve_member_emails

    pipeline_mod._USER_EMAIL_CACHE.clear()

    fake = MagicMock()
    fake.get_user_emails.return_value = [
        SlackChannelMember(user_id="U001", email="taro@x.jp", display_name="Taro"),
        SlackChannelMember(user_id="U002", email=None, display_name="NoMail"),
        SlackChannelMember(user_id="U003", email="bot@x.jp", display_name="Bot", is_bot=True),
        SlackChannelMember(user_id="U004", email="old@x.jp", display_name="Old", deleted=True),
    ]
    emails = _resolve_member_emails(fake, ["U001", "U002", "U003", "U004"], "r")
    # bot / deleted / email無し はすべて除外
    assert emails == ["taro@x.jp"]
    # 2 回目は cache hit で API 呼ばない
    fake.get_user_emails.reset_mock()
    emails2 = _resolve_member_emails(fake, ["U001"], "r")
    assert emails2 == ["taro@x.jp"]
    fake.get_user_emails.assert_not_called()


# -----------------------------------------------------------
# Sheet handler の最小確認（row → document 化）
# -----------------------------------------------------------
def test_ingest_gsheet_handler_row_per_document(monkeypatch: pytest.MonkeyPatch) -> None:
    from teamagent.adapters.gsheets_client import TabRows

    fake_client = MagicMock()
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1V",
        tab_name="フォーム回答 1",
        headers=("業界", "温度感"),
        rows=(("飲食", "高"), ("コスメ", "中")),
        row_count=2,
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1V",
        sheet_name="FB",
        description="",
        tabs=(GSheetsTabSpec(gid=537831563, tab_name="フォーム回答 1"),),
    )
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r",
    )
    assert docs_n == 2  # 2 行 = 2 document
    assert chunks_n == 2
    # external_id は <sheet_id>:<gid>:<row_idx> 形式（row_idx=2 から）
    assert repo.upsert_calls[0]["external_id"] == "1V:537831563:2"
    assert repo.upsert_calls[1]["external_id"] == "1V:537831563:3"


# -----------------------------------------------------------
# Sheet handler — 営業 FB フォーム回答シートの構造化メタ（2026-07-03）
# -----------------------------------------------------------
# ショート動画営業 FB のフォーム回答シートの実ヘッダ（ユーザー実データ確認済・仮名化不要）。
# 「顧客名・案件名」（'・' 区切り）と「顧客反応(ポジ・ネガ)」（半角括弧・統合列）が
# Slack フォームのラベルとの表記ゆれポイント。
_FB_SHEET_HEADERS = (
    "タイムスタンプ",
    "商流",
    "顧客名",
    "顧客名・案件名",
    "商談フェーズ",
    "商談感触（BANT）",
    "顧客反応(ポジ・ネガ)",
    "提案メニュー",
)


def test_ingest_gsheet_fb_sheet_rows_get_structured_metadata(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """FB シート行はヘッダ → 値が Slack FB 経路と同じキーでメタ化される。

    is_sales_fb / channel_type / agency_name / client_case / deal_phase /
    bant_score / client_reaction / proposed_menu / client_name（導出）。
    """
    from teamagent.adapters.gsheets_client import TabRows

    fake_client = MagicMock()
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1V",
        tab_name="フォーム回答 1",
        headers=_FB_SHEET_HEADERS,
        rows=(
            (
                "2026/06/30 10:00:00",
                "代理店",
                "アルファ広告社",
                "ベータ商事 / ガンマ製品",
                "ヒアリング",
                "B（前向き）",
                "ポジ: 適合性を評価 / ネガ: 予算感に懸念",
                "UGC（TTO、切り抜きなど）",
            ),
            # コア列が全部空（タイムスタンプのみ）の行 → FB メタは付かない
            ("2026/07/01 09:00:00", "", "", "", "", "", "", ""),
        ),
        row_count=2,
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1V",
        sheet_name="ショート動画営業 FB - フォーム回答",
        description="",
        tabs=(GSheetsTabSpec(gid=537831563, tab_name="フォーム回答 1"),),
        extra_metadata={"topic": "営業 FB", "priority": "high"},
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r-gsheet-fb",
    )
    assert docs_n == 2

    md = repo.upsert_calls[0]["metadata"]
    # Slack FB 経路と同じ first-class メタ
    assert md["is_sales_fb"] is True
    assert md["channel_type"] == "代理店"
    assert md["agency_name"] == "アルファ広告社"
    assert md["client_case"] == "ベータ商事 / ガンマ製品"
    assert md["deal_phase"] == "ヒアリング"
    assert md["bant_score"] == "B（前向き）"
    assert md["client_reaction"] == "ポジ: 適合性を評価 / ネガ: 予算感に懸念"
    assert md["proposed_menu"] == "UGC（TTO、切り抜きなど）"
    # client_name は client_case の '/' 左側から導出
    assert md["client_name"] == "ベータ商事"
    # 既存の固定キー / extra_metadata は破壊されない
    assert md["tab_name"] == "フォーム回答 1"
    assert md["row_idx"] == 2
    assert md["topic"] == "営業 FB"
    assert md["priority"] == "high"

    # コア列が全部空の行には FB メタを付けない
    md_empty = repo.upsert_calls[1]["metadata"]
    assert "is_sales_fb" not in md_empty
    assert "client_name" not in md_empty
    assert "deal_phase" not in md_empty

    # FB シートにはナレッジ共有メタは付かない（コアヘッダが交差しない・相互排他）
    assert "is_knowledge_share" not in md
    assert "client_company" not in md
    assert "knowledge_kind" not in md


def test_ingest_gsheet_non_fb_sheet_metadata_unchanged(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """非 FB シート（ナレッジ共有フォーム等・コア列 < 3）は従来どおりメタ追加なし。"""
    from teamagent.adapters.gsheets_client import TabRows

    fake_client = MagicMock()
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1K",
        tab_name="フォーム回答 1",
        headers=("タイムスタンプ", "共有者", "ナレッジ内容", "カテゴリ"),
        rows=(("2026/06/30 10:00:00", "山田太郎", "TikTok 新機能について", "メディア動向"),),
        row_count=1,
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1K",
        sheet_name="ナレッジ共有",
        description="",
        tabs=(GSheetsTabSpec(gid=1, tab_name="フォーム回答 1"),),
        extra_metadata={"topic": "提案ナレッジ"},
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r-gsheet-nonfb",
    )
    assert docs_n == 1
    md = repo.upsert_calls[0]["metadata"]
    # 従来キーのみ（FB メタもナレッジ共有メタも一切付かない）
    assert set(md) == {"topic", "tab_name", "row_idx"}
    assert "is_sales_fb" not in md
    assert "is_knowledge_share" not in md


# -----------------------------------------------------------
# Sheet handler — ナレッジ共有フォーム回答シートの構造化メタ（2026-07-06）
# -----------------------------------------------------------
# ナレッジ共有フォーム回答シートの実ヘッダ 13 列（Drive API 実データ 190 行で確認済・
# 仮名化不要）。値は実データの形（法人格/敬称/カンマ多選択/GAS 運用列）を模した仮名。
_KNOWLEDGE_SHEET_HEADERS = (
    "ファイルをアップ",
    "正式社名",
    "案件名",
    "クライアント種別",
    "提案プロダクト",
    "資料の概要",
    "このナレッジのポイントはここ！",
    "なぜそのナレッジ（資料）を共有したのか？",
    "フリーコメント",
    "送信者",
    "タイムスタンプ",
    "ドライブ格納",
    "保管先フォルダID記録（GAS処理)",
)


def test_ingest_gsheet_knowledge_sheet_rows_get_structured_metadata(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """ナレッジ共有シート行は form_mappings の写像で first-class メタ化される。

    is_knowledge_share / client_company / client_case / client_type / proposed_menu /
    knowledge_kind / submitter ＋ client_name（正式社名から FB と同品質で導出）。
    is_sales_fb は立てない（これは FB ではなくナレッジ共有）。
    """
    from teamagent.adapters.gsheets_client import TabRows

    fake_client = MagicMock()
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1J",
        tab_name="フォーム回答 1",
        headers=_KNOWLEDGE_SHEET_HEADERS,
        rows=(
            (
                "https://slack.example/files/F1/deck.pdf",
                "株式会社デルタ製薬様",  # 法人格 prefix ＋ 敬称 suffix（実データの揺れ形）
                "新製品プロモーション",
                "TOP500 or ベス10,メーカー",  # カンマ多選択（実値域は企業属性・業種ではない）
                "ビデオリリース,タテガタ",  # カンマ多選択（自社プロダクト名）
                "提案",
                "",  # このナレッジのポイントはここ！（実データでは全行空）
                "",  # なぜ共有したのか（同上）
                "参考までに共有します（長文自由記述）",
                "@山田太郎",
                "2025/06/17 13:14:45",
                "20250617_デルタ製薬",
                "-",
            ),
            # 正式社名がプレースホルダの行 → client_name は導出されない
            (
                "",
                "なし",
                "",
                "その他",
                "その他",
                "その他ナレッジ",
                "",
                "",
                "社内 Tips です",
                "@佐藤花子",
                "2025/07/01 09:00:00",
                "",
                "-",
            ),
        ),
        row_count=2,
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1J",
        sheet_name="ナレッジ共有 - フォーム回答",
        description="",
        tabs=(GSheetsTabSpec(gid=278789217, tab_name="フォーム回答 1"),),
        extra_metadata={"topic": "提案ナレッジ"},
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r-gsheet-knowledge",
    )
    assert docs_n == 2

    md = repo.upsert_calls[0]["metadata"]
    assert md["is_knowledge_share"] is True
    assert md["client_company"] == "株式会社デルタ製薬様"  # 生値（監査用）
    assert md["client_name"] == "デルタ製薬"  # 法人格・敬称を落とした first-class 名
    assert md["client_case"] == "新製品プロモーション"
    assert md["client_type"] == "TOP500 or ベス10,メーカー"
    assert md["proposed_menu"] == "ビデオリリース,タテガタ"
    assert md["knowledge_kind"] == "提案"
    assert md["submitter"] == "@山田太郎"
    # FB フラグは立てない・運用列（URL/フリーコメント/タイムスタンプ/GAS 列）は写像しない
    assert "is_sales_fb" not in md
    assert set(md) == {
        "topic",
        "tab_name",
        "row_idx",
        "is_knowledge_share",
        "client_company",
        "client_name",
        "client_case",
        "client_type",
        "proposed_menu",
        "knowledge_kind",
        "submitter",
    }
    # 既存の固定キー / extra_metadata は破壊されない
    assert md["tab_name"] == "フォーム回答 1"
    assert md["row_idx"] == 2
    assert md["topic"] == "提案ナレッジ"

    # プレースホルダ社名（なし）の行: 写像は付くが client_name は導出されない
    md2 = repo.upsert_calls[1]["metadata"]
    assert md2["is_knowledge_share"] is True
    assert md2["client_company"] == "なし"
    assert "client_name" not in md2
    assert "client_case" not in md2  # 空値は drop


def test_ingest_gsheet_knowledge_metadata_coexists_with_classification(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """ナレッジ共有メタと cls_* が併存する（合成順 knowledge → cls・キーは交差しない）。

    人間入力（client_type/proposed_menu/knowledge_kind）と Haiku
    （cls_industry/cls_solution/cls_doc_type）は別キー設計なので、cls を後置しても
    人間入力が上書きされることはない（根拠は form_mappings の module docstring）。
    """
    from teamagent.adapters.gsheets_client import TabRows

    fake_client = MagicMock()
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1J",
        tab_name="フォーム回答 1",
        headers=_KNOWLEDGE_SHEET_HEADERS,
        rows=(
            (
                "",
                "イプシロン食品株式会社",
                "レトルト新商品ローンチ",
                "上場企業",
                "ソリューションプラン",
                "提案,クロージング",
                "",
                "",
                "決め手になった構成です",
                "@田中一郎",
                "2025/07/02 10:00:00",
                "",
                "-",
            ),
        ),
        row_count=1,
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    stub = _StubClassifier()
    monkeypatch.setattr("teamagent.ingest.classify.build_classifier_from_env", lambda: stub)

    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1J",
        sheet_name="ナレッジ共有 - フォーム回答",
        description="",
        tabs=(GSheetsTabSpec(gid=278789217, tab_name="フォーム回答 1"),),
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r-gsheet-knowledge-cls",
    )
    assert docs_n == 1
    md = repo.upsert_calls[0]["metadata"]
    # cls_* はナレッジメタに上書きされない（stub の値のまま）
    _assert_classified(md)
    # 人間入力メタも同時に付く（proposed_menu=人間入力、cls_solution 系は Haiku で併存）
    assert md["is_knowledge_share"] is True
    assert md["client_name"] == "イプシロン食品"  # 法人格 suffix を除去
    assert md["client_type"] == "上場企業"
    assert md["proposed_menu"] == "ソリューションプラン"
    assert md["knowledge_kind"] == "提案,クロージング"
    assert md["cls_doc_type"] == "提案書"  # Haiku の doc_type はそのまま


# -----------------------------------------------------------
# Drive folder ingest — PDF 本文抽出 + ACL 解決
# -----------------------------------------------------------
def _make_drive_file(
    id: str,
    name: str = "test.pdf",
    mime: str = "application/pdf",
    owners: tuple[str, ...] = (),
) -> Any:
    from teamagent.adapters.gdrive_client import DriveFile

    return DriveFile(
        id=id,
        name=name,
        mime_type=mime,
        modified_time="2026-05-26T16:00:00Z",
        size=1234,
        parents=(),
        web_view_link=f"https://drive.google.com/file/d/{id}/view",
        owners_email=owners,
    )


def _make_drive_perm(type: str, role: str, email: str | None, deleted: bool = False) -> Any:
    from teamagent.adapters.gdrive_client import DrivePermission

    return DrivePermission(
        id="p1", type=type, role=role, email_address=email, domain=None, deleted=deleted
    )


def test_ingest_gdrive_folder_pdf_extracts_chunks_and_resolves_acl(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """PDF を download → 抽出 → 複数 chunk + ACL を解決して upsert する。"""
    fake_client = MagicMock()
    fake_client.list_files.return_value = (
        [_make_drive_file(id="F1", name="提案書.pdf", owners=("alice@x.jp",))],
        None,
    )
    fake_client.list_permissions.return_value = [
        _make_drive_perm("user", "owner", "alice@x.jp"),
        _make_drive_perm("user", "writer", "bob@x.jp"),
        _make_drive_perm("group", "reader", "sales@x.jp"),
    ]
    fake_client.download_file_bytes.return_value = b"<fake-pdf-bytes>"

    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    # pypdf を monkeypatch して 3 ページの PDF を擬似する
    fake_reader = MagicMock()
    p1 = MagicMock()
    p1.extract_text.return_value = "ページ 1 提案書本文" * 50  # 長文 → 複数 chunk
    p2 = MagicMock()
    p2.extract_text.return_value = "ページ 2"
    fake_reader.pages = [p1, p2]
    monkeypatch.setattr("pypdf.PdfReader", lambda _s: fake_reader)

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(
        folder_id="FOLDER1",
        folder_name="営業提案書",
        description="",
        mime_type_filter="application/pdf",
    )
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-drive",
    )
    assert docs_n == 1
    assert chunks_n >= 2  # 複数 chunks（page 1 が長文なので分割される）
    call = repo.upsert_calls[0]
    assert call["external_id"] == "F1"
    assert call["source_type"] == "gdrive"


def test_ingest_gdrive_folder_non_pdf_uses_title_only(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """未対応 mime（image/zip 等）は title だけで 1 chunk 作る（雛形動作）.

    Wave2-④ 以降、gdoc/docx/pptx/xlsx は実本文抽出に変わったため、
    本テストでは「未対応 mime」を使って title-only fallback を確認する。
    """
    fake_client = MagicMock()
    fake_client.list_files.return_value = (
        [
            _make_drive_file(
                id="DOC1",
                name="image.png",
                mime="image/png",
                owners=("alice@x.jp",),
            )
        ],
        None,
    )
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "alice@x.jp")]
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(
        folder_id="FOLDER2",
        folder_name="議事録",
        description="",
        mime_type_filter=None,
    )
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r",
    )
    assert docs_n == 1
    assert chunks_n == 1
    # download_file_bytes は呼ばれない（PDF 以外）
    fake_client.download_file_bytes.assert_not_called()


def _setup_fake_drive_pdf(monkeypatch: pytest.MonkeyPatch, *, name: str) -> Any:
    """分類テスト用: PDF 1 本を返す fake GDriveClient をセットアップする。"""
    fake_client = MagicMock()
    fake_client.list_files.return_value = (
        [_make_drive_file(id="F1", name=name, owners=("alice@x.jp",))],
        None,
    )
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "alice@x.jp")]
    fake_client.download_file_bytes.return_value = b"<fake-pdf>"
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    fake_reader = MagicMock()
    page = MagicMock()
    page.extract_text.return_value = "アース製薬向け SNS 提案の本文" * 10
    fake_reader.pages = [page]
    monkeypatch.setattr("pypdf.PdfReader", lambda _s: fake_reader)
    return fake_client


def test_ingest_gdrive_folder_applies_classification(monkeypatch: pytest.MonkeyPatch) -> None:
    """USE_DOC_CLASSIFY=1 のとき、本文を分類して cls_* を documents.metadata に付与する。"""
    monkeypatch.setenv("USE_DOC_CLASSIFY", "1")
    _setup_fake_drive_pdf(monkeypatch, name="アース製薬_提案.pdf")

    fake_bedrock = MagicMock()
    fake_bedrock.converse.return_value = MagicMock(
        text='{"project": "アース製薬", "industry": "日用品", "doc_type": "提案書", "phase": "提案"}'
    )
    monkeypatch.setattr(
        "teamagent.adapters.bedrock_client.BedrockClient.from_env",
        classmethod(lambda cls: fake_bedrock),
    )

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(
        folder_id="FOLDER1",
        folder_name="ナレッジ",
        description="",
        mime_type_filter="application/pdf",
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-cls",
    )
    assert docs_n == 1
    md = repo.upsert_calls[0]["metadata"]
    assert md["cls_project"] == "アース製薬"
    assert md["cls_industry"] == "日用品"
    assert md["cls_doc_type"] == "提案書"
    assert md["cls_phase"] == "提案"
    assert md["industry"] == "日用品"  # 既存の業界フィルタと整合


def test_ingest_gdrive_folder_no_classification_when_disabled(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """USE_DOC_CLASSIFY 未設定（既定）なら分類せず、Bedrock も呼ばない（完全後方互換）。"""
    monkeypatch.delenv("USE_DOC_CLASSIFY", raising=False)
    _setup_fake_drive_pdf(monkeypatch, name="提案.pdf")

    def _boom(cls: type) -> None:
        raise AssertionError("分類無効時に Bedrock を呼んではいけない")

    monkeypatch.setattr(
        "teamagent.adapters.bedrock_client.BedrockClient.from_env", classmethod(_boom)
    )

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(
        folder_id="FOLDER1",
        folder_name="ナレッジ",
        description="",
        mime_type_filter="application/pdf",
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-nocls",
    )
    assert docs_n == 1
    md = repo.upsert_calls[0]["metadata"]
    assert "cls_project" not in md
    assert "cls_doc_type" not in md


def test_ingest_gdrive_folder_docx_extracts_chunks(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Wave2-④: docx を download → office_extract → chunk 化して upsert."""
    from io import BytesIO

    from docx import Document

    from teamagent.ingest.office_extract import DOCX_MIME

    # 実 docx バイナリを作る
    doc = Document()
    doc.add_paragraph("提案書の本文サンプル。" * 60)
    buf = BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    fake_client = MagicMock()
    fake_client.list_files.return_value = (
        [_make_drive_file(id="DOCX1", name="提案書.docx", mime=DOCX_MIME, owners=("a@x.jp",))],
        None,
    )
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "a@x.jp")]
    fake_client.download_file_bytes.return_value = docx_bytes
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(folder_id="F", folder_name="x", description="", mime_type_filter=None)
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r",
    )
    assert docs_n == 1
    assert chunks_n >= 1
    assert repo.upsert_calls[0]["external_id"] == "DOCX1"
    fake_client.download_file_bytes.assert_called_once()


def test_ingest_gdrive_folder_pptx_extracts_per_slide(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Wave2-④: pptx は slide 単位で page_num を持つ chunk になる."""
    from io import BytesIO

    from pptx import Presentation

    from teamagent.ingest.office_extract import PPTX_MIME

    prs = Presentation()
    layout = prs.slide_layouts[5]
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "スライド1"
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "スライド2"
    buf = BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    fake_client = MagicMock()
    fake_client.list_files.return_value = (
        [_make_drive_file(id="P1", name="deck.pptx", mime=PPTX_MIME, owners=("a@x.jp",))],
        None,
    )
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "a@x.jp")]
    fake_client.download_file_bytes.return_value = pptx_bytes
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(folder_id="F", folder_name="x", description="", mime_type_filter=None)
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r",
    )
    assert docs_n == 1
    assert chunks_n == 2  # 2 slides → 2 chunks（小さいので分割なし）


def test_ingest_gdrive_folder_xlsx_extracts_per_sheet(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Wave2-④: xlsx は sheet 単位で page_num を持つ chunk になる."""
    from io import BytesIO

    from openpyxl import Workbook

    from teamagent.ingest.office_extract import XLSX_MIME

    wb = Workbook()
    wb.active.title = "Sheet1"
    wb.active["A1"] = "顧客"
    wb.active["B1"] = 100
    wb.create_sheet("Sheet2")["A1"] = "備考"
    buf = BytesIO()
    wb.save(buf)
    xlsx_bytes = buf.getvalue()

    fake_client = MagicMock()
    fake_client.list_files.return_value = (
        [_make_drive_file(id="X1", name="売上.xlsx", mime=XLSX_MIME, owners=("a@x.jp",))],
        None,
    )
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "a@x.jp")]
    fake_client.download_file_bytes.return_value = xlsx_bytes
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(folder_id="F", folder_name="x", description="", mime_type_filter=None)
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r",
    )
    assert docs_n == 1
    assert chunks_n == 2  # 2 sheets → 2 chunks


def test_ingest_gdrive_folder_gdoc_uses_gdocs_client(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Wave2-④: Google native gdoc は GDocsClient.get_document_text で本文取得 → chunk 化."""
    from teamagent.adapters.gdocs_client import DocContent

    fake_client = MagicMock()
    fake_client.list_files.return_value = (
        [
            _make_drive_file(
                id="GDOC1",
                name="議事録.gdoc",
                mime="application/vnd.google-apps.document",
                owners=("a@x.jp",),
            )
        ],
        None,
    )
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "a@x.jp")]
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    # GDocsClient.from_env を fake で差し替え（実 Google API は呼ばない）
    fake_gdocs = MagicMock()
    fake_gdocs.get_document_text.return_value = DocContent(
        document_id="GDOC1", title="議事録", text="議事録の本文サンプル。" * 30
    )
    monkeypatch.setattr(
        "teamagent.adapters.gdocs_client.GDocsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_gdocs),
    )

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(folder_id="F", folder_name="x", description="", mime_type_filter=None)
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r",
    )
    assert docs_n == 1
    assert chunks_n >= 1
    fake_gdocs.get_document_text.assert_called_once()
    # download_file_bytes は呼ばれない（gdoc は Docs API 経由）
    fake_client.download_file_bytes.assert_not_called()


def test_ingest_gdrive_folder_skips_pdf_when_download_fails(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """download_file_bytes が例外を投げると、その 1 件は skip され他はそのまま処理される。"""
    fake_client = MagicMock()
    fake_client.list_files.return_value = (
        [
            _make_drive_file(id="OK", name="ok.pdf", owners=("alice@x.jp",)),
            _make_drive_file(id="FAIL", name="fail.pdf", owners=("alice@x.jp",)),
        ],
        None,
    )
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "alice@x.jp")]

    def _dl(file_id: str, request_id: str) -> bytes:
        if file_id == "FAIL":
            raise RuntimeError("simulated 403")
        return b"<fake>"

    fake_client.download_file_bytes.side_effect = _dl
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    fake_reader = MagicMock()
    p = MagicMock()
    p.extract_text.return_value = "ok content"
    fake_reader.pages = [p]
    monkeypatch.setattr("pypdf.PdfReader", lambda _s: fake_reader)

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(
        folder_id="F",
        folder_name="x",
        description="",
        mime_type_filter="application/pdf",
    )
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r",
    )
    # OK 1 件のみ upsert される
    assert docs_n == 1
    assert chunks_n >= 1
    assert len(repo.upsert_calls) == 1
    assert repo.upsert_calls[0]["external_id"] == "OK"


def test_ingest_gdrive_folder_paginates_list_files(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """list_files の next_page_token を辿って全件取得する.

    pagination 自体の検証なので、本文抽出を伴わない未対応 mime（image/png）を使う
    （gdoc は GDocsClient credentials を要するため別テストで扱う）。
    """
    fake_client = MagicMock()
    fake_client.list_files.side_effect = [
        ([_make_drive_file(id="F1", mime="image/png")], "PAGE2"),
        ([_make_drive_file(id="F2", mime="image/png")], None),
    ]
    fake_client.list_permissions.return_value = []
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(folder_id="F", folder_name="x", description="", mime_type_filter=None)
    repo = _FakeRepository()
    docs_n, _chunks_n = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r",
    )
    assert docs_n == 2
    assert fake_client.list_files.call_count == 2


def test_resolve_drive_file_acl_extracts_owner_and_acl(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """_resolve_drive_file_acl が owner / user / group を正しく分解する。"""
    fake_client = MagicMock()
    fake_client.list_permissions.return_value = [
        _make_drive_perm("user", "owner", "owner@x.jp"),
        _make_drive_perm("user", "writer", "bob@x.jp"),
        _make_drive_perm("user", "reader", "carol@x.jp", deleted=True),  # 除外
        _make_drive_perm("group", "reader", "sales@x.jp"),
    ]

    from teamagent.ingest.pipeline import _resolve_drive_file_acl

    owner, emails, groups = _resolve_drive_file_acl(
        fake_client, "F1", "r", fallback_owner_email="fallback@x.jp"
    )
    assert owner == "owner@x.jp"
    # owner は ACL に含まれる（fail-safe）+ user role!='owner' な bob のみ
    assert set(emails) == {"owner@x.jp", "bob@x.jp"}
    assert groups == ["sales@x.jp"]


def test_resolve_drive_file_acl_returns_fallback_on_failure(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """permissions.list が例外を投げたら fallback owner を返す。"""
    fake_client = MagicMock()
    fake_client.list_permissions.side_effect = RuntimeError("403")

    from teamagent.ingest.pipeline import _resolve_drive_file_acl

    owner, emails, groups = _resolve_drive_file_acl(
        fake_client, "F1", "r", fallback_owner_email="fallback@x.jp"
    )
    assert owner == "fallback@x.jp"
    assert emails == ["fallback@x.jp"]
    assert groups == []


# -----------------------------------------------------------
# Contextual Retrieval 配線（P1）
# -----------------------------------------------------------
class _FakeContextualizer:
    """各 chunk の contextualized / embedding を埋める fake（実 Bedrock/embed なし）。"""

    def __init__(self) -> None:
        self.calls: list[tuple[str, str, int]] = []

    def contextualize_chunks(
        self,
        doc_title: str,
        full_text: str,
        chunks: list[Any],
        request_id: str,
    ) -> list[Any]:
        from dataclasses import replace

        self.calls.append((doc_title, full_text, len(chunks)))
        return [
            replace(
                c,
                contextualized=f"[ctx] {c.content}",
                embedding=[0.5] * len(c.embedding),
            )
            for c in chunks
        ]


def test_ingest_gdrive_folder_contextualizes_when_injected(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """USE_CONTEXTUAL_INGEST 相当: contextualizer 注入で chunks.contextualized が埋まる。"""
    _setup_fake_drive_pdf(monkeypatch, name="提案.pdf")

    fake_ctx = _FakeContextualizer()
    monkeypatch.setattr(
        "teamagent.ingest.contextualize.build_contextualizer_from_env",
        lambda: fake_ctx,
    )

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(
        folder_id="FOLDER1",
        folder_name="ナレッジ",
        description="",
        mime_type_filter="application/pdf",
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-ctx",
    )
    assert docs_n == 1
    assert fake_ctx.calls  # contextualizer が呼ばれた
    chunks = repo.upsert_calls[0]["chunks"]
    assert chunks
    assert all(c.contextualized is not None for c in chunks)
    assert all(c.contextualized.startswith("[ctx] ") for c in chunks)


def test_ingest_gdrive_folder_no_contextualize_when_disabled(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """contextualizer None（既定）なら contextualized は None のまま（完全後方互換）。"""
    _setup_fake_drive_pdf(monkeypatch, name="提案.pdf")

    monkeypatch.setattr(
        "teamagent.ingest.contextualize.build_contextualizer_from_env",
        lambda: None,
    )

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(
        folder_id="FOLDER1",
        folder_name="ナレッジ",
        description="",
        mime_type_filter="application/pdf",
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-noctx",
    )
    assert docs_n == 1
    chunks = repo.upsert_calls[0]["chunks"]
    assert chunks
    assert all(c.contextualized is None for c in chunks)


def test_ingest_slack_channel_contextualizes_when_injected(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Slack 経路でも contextualizer 注入で chunks.contextualized が埋まる。"""
    from teamagent.adapters.slack_channel_ingest_client import (
        HistoryBatch,
        SlackChannelMember,
        SlackMessage,
    )

    parent = SlackMessage(ts="1700.000001", user="U1", text="スレッド本文サンプル")
    fake_client = MagicMock()
    fake_client.list_channel_history.return_value = HistoryBatch(
        messages=(parent,), next_cursor=None, has_more=False
    )
    fake_client.list_channel_members.return_value = (["U001"], None)
    fake_client.get_user_emails.return_value = [
        SlackChannelMember(user_id="U001", email="taro@x.jp", display_name="Taro"),
    ]
    monkeypatch.setattr(
        "teamagent.adapters.slack_channel_ingest_client.SlackChannelIngestClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    fake_ctx = _FakeContextualizer()
    monkeypatch.setattr(
        "teamagent.ingest.contextualize.build_contextualizer_from_env",
        lambda: fake_ctx,
    )

    from teamagent.ingest import pipeline as pipeline_mod

    pipeline_mod._USER_EMAIL_CACHE.clear()

    from teamagent.ingest.pipeline import _ingest_slack_channel

    spec = SlackChannelSpec(channel_id="C0CTX", channel_name="#ctx", description="")
    repo = _FakeRepository()
    docs_n, _ = _ingest_slack_channel(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r-slack-ctx",
    )
    assert docs_n == 1
    assert fake_ctx.calls
    chunks = repo.upsert_calls[0]["chunks"]
    assert chunks
    assert all(c.contextualized is not None for c in chunks)


def test_ingest_slack_channel_no_contextualize_when_disabled(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Slack 経路: contextualizer None（既定）なら contextualized は None のまま。"""
    from teamagent.adapters.slack_channel_ingest_client import (
        HistoryBatch,
        SlackChannelMember,
        SlackMessage,
    )

    parent = SlackMessage(ts="1700.000001", user="U1", text="スレッド本文")
    fake_client = MagicMock()
    fake_client.list_channel_history.return_value = HistoryBatch(
        messages=(parent,), next_cursor=None, has_more=False
    )
    fake_client.list_channel_members.return_value = (["U001"], None)
    fake_client.get_user_emails.return_value = [
        SlackChannelMember(user_id="U001", email="taro@x.jp", display_name="Taro"),
    ]
    monkeypatch.setattr(
        "teamagent.adapters.slack_channel_ingest_client.SlackChannelIngestClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    monkeypatch.setattr(
        "teamagent.ingest.contextualize.build_contextualizer_from_env",
        lambda: None,
    )

    from teamagent.ingest import pipeline as pipeline_mod

    pipeline_mod._USER_EMAIL_CACHE.clear()

    from teamagent.ingest.pipeline import _ingest_slack_channel

    spec = SlackChannelSpec(channel_id="C0NOCTX", channel_name="#noctx", description="")
    repo = _FakeRepository()
    docs_n, _ = _ingest_slack_channel(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r-slack-noctx",
    )
    assert docs_n == 1
    chunks = repo.upsert_calls[0]["chunks"]
    assert chunks
    assert all(c.contextualized is None for c in chunks)


# -----------------------------------------------------------
# 自動分類の配線: crawl / slack / gsheet 経路にも cls_* が乗ることを検証
# （USE_DOC_CLASSIFY=1 が gdrive folder だけでなく全経路で効くことの回帰）
# -----------------------------------------------------------
class _StubClassifier:
    """既知の DocClassification を返すスタブ分類器（実 Bedrock なし）。"""

    def __init__(self) -> None:
        from teamagent.ingest.classify import DocClassification

        self._result = DocClassification(
            project="アース製薬", industry="日用品", doc_type="提案書", phase="提案"
        )
        self.calls: list[tuple[str, str]] = []
        # ④ フォルダ文脈: classify に渡された folder_name を記録（未指定は ""）。
        self.folder_names: list[str] = []

    def classify(self, *, title: str, text: str, request_id: str, folder_name: str = "") -> Any:
        self.calls.append((title, text))
        self.folder_names.append(folder_name)
        return self._result


def _assert_classified(md: dict[str, Any]) -> None:
    assert md["cls_project"] == "アース製薬"
    assert md["cls_industry"] == "日用品"
    assert md["cls_doc_type"] == "提案書"
    assert md["cls_phase"] == "提案"
    assert md["industry"] == "日用品"  # 既存の業界フィルタと整合


def test_ingest_shared_drives_crawl_applies_classification(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """高ボリューム crawl 経路でも cls_* が documents.metadata に乗る。"""
    from teamagent.adapters.gdrive_client import SharedDrive

    fake_client = MagicMock()
    fake_client.list_shared_drives.return_value = [SharedDrive(id="D1", name="営業ナレッジ")]
    fake_client.walk_files_recursive.return_value = [
        _make_drive_file(id="F1", name="アース製薬_提案.pdf", owners=("alice@x.jp",))
    ]
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "alice@x.jp")]
    fake_client.download_file_bytes.return_value = b"<fake-pdf>"
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    fake_reader = MagicMock()
    page = MagicMock()
    page.extract_text.return_value = "アース製薬向け SNS 提案の本文" * 10
    fake_reader.pages = [page]
    monkeypatch.setattr("pypdf.PdfReader", lambda _s: fake_reader)

    stub = _StubClassifier()
    monkeypatch.setattr("teamagent.ingest.classify.build_classifier_from_env", lambda: stub)

    from teamagent.ingest.loader import SharedDriveCrawlSpec
    from teamagent.ingest.pipeline import _ingest_shared_drives_crawl

    spec = SharedDriveCrawlSpec(
        enabled=True,
        name_filter=("営業",),
        sales_relevance_filter=False,  # フィルタを外して fake PDF を確実に通す
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_shared_drives_crawl(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-crawl-cls",
    )
    assert docs_n == 1
    assert stub.calls  # 分類器が呼ばれた
    _assert_classified(repo.upsert_calls[0]["metadata"])
    # ④ crawl 経路はフォルダ文脈を渡さない（walk が親フォルダ名を保持しない・
    # drive.name はドライブ名でフォルダ名ではない＝正直なスコープ制限）
    assert stub.folder_names == [""]


def test_ingest_gdrive_folder_passes_folder_name_to_classifier(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """④ gdrive_folders 経路は spec.folder_name を classify のフォルダ文脈として渡す。"""
    fake_client = MagicMock()
    fake_client.list_files.return_value = (
        [_make_drive_file(id="F1", name="提案書.pdf", owners=("alice@x.jp",))],
        None,
    )
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "alice@x.jp")]
    fake_client.download_file_bytes.return_value = b"<fake-pdf>"
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    fake_reader = MagicMock()
    page = MagicMock()
    page.extract_text.return_value = "アース製薬向け SNS 提案の本文" * 10
    fake_reader.pages = [page]
    monkeypatch.setattr("pypdf.PdfReader", lambda _s: fake_reader)

    stub = _StubClassifier()
    monkeypatch.setattr("teamagent.ingest.classify.build_classifier_from_env", lambda: stub)

    from teamagent.ingest.pipeline import _ingest_gdrive_folder

    spec = GDriveFolderSpec(
        folder_id="FOLDER1",
        folder_name="ショート動画資料全般",
        description="",
        mime_type_filter="application/pdf",
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gdrive_folder(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-drive-folder-cls",
    )
    assert docs_n == 1
    assert stub.calls
    # フォルダ文脈が classify に届く（決定論ルール＋Haiku ヒントの入口）
    assert stub.folder_names == ["ショート動画資料全般"]
    _assert_classified(repo.upsert_calls[0]["metadata"])


def test_ingest_slack_channel_applies_classification(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Slack 経路でも cls_* が documents.metadata に乗る（doc_metadata に update）。"""
    from teamagent.adapters.slack_channel_ingest_client import (
        HistoryBatch,
        SlackChannelMember,
        SlackMessage,
    )

    parent = SlackMessage(ts="1700.000001", user="U1", text="スレッド本文サンプル")
    fake_client = MagicMock()
    fake_client.list_channel_history.return_value = HistoryBatch(
        messages=(parent,), next_cursor=None, has_more=False
    )
    fake_client.list_channel_members.return_value = (["U001"], None)
    fake_client.get_user_emails.return_value = [
        SlackChannelMember(user_id="U001", email="taro@x.jp", display_name="Taro"),
    ]
    monkeypatch.setattr(
        "teamagent.adapters.slack_channel_ingest_client.SlackChannelIngestClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    stub = _StubClassifier()
    monkeypatch.setattr("teamagent.ingest.classify.build_classifier_from_env", lambda: stub)

    from teamagent.ingest import pipeline as pipeline_mod

    pipeline_mod._USER_EMAIL_CACHE.clear()

    from teamagent.ingest.pipeline import _ingest_slack_channel

    spec = SlackChannelSpec(channel_id="C0CLS", channel_name="#cls", description="")
    repo = _FakeRepository()
    docs_n, _ = _ingest_slack_channel(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r-slack-cls",
    )
    assert docs_n == 1
    assert stub.calls
    md = repo.upsert_calls[0]["metadata"]
    _assert_classified(md)
    # 既存キーは cls_* マージで破壊されない（後方互換）
    assert md["channel_name"] == "#cls"


def test_ingest_gsheet_applies_classification(monkeypatch: pytest.MonkeyPatch) -> None:
    """gsheet 経路でも各行に cls_* が乗る（contextualizer は付けない）。"""
    from teamagent.adapters.gsheets_client import TabRows

    fake_client = MagicMock()
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1V",
        tab_name="フォーム回答 1",
        headers=("業界", "温度感"),
        rows=(("飲食", "高"),),
        row_count=1,
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    stub = _StubClassifier()
    monkeypatch.setattr("teamagent.ingest.classify.build_classifier_from_env", lambda: stub)

    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1V",
        sheet_name="FB",
        description="",
        tabs=(GSheetsTabSpec(gid=537831563, tab_name="フォーム回答 1"),),
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r-gsheet-cls",
    )
    assert docs_n == 1
    assert stub.calls
    md = repo.upsert_calls[0]["metadata"]
    _assert_classified(md)
    # 既存キーは cls_* マージで破壊されない（後方互換）
    assert md["tab_name"] == "フォーム回答 1"
    assert md["row_idx"] == 2


def test_ingest_gsheet_fb_metadata_coexists_with_classification(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """FB シートで FB メタと cls_* が共存する（fb → cls の順で合成・cls_* を上書きしない）。

    deal_phase（フォーム実値）と cls_phase（Haiku 分類）は別キーなので両立する。
    """
    from teamagent.adapters.gsheets_client import TabRows

    fake_client = MagicMock()
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1V",
        tab_name="フォーム回答 1",
        headers=_FB_SHEET_HEADERS,
        rows=(
            (
                "2026/06/30 10:00:00",
                "直販",
                "アルファ広告社",
                "シータ食品",
                "2回目以降提案",
                "B（前向き）",
                "施策の考え方は好評",
                "【メディア】グレースモード",
            ),
        ),
        row_count=1,
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    stub = _StubClassifier()
    monkeypatch.setattr("teamagent.ingest.classify.build_classifier_from_env", lambda: stub)

    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1V",
        sheet_name="ショート動画営業 FB - フォーム回答",
        description="",
        tabs=(GSheetsTabSpec(gid=537831563, tab_name="フォーム回答 1"),),
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r-gsheet-fb-cls",
    )
    assert docs_n == 1
    md = repo.upsert_calls[0]["metadata"]
    # cls_* は FB メタに上書きされない（stub の値のまま）
    _assert_classified(md)
    # FB メタも同時に付く（deal_phase=フォーム実値、cls_phase=分類値で両立）
    assert md["is_sales_fb"] is True
    assert md["channel_type"] == "直販"
    assert md["client_case"] == "シータ食品"
    assert md["deal_phase"] == "2回目以降提案"
    assert md["cls_phase"] == "提案"
    assert md["client_name"] == "シータ食品"


def test_ingest_crawl_no_classification_when_disabled(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """classifier None（USE_DOC_CLASSIFY OFF 相当）なら crawl 経路に cls_* は出ない（後方互換）。"""
    from teamagent.adapters.gdrive_client import SharedDrive

    fake_client = MagicMock()
    fake_client.list_shared_drives.return_value = [SharedDrive(id="D1", name="営業ナレッジ")]
    fake_client.walk_files_recursive.return_value = [
        _make_drive_file(id="F1", name="提案.pdf", owners=("alice@x.jp",))
    ]
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "alice@x.jp")]
    fake_client.download_file_bytes.return_value = b"<fake-pdf>"
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    fake_reader = MagicMock()
    page = MagicMock()
    page.extract_text.return_value = "本文" * 50
    fake_reader.pages = [page]
    monkeypatch.setattr("pypdf.PdfReader", lambda _s: fake_reader)

    monkeypatch.setattr("teamagent.ingest.classify.build_classifier_from_env", lambda: None)

    from teamagent.ingest.loader import SharedDriveCrawlSpec
    from teamagent.ingest.pipeline import _ingest_shared_drives_crawl

    spec = SharedDriveCrawlSpec(
        enabled=True,
        name_filter=("営業",),
        sales_relevance_filter=False,
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_shared_drives_crawl(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-crawl-nocls",
    )
    assert docs_n == 1
    md = repo.upsert_calls[0]["metadata"]
    assert "cls_project" not in md
    assert "cls_doc_type" not in md


# -----------------------------------------------------------
# crawl 経路の Office (pptx/docx/xlsx) 本文抽出の配線
# （高ボリューム crawl が title だけでなく中身も index 化することの回帰。
#   営業提案書は大半が pptx なので、ここが title_only に落ちると検索が死ぬ。）
# -----------------------------------------------------------
def test_ingest_shared_drives_crawl_pptx_extracts_content_chunks(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """crawl 経路でも pptx は slide 単位の本文 chunk になる（title_only に落ちない）。"""
    from io import BytesIO

    from pptx import Presentation

    from teamagent.adapters.gdrive_client import SharedDrive
    from teamagent.ingest.office_extract import PPTX_MIME

    prs = Presentation()
    layout = prs.slide_layouts[5]
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "アース製薬向け SNS 提案 スライド1"
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "施策概要 スライド2"
    buf = BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    fake_client = MagicMock()
    fake_client.list_shared_drives.return_value = [SharedDrive(id="D1", name="営業ナレッジ")]
    fake_client.walk_files_recursive.return_value = [
        _make_drive_file(id="P1", name="提案.pptx", mime=PPTX_MIME, owners=("alice@x.jp",))
    ]
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "alice@x.jp")]
    fake_client.download_file_bytes.return_value = pptx_bytes
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.loader import SharedDriveCrawlSpec
    from teamagent.ingest.pipeline import _ingest_shared_drives_crawl

    spec = SharedDriveCrawlSpec(
        enabled=True,
        name_filter=("営業",),
        sales_relevance_filter=False,  # フィルタを外して fake pptx を確実に通す
    )
    repo = _FakeRepository()
    docs_n, chunks_n = _ingest_shared_drives_crawl(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-crawl-pptx",
    )
    assert docs_n == 1
    assert chunks_n == 2  # 2 slides → 2 本文 chunk（小さいので分割なし）
    fake_client.download_file_bytes.assert_called_once()
    chunks = repo.upsert_calls[0]["chunks"]
    assert len(chunks) == 2
    # 本文抽出 chunk であることの検証: page_num を持ち、title_only に落ちていない
    for c in chunks:
        assert "page_num" in c.metadata
        assert not c.metadata.get("title_only")


def test_ingest_shared_drives_crawl_office_badzip_is_skipped(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """壊れた pptx (BadZipFile) は fail-open で skip され、crawl は他ファイルを処理し続ける。"""
    import zipfile

    from teamagent.adapters.gdrive_client import SharedDrive
    from teamagent.ingest.office_extract import PPTX_MIME

    fake_client = MagicMock()
    fake_client.list_shared_drives.return_value = [SharedDrive(id="D1", name="営業ナレッジ")]
    # 1 件目: 壊れた pptx（skip される）、2 件目: 正常 PDF（処理される）
    fake_client.walk_files_recursive.return_value = [
        _make_drive_file(id="BAD", name="corrupt.pptx", mime=PPTX_MIME, owners=("alice@x.jp",)),
        _make_drive_file(id="OK", name="ok.pdf", owners=("alice@x.jp",)),
    ]
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "alice@x.jp")]
    fake_client.download_file_bytes.return_value = b"<bytes>"
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    # extract_office_pages は pipeline 内で lazy import されるので、元モジュールを差し替える。
    def _boom(_data: bytes, *, mime_type: str) -> list[tuple[int, str]]:
        raise zipfile.BadZipFile("File is not a zip file")

    monkeypatch.setattr("teamagent.ingest.office_extract.extract_office_pages", _boom)

    # 正常 PDF 側の pypdf を擬似
    fake_reader = MagicMock()
    page = MagicMock()
    page.extract_text.return_value = "正常な PDF 本文" * 20
    fake_reader.pages = [page]
    monkeypatch.setattr("pypdf.PdfReader", lambda _s: fake_reader)

    from teamagent.ingest.loader import SharedDriveCrawlSpec
    from teamagent.ingest.pipeline import _ingest_shared_drives_crawl

    spec = SharedDriveCrawlSpec(
        enabled=True,
        name_filter=("営業",),
        sales_relevance_filter=False,
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_shared_drives_crawl(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-crawl-badzip",
    )
    # 壊れた pptx は skip、正常 PDF だけが doc 化されて crawl は止まらない
    assert docs_n == 1
    assert repo.upsert_calls[0]["external_id"] == "OK"


# -----------------------------------------------------------
# テンプレ検出（boilerplate）配線: env ゲートで mark_boilerplate を呼ぶ／呼ばない
# -----------------------------------------------------------
class _NoopOpsCursor:
    """_disable_statement_timeout が ``SET LOCAL statement_timeout`` を発行できるよう、
    cursor の context manager を最小実装した fake（実行内容は記録するだけ）。"""

    def __init__(self, executed: list[str]) -> None:
        self._executed = executed

    def __enter__(self) -> _NoopOpsCursor:
        return self

    def __exit__(self, *exc: object) -> None:
        return None

    def execute(self, sql: str, params: object = None) -> None:
        self._executed.append(sql)


class _FakeOpsConn:
    """``_ops_connection()`` が返す接続本体の fake。``cursor()`` で no-op cursor を返す。

    H1 の ``_disable_statement_timeout(conn)`` が ``conn.cursor()`` を呼ぶため、
    plain ``object()`` では落ちる。実 SQL は走らせないが、発行された SQL は
    ``executed`` に記録して検証可能にする。"""

    def __init__(self) -> None:
        self.executed: list[str] = []

    def cursor(self) -> _NoopOpsCursor:
        return _NoopOpsCursor(self.executed)


class _OpsConnCtx:
    """repository._ops_connection() が返す context manager の fake。"""

    def __init__(self, conn: object) -> None:
        self._conn = conn

    def __enter__(self) -> object:
        return self._conn

    def __exit__(self, *exc: object) -> None:
        return None


class _RepoWithOpsConn(_FakeRepository):
    """_ops_connection() を持つ fake repository（boilerplate 配線テスト用）。"""

    def __init__(self) -> None:
        super().__init__()
        self.ops_conn = _FakeOpsConn()
        self.ops_conn_calls = 0

    def _ops_connection(self) -> _OpsConnCtx:
        self.ops_conn_calls += 1
        return _OpsConnCtx(self.ops_conn)


def _empty_sources() -> IngestSources:
    return IngestSources(version=1, slack_channels=(), gdrive_folders=(), gsheets=())


def test_boilerplate_not_called_when_env_off(monkeypatch: pytest.MonkeyPatch) -> None:
    """BOILERPLATE_DETECT 未設定なら mark_boilerplate は呼ばれない＝現行と完全一致。"""
    monkeypatch.delenv("BOILERPLATE_DETECT", raising=False)
    calls: list[int] = []
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_boilerplate",
        lambda conn, *, min_docs: calls.append(min_docs) or 0,
    )
    repo = _RepoWithOpsConn()
    runner = IngestRunner(
        repository=repo,  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
    )
    runner.run(_empty_sources())
    assert calls == []
    assert repo.ops_conn_calls == 0


def test_boilerplate_not_called_in_dry_run(monkeypatch: pytest.MonkeyPatch) -> None:
    """dry-run では env ON でも DB を書かないので呼ばない。"""
    monkeypatch.setenv("BOILERPLATE_DETECT", "1")
    calls: list[int] = []
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_boilerplate",
        lambda conn, *, min_docs: calls.append(min_docs) or 0,
    )
    runner = IngestRunner(
        repository=_RepoWithOpsConn(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=True,
    )
    runner.run(_empty_sources())
    assert calls == []


def test_boilerplate_called_when_env_on(monkeypatch: pytest.MonkeyPatch) -> None:
    """BOILERPLATE_DETECT=1 + commit で mark_boilerplate が ops 接続付きで呼ばれる。"""
    monkeypatch.setenv("BOILERPLATE_DETECT", "1")
    monkeypatch.delenv("BOILERPLATE_MIN_DOCS", raising=False)
    seen: list[tuple[object, int]] = []
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_boilerplate",
        lambda conn, *, min_docs, **_: seen.append((conn, min_docs)) or 7,
    )
    repo = _RepoWithOpsConn()
    runner = IngestRunner(
        repository=repo,  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
    )
    runner.run(_empty_sources())
    assert len(seen) == 1
    conn, min_docs = seen[0]
    assert conn is repo.ops_conn  # 既存の admin role 接続を再利用
    assert min_docs == 3  # 既定値
    assert repo.ops_conn_calls == 1
    # H1: mark 実行前に当該 tx 内で statement_timeout を無制限にしている。
    assert any("SET LOCAL statement_timeout" in s for s in repo.ops_conn.executed)


def test_boilerplate_min_docs_env_override(monkeypatch: pytest.MonkeyPatch) -> None:
    """BOILERPLATE_MIN_DOCS で閾値を上書きできる。"""
    monkeypatch.setenv("BOILERPLATE_DETECT", "true")
    monkeypatch.setenv("BOILERPLATE_MIN_DOCS", "5")
    seen: list[int] = []
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_boilerplate",
        lambda conn, *, min_docs, **_: seen.append(min_docs) or 0,
    )
    runner = IngestRunner(
        repository=_RepoWithOpsConn(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
    )
    runner.run(_empty_sources())
    assert seen == [5]


def test_boilerplate_failure_is_fail_open(monkeypatch: pytest.MonkeyPatch) -> None:
    """mark_boilerplate が例外でも run() は成功して結果を返す（fail-open）。"""
    monkeypatch.setenv("BOILERPLATE_DETECT", "1")

    def _boom(conn: object, *, min_docs: int, **_: object) -> int:
        raise RuntimeError("db down")

    monkeypatch.setattr("teamagent.ingest.pipeline.mark_boilerplate", _boom)
    runner = IngestRunner(
        repository=_RepoWithOpsConn(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
    )
    result = runner.run(_empty_sources())  # 例外を投げない
    assert isinstance(result, IngestResult)
    assert result.total_documents() == 0


# -----------------------------------------------------------
# 資料まるごと重複排除（docdedup）配線: env ゲートで mark_duplicate_documents を呼ぶ／呼ばない
# -----------------------------------------------------------
def test_docdedup_not_called_when_env_off(monkeypatch: pytest.MonkeyPatch) -> None:
    """DOC_DEDUP_DETECT 未設定なら mark_duplicate_documents は呼ばれない＝現行と完全一致。"""
    monkeypatch.delenv("DOC_DEDUP_DETECT", raising=False)
    monkeypatch.delenv("BOILERPLATE_DETECT", raising=False)
    calls: list[float] = []
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_duplicate_documents",
        lambda conn, *, jaccard_threshold: calls.append(jaccard_threshold) or 0,
    )
    repo = _RepoWithOpsConn()
    runner = IngestRunner(
        repository=repo,  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
    )
    runner.run(_empty_sources())
    assert calls == []
    assert repo.ops_conn_calls == 0


def test_docdedup_not_called_in_dry_run(monkeypatch: pytest.MonkeyPatch) -> None:
    """dry-run では env ON でも DB を書かないので呼ばない。"""
    monkeypatch.setenv("DOC_DEDUP_DETECT", "1")
    calls: list[float] = []
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_duplicate_documents",
        lambda conn, *, jaccard_threshold: calls.append(jaccard_threshold) or 0,
    )
    runner = IngestRunner(
        repository=_RepoWithOpsConn(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=True,
    )
    runner.run(_empty_sources())
    assert calls == []


def test_docdedup_called_when_env_on(monkeypatch: pytest.MonkeyPatch) -> None:
    """DOC_DEDUP_DETECT=1 + commit で mark_duplicate_documents が ops 接続付きで呼ばれる。"""
    monkeypatch.setenv("DOC_DEDUP_DETECT", "1")
    monkeypatch.delenv("DOC_DEDUP_JACCARD", raising=False)
    monkeypatch.delenv("BOILERPLATE_DETECT", raising=False)
    seen: list[tuple[object, float]] = []
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_duplicate_documents",
        lambda conn, *, jaccard_threshold, **_: seen.append((conn, jaccard_threshold)) or 4,
    )
    repo = _RepoWithOpsConn()
    runner = IngestRunner(
        repository=repo,  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
    )
    runner.run(_empty_sources())
    assert len(seen) == 1
    conn, jaccard_threshold = seen[0]
    assert conn is repo.ops_conn  # 既存の admin role 接続を再利用
    assert jaccard_threshold == 0.7  # 既定値
    assert repo.ops_conn_calls == 1
    # H1: mark 実行前に当該 tx 内で statement_timeout を無制限にしている。
    assert any("SET LOCAL statement_timeout" in s for s in repo.ops_conn.executed)


def test_docdedup_jaccard_env_override(monkeypatch: pytest.MonkeyPatch) -> None:
    """DOC_DEDUP_JACCARD でしきい値を上書きできる。"""
    monkeypatch.setenv("DOC_DEDUP_DETECT", "true")
    monkeypatch.setenv("DOC_DEDUP_JACCARD", "0.85")
    seen: list[float] = []
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_duplicate_documents",
        lambda conn, *, jaccard_threshold, **_: seen.append(jaccard_threshold) or 0,
    )
    runner = IngestRunner(
        repository=_RepoWithOpsConn(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
    )
    runner.run(_empty_sources())
    assert seen == [0.85]


def test_docdedup_failure_is_fail_open(monkeypatch: pytest.MonkeyPatch) -> None:
    """mark_duplicate_documents が例外でも run() は成功して結果を返す（fail-open）。"""
    monkeypatch.setenv("DOC_DEDUP_DETECT", "1")

    def _boom(conn: object, *, jaccard_threshold: float, **_: object) -> int:
        raise RuntimeError("db down")

    monkeypatch.setattr("teamagent.ingest.pipeline.mark_duplicate_documents", _boom)
    runner = IngestRunner(
        repository=_RepoWithOpsConn(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
    )
    result = runner.run(_empty_sources())  # 例外を投げない
    assert isinstance(result, IngestResult)
    assert result.total_documents() == 0


def test_boilerplate_runs_after_docdedup(monkeypatch: pytest.MonkeyPatch) -> None:
    """両 env ON のとき docdedup → boilerplate の順で呼ばれる（M3・同じ ops 接続を各々取得）。

    M3: boilerplate の指紋集計は suppressed（非正本）doc を母数から外すため、同一 run 内で
    先に docdedup が suppressed を確定させてから boilerplate を走らせる必要がある。
    """
    monkeypatch.setenv("BOILERPLATE_DETECT", "1")
    monkeypatch.setenv("DOC_DEDUP_DETECT", "1")
    order: list[str] = []
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_boilerplate",
        lambda conn, *, min_docs, **_: order.append("boilerplate") or 0,
    )
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_duplicate_documents",
        lambda conn, *, jaccard_threshold, **_: order.append("docdedup") or 0,
    )
    repo = _RepoWithOpsConn()
    runner = IngestRunner(
        repository=repo,  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
    )
    runner.run(_empty_sources())
    assert order == ["docdedup", "boilerplate"]
    assert repo.ops_conn_calls == 2


# -----------------------------------------------------------
# M5: 非数値 env でも crash せず default にフォールバック（_envint / _envfloat）
# -----------------------------------------------------------
def test_envint_falls_back_on_blank_and_nonnumeric(monkeypatch: pytest.MonkeyPatch) -> None:
    """空文字・非数値・未設定はすべて default に倒れ、有効値は反映される（ingest CRASH 防止）。"""
    from teamagent.ingest.pipeline import _envint

    monkeypatch.delenv("X_INT", raising=False)
    assert _envint("X_INT", 7) == 7  # 未設定
    monkeypatch.setenv("X_INT", "")
    assert _envint("X_INT", 7) == 7  # 空文字
    monkeypatch.setenv("X_INT", "   ")
    assert _envint("X_INT", 7) == 7  # 空白のみ
    monkeypatch.setenv("X_INT", "3x")
    assert _envint("X_INT", 7) == 7  # 非数値
    monkeypatch.setenv("X_INT", "11")
    assert _envint("X_INT", 7) == 11  # 有効値


def test_envfloat_falls_back_on_blank_and_nonnumeric(monkeypatch: pytest.MonkeyPatch) -> None:
    """空文字・非数値・未設定はすべて default に倒れ、有効値は反映される（ingest CRASH 防止）。"""
    from teamagent.ingest.pipeline import _envfloat

    monkeypatch.delenv("X_FLOAT", raising=False)
    assert _envfloat("X_FLOAT", 0.7) == 0.7  # 未設定
    monkeypatch.setenv("X_FLOAT", "")
    assert _envfloat("X_FLOAT", 0.7) == 0.7  # 空文字
    monkeypatch.setenv("X_FLOAT", "abc")
    assert _envfloat("X_FLOAT", 0.7) == 0.7  # 非数値
    monkeypatch.setenv("X_FLOAT", "0.85")
    assert _envfloat("X_FLOAT", 0.7) == 0.85  # 有効値


def test_runner_does_not_crash_on_nonnumeric_detection_envs(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """非数値 BOILERPLATE_MIN_DOCS / DOC_DEDUP_JACCARD でも run() は CRASH せず default で進む（M5）。"""
    monkeypatch.setenv("BOILERPLATE_DETECT", "1")
    monkeypatch.setenv("BOILERPLATE_MIN_DOCS", "")  # 空文字 → 3
    monkeypatch.setenv("BOILERPLATE_MIN_CHARS", "x")  # 非数値 → 40
    monkeypatch.setenv("DOC_DEDUP_DETECT", "1")
    monkeypatch.setenv("DOC_DEDUP_JACCARD", "3x")  # 非数値 → 0.7
    monkeypatch.setenv("DOC_DEDUP_MAX_DOCS", "")  # 空文字 → 5000
    seen_bp: list[tuple[int, int]] = []
    seen_dd: list[float] = []
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_boilerplate",
        lambda conn, *, min_docs, min_chars, **_: seen_bp.append((min_docs, min_chars)) or 0,
    )
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.mark_duplicate_documents",
        lambda conn, *, jaccard_threshold, **_: seen_dd.append(jaccard_threshold) or 0,
    )
    runner = IngestRunner(
        repository=_RepoWithOpsConn(),  # type: ignore[arg-type]
        embedder=_FakeEmbedder(),
        owner_email="x@y.jp",
        dry_run=False,
    )
    runner.run(_empty_sources())  # 例外を投げない
    assert seen_bp == [(3, 40)]  # 既定値にフォールバック
    assert seen_dd == [0.7]


# -----------------------------------------------------------
# L4: _envflag が末尾空白を strip して判定する
# -----------------------------------------------------------
def test_envflag_strips_trailing_whitespace(monkeypatch: pytest.MonkeyPatch) -> None:
    """末尾改行/空白付きの "true\\n" 等でも ON と判定される（skill._envflag と同流儀）。"""
    from teamagent.ingest.pipeline import _envflag

    monkeypatch.setenv("X_FLAG", "true\n")
    assert _envflag("X_FLAG") is True
    monkeypatch.setenv("X_FLAG", "  1  ")
    assert _envflag("X_FLAG") is True
    monkeypatch.setenv("X_FLAG", "yes\t")
    assert _envflag("X_FLAG") is True
    monkeypatch.setenv("X_FLAG", " false ")
    assert _envflag("X_FLAG") is False


# -----------------------------------------------------------
# M4: crawl の file ループは file 単位 try/except で 1 file 例外を吸収する
# -----------------------------------------------------------
def test_crawl_one_file_exception_does_not_kill_source(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """crawl で 1 file の embed 例外が source 全体を落とさず、後続 file は処理される（M4）。

    folder 経路（pipeline.py ~838）には file 単位 try/except があるが crawl 経路には無く、
    embed / DB upsert の例外で残り全ファイルが道連れになっていた。file 単位で skip して継続する。
    """
    from teamagent.adapters.gdrive_client import SharedDrive
    from teamagent.ingest.loader import SharedDriveCrawlSpec
    from teamagent.ingest.pipeline import _ingest_shared_drives_crawl

    fake_client = MagicMock()
    fake_client.list_shared_drives.return_value = [SharedDrive(id="D1", name="営業ナレッジ")]
    # 1 件目: embed が落ちる title_only file（旧バイナリ等）、2 件目: 正常 file。
    fake_client.walk_files_recursive.return_value = [
        _make_drive_file(
            id="BOOM", name="boom.bin", mime="application/x-unknown", owners=("a@x.jp",)
        ),
        _make_drive_file(id="OK", name="ok.bin", mime="application/x-unknown", owners=("a@x.jp",)),
    ]
    fake_client.list_permissions.return_value = [_make_drive_perm("user", "owner", "a@x.jp")]
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    class _BoomOnFirstEmbedder:
        """最初の file（external_id=BOOM）の title_only テキストだけ embed で例外を投げる。"""

        def embed(self, text: str) -> list[float]:
            if "boom.bin" in text:
                raise RuntimeError("embed exploded")
            return [0.1] * 1024

        def embed_passage(self, text: str) -> list[float]:
            return self.embed(text)

    spec = SharedDriveCrawlSpec(
        enabled=True,
        name_filter=("営業",),
        sales_relevance_filter=False,
    )
    repo = _FakeRepository()
    # 例外を投げず、正常 file (OK) だけが doc 化される。
    docs_n, _ = _ingest_shared_drives_crawl(
        spec,
        embedder=_BoomOnFirstEmbedder(),  # type: ignore[arg-type]
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@x.jp",
        dry_run=False,
        request_id="r-crawl-boom",
    )
    assert docs_n == 1
    assert [c["external_id"] for c in repo.upsert_calls] == ["OK"]


def test_ingest_gsheet_resolves_renamed_tab_by_gid(monkeypatch: pytest.MonkeyPatch) -> None:
    """タブがリネームされても gid で現タイトルを解決して取り込める(本番障害の回帰テスト)。

    2026-07-06: 命名ルール導入でタブ名が変わり "Unable to parse range: 'フォーム回答 1'"
    で全シート取り込みが失敗した。gid は不変なので metadata から現タイトルを解決する。
    """
    from teamagent.adapters.gsheets_client import SheetMetadata, SheetTab, TabRows

    fake_client = MagicMock()
    # yaml 上は「フォーム回答 1」だが、実シートは「営業FB_回答」にリネーム済み
    fake_client.get_sheet_metadata.return_value = SheetMetadata(
        sheet_id="1V",
        title="FB",
        tabs=(
            SheetTab(sheet_id="1V", gid=537831563, title="営業FB_回答", row_count=2, col_count=2),
        ),
    )
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1V",
        tab_name="営業FB_回答",
        headers=("業界", "温度感"),
        rows=(("飲食", "高"),),
        row_count=1,
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )

    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1V",
        sheet_name="FB",
        description="",
        tabs=(GSheetsTabSpec(gid=537831563, tab_name="フォーム回答 1"),),
    )
    repo = _FakeRepository()
    docs_n, _ = _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r",
    )
    assert docs_n == 1
    # get_tab_rows は yaml の旧名でなく gid 解決後の現タイトルで呼ばれる
    _, kwargs = fake_client.get_tab_rows.call_args
    assert kwargs["tab_name"] == "営業FB_回答"
    # external_id は gid ベースのまま(リネームしても冪等 upsert)
    assert "537831563" in repo.upsert_calls[0]["external_id"]


# 実シート dump(2026-07-15・gid 1962561294) のヘッダ準拠。連番/保存ファイル（リンク付き）/
# 旧_保管先フォルダ は ファイル記録タブにのみ存在し、フォーム回答タブには無い。
_FILEREC_HEADERS = (
    "ファイルをアップ",
    "正式社名",
    "案件名",
    "クライアント種別",
    "提案プロダクト",
    "資料の概要_メイン",
    "送信者",
    "連番",
    "保存ファイル（リンク付き）",
    "旧_保管先フォルダ",
)


def _filerec_sheet(monkeypatch: pytest.MonkeyPatch, rows: tuple[tuple[str, ...], ...]) -> None:
    from teamagent.adapters.gsheets_client import SheetMetadata, SheetTab, TabRows

    fake_client = MagicMock()
    fake_client.get_sheet_metadata.return_value = SheetMetadata(
        sheet_id="1J",
        title="ナレッジ共有",
        tabs=(
            SheetTab(
                sheet_id="1J",
                gid=1962561294,
                title="ファイル記録",
                row_count=len(rows) + 1,
                col_count=len(_FILEREC_HEADERS),
            ),
        ),
    )
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1J",
        tab_name="ファイル記録",
        headers=_FILEREC_HEADERS,
        rows=rows,
        row_count=len(rows),
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    monkeypatch.delenv("USE_DOC_CLASSIFY", raising=False)


def _run_filerec() -> _FakeRepository:
    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1J",
        sheet_name="ナレッジ共有",
        description="",
        tabs=(GSheetsTabSpec(gid=1962561294, tab_name="ファイル記録"),),
    )
    repo = _FakeRepository()
    _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r",
    )
    return repo


_ROW_A = (
    "https://slack/f/A",
    "株式会社デルタ製薬",
    "新製品PR",
    "メーカー",
    "タテガタ",
    "提案",
    "@yamada",
    "20250617-001",
    "提案_デルタ製薬_x.pdf",
    "20250617_デルタ製薬",
)
_ROW_B = (
    "https://slack/f/B",
    "ゼータ工業",
    "採用支援",
    "メーカー",
    "NCS",
    "レポート",
    "@sato",
    "20250625-002",
    "レポート_ゼータ_y.pdf",
    "20250625_ゼータ工業",
)


def test_ingest_gsheet_uses_seq_as_stable_id(monkeypatch: pytest.MonkeyPatch) -> None:
    """ファイル記録行は「連番」を identity にする＝行位置が変わっても同一 document（Codex #215-4）。

    実 dump で 連番 は 142/142 ユニーク・重複 0。行番号 ID だと並べ替え/行挿入で上書き先がズレた。
    """
    _filerec_sheet(monkeypatch, (_ROW_A, _ROW_B))
    first = {c["metadata"]["client_company"]: c["external_id"] for c in _run_filerec().upsert_calls}
    # 行を入れ替えて再取込しても external_id は不変（＝上書き先がズレない）
    _filerec_sheet(monkeypatch, (_ROW_B, _ROW_A))
    second = {
        c["metadata"]["client_company"]: c["external_id"] for c in _run_filerec().upsert_calls
    }
    assert first == second
    assert all(v.startswith("1J:1962561294:k") for v in first.values())
    assert len(set(first.values())) == 2  # 別行は別 document（衝突で潰れない）


def test_ingest_gsheet_title_uses_company_and_case_not_row_n(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """title は "row N" でなく「正式社名 案件名」（法人格は除去・Codex #215-2）。"""
    _filerec_sheet(monkeypatch, (_ROW_A,))
    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1J",
        sheet_name="ナレッジ共有",
        description="",
        tabs=(GSheetsTabSpec(gid=1962561294, tab_name="ファイル記録"),),
    )
    seen: list[str] = []

    class _TitleRepo(_FakeRepository):
        def upsert_document_with_chunks(self, doc, chunks, request_id, **kw):  # type: ignore[no-untyped-def]
            seen.append(doc.title)
            return super().upsert_document_with_chunks(doc, chunks, request_id, **kw)

    _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=_TitleRepo(),  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r",
    )
    assert seen == ["デルタ製薬 新製品PR"]
    assert "row " not in seen[0]


def test_ingest_gsheet_skips_ichiji_souko_rows(monkeypatch: pytest.MonkeyPatch) -> None:
    """99_一次倉庫系（検索対象外の生データ置き場）の行は取り込まない（Codex #215-5）。

    実 dump では 99_ は 0 件＝将来行への保険。判定は gdrive と同一 regex を再利用する。
    """
    drop = (
        "https://slack/f/B",
        "ゼータ工業",
        "採用",
        "メーカー",
        "NCS",
        "提案",
        "@sato",
        "20250625-002",
        "raw.pdf",
        "99_一次倉庫",  # 旧_保管先フォルダ が一次倉庫 → 取り込まない
    )
    _filerec_sheet(monkeypatch, (_ROW_A, drop))
    repo = _run_filerec()
    assert len(repo.upsert_calls) == 1
    assert repo.upsert_calls[0]["metadata"]["client_company"] == "株式会社デルタ製薬"


def test_ingest_gsheet_without_seq_keeps_row_id(monkeypatch: pytest.MonkeyPatch) -> None:
    """連番を持たないシート（フォーム回答タブ等）は従来の行番号 ID のまま＝既存 document を孤児化しない。

    gsheets には stale 検出が無いため、ID 形式を変えると旧 doc が永久残存し二重表示になる。
    """
    from teamagent.adapters.gsheets_client import SheetMetadata, SheetTab, TabRows

    headers = ("正式社名", "案件名", "クライアント種別", "提案プロダクト", "資料の概要", "送信者")
    rows = (("株式会社デルタ製薬", "新製品PR", "メーカー", "タテガタ", "提案", "@yamada"),)
    fake_client = MagicMock()
    fake_client.get_sheet_metadata.return_value = SheetMetadata(
        sheet_id="1J",
        title="ナレッジ共有",
        tabs=(
            SheetTab(
                sheet_id="1J",
                gid=278789217,
                title="フォーム回答",
                row_count=2,
                col_count=len(headers),
            ),
        ),
    )
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1J", tab_name="フォーム回答", headers=headers, rows=rows, row_count=1
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    monkeypatch.delenv("USE_DOC_CLASSIFY", raising=False)

    from teamagent.ingest.pipeline import _ingest_gsheet

    spec = GSheetSpec(
        sheet_id="1J",
        sheet_name="ナレッジ共有",
        description="",
        tabs=(GSheetsTabSpec(gid=278789217, tab_name="フォーム回答"),),
    )
    repo = _FakeRepository()
    _ingest_gsheet(
        spec,
        embedder=_FakeEmbedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="x@y.jp",
        dry_run=False,
        request_id="r",
    )
    assert repo.upsert_calls[0]["external_id"] == "1J:278789217:2"  # 行番号 ID のまま（形式不変）


def test_ingest_gsheet_body_drops_ops_columns_so_excerpt_keeps_knowledge(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """本文から運用列(Slack URL 等)を外す＝下流の 160 字抜粋に「知見の列」が載る（Codex #215-160字）。

    ファイル記録は先頭列が長大な Slack file URL のため、素の整形だと export_vault の
    left(...,160) 抜粋が URL で埋まり、ポイント/なぜ/フリーコメント が /app のタグ源から落ちる
    （実測で施策手法タグが 6→1 に激減した）。
    """
    headers = (*_FILEREC_HEADERS, "このナレッジのポイントはここ！")
    long_url = "https://news-tv.slack.com/files/UBNDFMCQ3/F091USMAVPW/" + "x" * 90
    row = (
        long_url,
        "デルタ製薬",
        "新製品PR",
        "メーカー",
        "タテガタ",
        "提案",
        "@yamada",
        "20250617-001",
        "提案_x.pdf",
        "20250617_デルタ製薬",
        "UGCと切り抜きが刺さった",
    )
    from teamagent.adapters.gsheets_client import SheetMetadata, SheetTab, TabRows

    fake_client = MagicMock()
    fake_client.get_sheet_metadata.return_value = SheetMetadata(
        sheet_id="1J",
        title="ナレッジ共有",
        tabs=(
            SheetTab(
                sheet_id="1J",
                gid=1962561294,
                title="ファイル記録",
                row_count=2,
                col_count=len(headers),
            ),
        ),
    )
    fake_client.get_tab_rows.return_value = TabRows(
        sheet_id="1J", tab_name="ファイル記録", headers=headers, rows=(row,), row_count=1
    )
    monkeypatch.setattr(
        "teamagent.adapters.gsheets_client.GSheetsClient.from_env",
        classmethod(lambda cls, **kwargs: fake_client),
    )
    monkeypatch.delenv("USE_DOC_CLASSIFY", raising=False)
    repo = _run_filerec()
    body = repo.upsert_calls[0]["chunks"][0].content
    assert long_url not in body  # Slack file URL は本文に載せない
    assert "連番:" not in body and "保存ファイル" not in body  # GAS 運用列も載せない
    # 知見の列が先頭 160 字（＝export_vault の抜粋窓）に収まる
    assert "UGCと切り抜きが刺さった" in body[:160]
    assert "正式社名: デルタ製薬" in body
