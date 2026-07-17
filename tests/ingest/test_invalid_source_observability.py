"""Drive Office invalid payload の全経路・再試行・warning状態回帰テスト。"""

from __future__ import annotations

import hashlib
from dataclasses import replace
from io import BytesIO
from typing import Any
from unittest.mock import MagicMock, patch

import pytest
from structlog.testing import capture_logs

from teamagent.adapters.gdrive_client import (
    DriveFile,
    DrivePermission,
    GDriveTraversalIncompleteError,
    SharedDrive,
)
from teamagent.ingest.loader import (
    GDriveFolderSpec,
    IngestSources,
    SharedDriveCrawlSpec,
)
from teamagent.ingest.office_extract import (
    DOCX_MIME,
    GDOC_NATIVE_MIME,
    OFFICE_VALIDATOR_SCHEMA_VERSION,
    PPTX_MIME,
)
from teamagent.ingest.ops_alert import IngestOpsAlerter
from teamagent.ingest.pipeline import (
    ChunkUpsert,
    DocumentUpsert,
    IngestDurabilityError,
    IngestRunner,
    _guarded_upsert,
    _ingest_gdrive_folder,
    _ingest_shared_drives_crawl,
    _IngestWarningCollector,
    _process_one_gdrive_file,
    _send_ops_warning_summary,
)
from teamagent.ingest.repository import ConnectorState, SourceRetry


class _Embedder:
    def embed(self, text: str) -> list[float]:
        return [0.1] * 1024

    def embed_passage(self, text: str) -> list[float]:
        return self.embed(text)


class _Repository:
    def __init__(
        self,
        known: set[tuple[str, str, int, str, str]] | None = None,
    ) -> None:
        self.known = known or set()
        self.lookup_calls: list[
            tuple[str, str, str | None, int | None, str | None, str | None]
        ] = []
        self.invalid_records: list[dict[str, Any]] = []
        self.upserts: list[tuple[Any, list[Any]]] = []
        self.saved_states: list[dict[str, Any]] = []
        self.connector_runs: list[dict[str, Any]] = []
        self.retry_records: list[dict[str, Any]] = []
        self.retry_resolutions: list[dict[str, Any]] = []
        self.retry_claims: list[SourceRetry] = []
        self.retry_record_success = True
        self.retry_lease_success = True
        self.retry_lease_renewals: list[dict[str, Any]] = []
        self.connector_state: ConnectorState | None = None
        self.reconciliation_counts: dict[str, int] = {}
        self.reconciliation_resolutions: list[str] = []

    def find_invalid_source_reason(
        self,
        source_type: str,
        external_id: str,
        md5_checksum: str | None,
        size_bytes: int | None,
        mime_type: str | None = None,
        validator_schema_version: str | None = None,
    ) -> str | None:
        self.lookup_calls.append(
            (
                source_type,
                external_id,
                md5_checksum,
                size_bytes,
                mime_type,
                validator_schema_version,
            )
        )
        if md5_checksum is not None and size_bytes is not None:
            if (
                external_id,
                md5_checksum,
                size_bytes,
                str(mime_type),
                str(validator_schema_version),
            ) in self.known:
                return "corrupt_zip"
        return None

    def record_invalid_source(
        self,
        source_type: str,
        external_id: str,
        **kwargs: Any,
    ) -> bool:
        self.invalid_records.append(
            {"source_type": source_type, "external_id": external_id, **kwargs}
        )
        return True

    def upsert_document_with_chunks(
        self,
        doc: Any,
        chunks: list[Any],
        request_id: str,
        *,
        replace_existing_chunks: bool = True,
    ) -> str:
        self.upserts.append((doc, list(chunks)))
        return "doc-id"

    def load_connector_state(self, source_kind: str, source_id: str) -> ConnectorState | None:
        return self.connector_state

    def save_connector_state(self, source_kind: str, source_id: str, **kwargs: Any) -> None:
        self.saved_states.append({"source_kind": source_kind, "source_id": source_id, **kwargs})
        if kwargs.get("cursor"):
            self.connector_state = ConnectorState(
                source_kind=source_kind,
                source_id=source_id,
                cursor=str(kwargs["cursor"]),
                metadata=dict(kwargs.get("metadata") or {}),
            )

    def record_ingest_job(self, *args: Any, **kwargs: Any) -> None:
        return None

    def record_connector_run(self, **kwargs: Any) -> bool:
        self.connector_runs.append(kwargs)
        return True

    def claim_due_source_retries(self, **kwargs: Any) -> list[SourceRetry]:
        claimed = list(self.retry_claims)
        self.retry_claims.clear()
        return claimed

    def record_source_retry(self, **kwargs: Any) -> bool:
        self.retry_records.append(kwargs)
        if not self.retry_record_success:
            return False
        self.retry_claims = [
            SourceRetry(
                external_id=str(kwargs["external_id"]),
                md5_checksum=kwargs.get("md5_checksum"),
                size_bytes=kwargs.get("size_bytes"),
                mime_type=str(kwargs["mime_type"]),
                validator_schema_version=str(kwargs["validator_schema_version"]),
                attempt_count=len(self.retry_records),
                reason=str(kwargs["reason"]),
            )
        ]
        return True

    def renew_source_retry_lease(self, **kwargs: Any) -> bool:
        self.retry_lease_renewals.append(kwargs)
        return self.retry_lease_success

    def resolve_source_retry(self, **kwargs: Any) -> bool:
        self.retry_resolutions.append(kwargs)
        self.retry_claims.clear()
        return True

    def unresolved_reconciliation_counts(self, source_kind: str) -> dict[str, int]:
        return dict(self.reconciliation_counts)

    def resolve_reconciliation_gaps(
        self, *, source_kind: str, external_id: str, request_id: str
    ) -> int:
        self.reconciliation_resolutions.append(external_id)
        return 1


def _pptx_bytes() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "safe content"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _docx_bytes() -> bytes:
    from docx import Document

    document = Document()
    document.add_paragraph("safe document content")
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _file(
    file_id: str,
    data: bytes,
    *,
    md5_checksum: str | None = None,
) -> DriveFile:
    return DriveFile(
        id=file_id,
        name="confidential.pptx",
        mime_type=PPTX_MIME,
        modified_time="2026-07-17T00:00:00Z",
        size=len(data),
        md5_checksum=md5_checksum or hashlib.md5(data, usedforsecurity=False).hexdigest(),
        web_view_link=f"https://drive.google.com/file/d/{file_id}/view",
        owners_email=("owner@example.jp",),
    )


def _owner_permission() -> DrivePermission:
    return DrivePermission(
        id="permission",
        type="user",
        role="owner",
        email_address="owner@example.jp",
    )


def _folder_spec() -> GDriveFolderSpec:
    return GDriveFolderSpec(
        folder_id="FOLDER",
        folder_name="folder",
        description="",
        mime_type_filter=None,
    )


def test_corrupt_payload_is_persisted_without_document_mutation_or_identifier_logs() -> None:
    broken = b"PK\x03\x04" + (b"x" * 64)
    drive_file = _file("SENSITIVE-DRIVE-ID", broken)
    client = MagicMock()
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.return_value = broken
    repo = _Repository()
    collector = _IngestWarningCollector()
    skipped: list[str] = []

    with capture_logs() as logs:
        result = _process_one_gdrive_file(
            drive_file,
            _folder_spec(),
            client=client,
            embedder=_Embedder(),
            repository=repo,  # type: ignore[arg-type]
            owner_email="bot@example.jp",
            dry_run=False,
            request_id="req",
            skipped=skipped,
            warning_collector=collector,
            warning_source_id="FOLDER",
        )

    assert result == (0, 0)
    assert repo.upserts == []
    assert len(repo.invalid_records) == 1
    assert repo.invalid_records[0]["reason"] == "corrupt_zip"
    assert repo.invalid_records[0]["md5_checksum"] == drive_file.md5_checksum
    warning = next(log for log in logs if log["event"] == "gdrive_office_payload_invalid")
    assert warning["existing_document_preserved"] is True
    assert "file_id" not in warning
    assert "file_name" not in warning
    assert drive_file.id not in str(warning)
    snapshot = collector.snapshot("gdrive", "FOLDER")
    assert snapshot.reasons == {"corrupt_zip": 1}


def test_known_invalid_suppresses_acl_and_body_download_and_touches_observation() -> None:
    broken = b"PK\x03\x04" + (b"x" * 64)
    drive_file = _file("KNOWN", broken)
    repo = _Repository(
        known={
            (
                "KNOWN",
                str(drive_file.md5_checksum),
                len(broken),
                drive_file.mime_type,
                OFFICE_VALIDATOR_SCHEMA_VERSION,
            )
        }
    )
    client = MagicMock()
    collector = _IngestWarningCollector()

    result = _process_one_gdrive_file(
        drive_file,
        _folder_spec(),
        client=client,
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="req",
        skipped=[],
        warning_collector=collector,
        warning_source_id="FOLDER",
    )

    assert result == (0, 0)
    client.list_permissions.assert_not_called()
    client.download_file_bytes.assert_not_called()
    assert len(repo.invalid_records) == 1
    snapshot = collector.snapshot("gdrive", "FOLDER")
    assert snapshot.reasons == {"corrupt_zip": 1}
    assert snapshot.suppressed == 1


@pytest.mark.parametrize("change", ["md5", "size"])
def test_fingerprint_change_forces_revalidation(change: str) -> None:
    data = _pptx_bytes()
    drive_file = _file("RECOVERED", data)
    current_md5 = str(drive_file.md5_checksum)
    known_md5 = "0" * 32 if change == "md5" else current_md5
    known_size = len(data) if change == "md5" else len(data) - 1
    repo = _Repository(
        known={
            (
                "RECOVERED",
                known_md5,
                known_size,
                drive_file.mime_type,
                OFFICE_VALIDATOR_SCHEMA_VERSION,
            )
        }
    )
    client = MagicMock()
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.return_value = data

    result = _process_one_gdrive_file(
        drive_file,
        _folder_spec(),
        client=client,
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="req",
        skipped=[],
    )

    assert result[0] == 1
    client.download_file_bytes.assert_called_once()
    assert len(repo.upserts) == 1
    assert repo.lookup_calls[-1][2:] == (
        current_md5,
        len(data),
        drive_file.mime_type,
        OFFICE_VALIDATOR_SCHEMA_VERSION,
    )


@pytest.mark.parametrize("change", ["mime", "validator"])
def test_mime_or_validator_schema_change_forces_revalidation(change: str) -> None:
    data = _docx_bytes()
    drive_file = replace(
        _file("MIME-RECOVERED", data),
        name="confidential.docx",
        mime_type=DOCX_MIME,
    )
    known_mime = PPTX_MIME if change == "mime" else DOCX_MIME
    known_validator = OFFICE_VALIDATOR_SCHEMA_VERSION if change == "mime" else "ooxml-safe-v1"
    repo = _Repository(
        known={
            (
                drive_file.id,
                str(drive_file.md5_checksum),
                len(data),
                known_mime,
                known_validator,
            )
        }
    )
    client = MagicMock()
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.return_value = data

    result = _process_one_gdrive_file(
        drive_file,
        _folder_spec(),
        client=client,
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="req",
        skipped=[],
    )

    assert result[0] == 1
    client.download_file_bytes.assert_called_once()
    assert repo.lookup_calls[-1][-2:] == (
        DOCX_MIME,
        OFFICE_VALIDATOR_SCHEMA_VERSION,
    )


def test_checksum_mismatch_is_warning_but_not_cached_as_known_invalid() -> None:
    data = _pptx_bytes()
    drive_file = _file("TRANSIENT", data, md5_checksum="0" * 32)
    repo = _Repository()
    client = MagicMock()
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.return_value = data

    result = _process_one_gdrive_file(
        drive_file,
        _folder_spec(),
        client=client,
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="req",
        skipped=[],
    )
    assert result == (0, 0)
    assert repo.invalid_records == []
    assert repo.upserts == []


def test_office_download_failure_is_nonpersistent_connector_warning() -> None:
    data = _pptx_bytes()
    drive_file = _file("DOWNLOAD-FAIL", data)
    repo = _Repository()
    client = MagicMock()
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.side_effect = TimeoutError("temporary")
    collector = _IngestWarningCollector()

    result = _process_one_gdrive_file(
        drive_file,
        _folder_spec(),
        client=client,
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="req",
        skipped=[],
        warning_collector=collector,
        warning_source_id="FOLDER",
    )
    assert result == (0, 0)
    assert repo.invalid_records == []
    assert collector.snapshot("gdrive", "FOLDER").reasons == {"office_download_failed": 1}


def test_incremental_transient_failure_is_retried_next_run_after_cursor_advance(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from teamagent.adapters.gdrive_client import ChangeBatch

    monkeypatch.setenv("USE_INCREMENTAL_SYNC", "1")
    data = _pptx_bytes()
    drive_file = _file("RETRY-AFTER-CURSOR", data)
    client = MagicMock()
    client.list_files.return_value = ([drive_file], None)
    client.list_permissions.return_value = [_owner_permission()]
    client.get_start_page_token.return_value = "CURSOR-1"
    client.get_changes.return_value = ChangeBatch(
        changes=(),
        next_page_token=None,
        new_start_page_token="CURSOR-2",
    )
    client.download_file_bytes.side_effect = [TimeoutError("temporary"), data]
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )
    repo = _Repository()

    first = _ingest_gdrive_folder(
        _folder_spec(),
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="run-1",
        warning_collector=_IngestWarningCollector(),
    )
    second = _ingest_gdrive_folder(
        _folder_spec(),
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="run-2",
        warning_collector=_IngestWarningCollector(),
    )

    assert first == (0, 0)
    assert second[0] == 1
    assert [state["cursor"] for state in repo.saved_states] == ["CURSOR-1", "CURSOR-2"]
    assert repo.retry_records[0]["reason"] == "office_download_failed"
    assert repo.retry_resolutions[-1]["external_id"] == drive_file.id
    assert client.download_file_bytes.call_count == 2
    assert len(repo.upserts) == 1
    # quick retryはclaim直後とupsert直前だけ。progress callbackは120秒throttleでDBを叩かない。
    assert len(repo.retry_lease_renewals) == 2


def test_validator_generation_change_forces_incremental_full_revalidation(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("USE_INCREMENTAL_SYNC", "1")
    data = _pptx_bytes()
    drive_file = _file("REVALIDATE", data)
    client = MagicMock()
    client.list_files.return_value = ([drive_file], None)
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.return_value = data
    client.get_start_page_token.return_value = "CURSOR-AFTER-REVALIDATION"
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )
    repo = _Repository()
    repo.connector_state = ConnectorState(
        source_kind="gdrive",
        source_id="FOLDER",
        cursor="OLD-CURSOR",
        metadata={"office_validator_schema_version": "ooxml-safe-v2"},
    )

    result = _ingest_gdrive_folder(
        _folder_spec(),
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="revalidate",
        warning_collector=_IngestWarningCollector(),
    )

    assert result[0] == 1
    client.get_changes.assert_not_called()
    client.get_start_page_token.assert_called_once_with("revalidate")
    client.download_file_bytes.assert_called_once()
    assert repo.saved_states[-1]["cursor"] == "CURSOR-AFTER-REVALIDATION"
    assert (
        repo.saved_states[-1]["metadata"]["office_validator_schema_version"]
        == OFFICE_VALIDATOR_SCHEMA_VERSION
    )


def test_retry_persistence_false_blocks_incremental_cursor(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("USE_INCREMENTAL_SYNC", "1")
    data = _pptx_bytes()
    drive_file = _file("RETRY-NOT-DURABLE", data)
    client = MagicMock()
    client.list_files.return_value = ([drive_file], None)
    client.list_permissions.return_value = [_owner_permission()]
    client.get_start_page_token.return_value = "MUST-NOT-SAVE"
    client.download_file_bytes.side_effect = TimeoutError("temporary")
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )
    repo = _Repository()
    repo.retry_record_success = False

    with pytest.raises(IngestDurabilityError, match="cursor not advanced"):
        _ingest_gdrive_folder(
            _folder_spec(),
            embedder=_Embedder(),
            repository=repo,  # type: ignore[arg-type]
            owner_email="bot@example.jp",
            dry_run=False,
            request_id="retry-false",
            warning_collector=_IngestWarningCollector(),
        )

    assert repo.retry_records[-1]["reason"] == "office_download_failed"
    assert repo.saved_states == []
    assert repo.upserts == []


def test_retry_claim_failure_blocks_incremental_cursor(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    class _ClaimFailureRepository(_Repository):
        def claim_due_source_retries(self, **kwargs: Any) -> list[SourceRetry]:
            raise RuntimeError("retry table unavailable")

    monkeypatch.setenv("USE_INCREMENTAL_SYNC", "1")
    client = MagicMock()
    client.list_files.return_value = ([], None)
    client.get_start_page_token.return_value = "MUST-NOT-SAVE"
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )
    repo = _ClaimFailureRepository()

    with pytest.raises(IngestDurabilityError, match="cursor not advanced"):
        _ingest_gdrive_folder(
            _folder_spec(),
            embedder=_Embedder(),
            repository=repo,  # type: ignore[arg-type]
            owner_email="bot@example.jp",
            dry_run=False,
            request_id="claim-failed",
            warning_collector=_IngestWarningCollector(),
        )

    assert repo.saved_states == []


def test_native_gdoc_failure_requires_durable_retry_before_cursor(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("USE_INCREMENTAL_SYNC", "1")
    data = b"native-placeholder"
    drive_file = replace(
        _file("GDOC-FAIL", data),
        name="confidential-native",
        mime_type=GDOC_NATIVE_MIME,
    )
    client = MagicMock()
    client.list_files.return_value = ([drive_file], None)
    client.list_permissions.return_value = [_owner_permission()]
    client.get_start_page_token.return_value = "MUST-NOT-SAVE"
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )
    gdocs = MagicMock()
    gdocs.get_document_text.side_effect = TimeoutError("temporary")
    monkeypatch.setattr(
        "teamagent.adapters.gdocs_client.GDocsClient.from_env",
        classmethod(lambda cls, **kwargs: gdocs),
    )
    repo = _Repository()
    repo.retry_record_success = False

    with pytest.raises(IngestDurabilityError, match="cursor not advanced"):
        _ingest_gdrive_folder(
            _folder_spec(),
            embedder=_Embedder(),
            repository=repo,  # type: ignore[arg-type]
            owner_email="bot@example.jp",
            dry_run=False,
            request_id="gdoc-false",
            warning_collector=_IngestWarningCollector(),
        )

    assert repo.retry_records[-1]["reason"] == "gdoc_extract_failed"
    assert repo.saved_states == []
    assert repo.upserts == []


@pytest.mark.parametrize(
    "limit_name",
    [
        "MAX_INGEST_EXTRACTED_CHARACTERS",
        "MAX_INGEST_CHUNKS_PER_FILE",
        "MAX_INGEST_EMBEDDINGS_PER_FILE",
    ],
)
def test_office_chunk_and_embedding_limits_preserve_existing_document(
    monkeypatch: pytest.MonkeyPatch,
    limit_name: str,
) -> None:
    import teamagent.ingest.office_extract as office_extract
    import teamagent.ingest.pipeline as pipeline

    data = _pptx_bytes()
    drive_file = _file(f"LIMIT-{limit_name}", data)
    client = MagicMock()
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.return_value = data
    monkeypatch.setattr(
        office_extract,
        "extract_office_pages",
        lambda *args, **kwargs: [(1, "x" * 1200)],
    )
    monkeypatch.setattr(pipeline, limit_name, 1)
    repo = _Repository()
    collector = _IngestWarningCollector()

    result = _process_one_gdrive_file(
        drive_file,
        _folder_spec(),
        client=client,
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="bounded",
        skipped=[],
        warning_collector=collector,
        warning_source_id="FOLDER",
    )

    assert result == (0, 0)
    assert repo.upserts == []
    assert repo.invalid_records[-1]["reason"] == "unsafe_content_volume"
    assert collector.snapshot("gdrive", "FOLDER").reasons == {"unsafe_content_volume": 1}


def test_contextualizer_cannot_expand_past_per_file_limit(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import teamagent.ingest.pipeline as pipeline

    class _ExpandingContextualizer:
        def contextualize_chunks(
            self,
            title: str,
            full_text: str,
            chunks: list[ChunkUpsert],
            request_id: str,
        ) -> list[ChunkUpsert]:
            return [*chunks, replace(chunks[0], chunk_idx=1)]

    data = _pptx_bytes()
    drive_file = _file("CONTEXT-LIMIT", data)
    client = MagicMock()
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.return_value = data
    monkeypatch.setattr(pipeline, "MAX_INGEST_EMBEDDINGS_PER_FILE", 1)
    repo = _Repository()
    collector = _IngestWarningCollector()

    result = _process_one_gdrive_file(
        drive_file,
        _folder_spec(),
        client=client,
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="context-limit",
        skipped=[],
        contextualizer=_ExpandingContextualizer(),  # type: ignore[arg-type]
        warning_collector=collector,
        warning_source_id="FOLDER",
        durable_retry=True,
    )

    assert result == (0, 0)
    assert repo.upserts == []
    assert repo.retry_records[-1]["reason"] == "contextualized_content_too_large"
    assert collector.snapshot("gdrive", "FOLDER").reasons == {"contextualized_content_too_large": 1}


def test_retry_lease_loss_stops_upsert_and_cursor(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from teamagent.adapters.gdrive_client import ChangeBatch

    monkeypatch.setenv("USE_INCREMENTAL_SYNC", "1")
    data = _pptx_bytes()
    drive_file = _file("LEASE-LOST", data)
    client = MagicMock()
    client.list_files.return_value = ([drive_file], None)
    client.get_changes.return_value = ChangeBatch(
        changes=(),
        next_page_token=None,
        new_start_page_token="MUST-NOT-SAVE",
    )
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )
    repo = _Repository()
    repo.connector_state = ConnectorState(
        source_kind="gdrive",
        source_id="FOLDER",
        cursor="OLD-CURSOR",
        metadata={"office_validator_schema_version": OFFICE_VALIDATOR_SCHEMA_VERSION},
    )
    repo.retry_claims = [
        SourceRetry(
            external_id=drive_file.id,
            md5_checksum=drive_file.md5_checksum,
            size_bytes=drive_file.size,
            mime_type=drive_file.mime_type,
            validator_schema_version=OFFICE_VALIDATOR_SCHEMA_VERSION,
            attempt_count=2,
            reason="office_download_failed",
        )
    ]
    repo.retry_lease_success = False

    with pytest.raises(IngestDurabilityError, match="cursor not advanced"):
        _ingest_gdrive_folder(
            _folder_spec(),
            embedder=_Embedder(),
            repository=repo,  # type: ignore[arg-type]
            owner_email="bot@example.jp",
            dry_run=False,
            request_id="lease-lost",
            warning_collector=_IngestWarningCollector(),
        )

    assert repo.retry_lease_renewals
    assert repo.upserts == []
    assert repo.saved_states == []


@pytest.mark.parametrize(
    "reason",
    [
        "page limit reached with remaining token",
        "pagination token cycle",
        "max depth reached with child folders remaining",
        "shared drive page limit reached with remaining token",
    ],
)
def test_incomplete_traversal_commits_neither_cursor_nor_stale_cleanup(
    monkeypatch: pytest.MonkeyPatch,
    reason: str,
) -> None:
    monkeypatch.setenv("USE_INCREMENTAL_SYNC", "1")
    monkeypatch.setenv("INGEST_MARK_STALE", "true")
    monkeypatch.delenv("BOILERPLATE_DETECT", raising=False)
    monkeypatch.delenv("DOC_DEDUP_DETECT", raising=False)

    client = MagicMock()
    client.get_start_page_token.return_value = "MUST-NOT-SAVE"
    client.walk_files_recursive.side_effect = GDriveTraversalIncompleteError(
        "walk_files_recursive",
        reason,
        root_id="FOLDER",
    )
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )

    class _TraversalRepository(_Repository):
        def __init__(self) -> None:
            super().__init__()
            self.stale_cleanup_calls: list[str] = []

        def clear_documents_stale(self, external_ids: list[str]) -> int:
            self.stale_cleanup_calls.append("clear")
            return 0

        def list_gdrive_external_ids_with_stale(self) -> list[tuple[str, bool]]:
            self.stale_cleanup_calls.append("list")
            return [("existing", False)]

        def mark_documents_stale(self, external_ids: list[str]) -> int:
            self.stale_cleanup_calls.append("mark")
            return 0

    repo = _TraversalRepository()
    sources = IngestSources(
        version=1,
        slack_channels=(),
        gdrive_folders=(replace(_folder_spec(), include_subfolders=True),),
        gsheets=(),
    )
    result = IngestRunner(
        repo,  # type: ignore[arg-type]
        _Embedder(),
        owner_email="bot@example.jp",
        dry_run=False,
        alerter=IngestOpsAlerter(webhook_url=None),
    ).run(sources, kinds=["gdrive"])

    assert result.by_kind["gdrive"].errors
    assert not any(state.get("success") is True for state in repo.saved_states)
    assert not any(state.get("cursor") for state in repo.saved_states)
    assert repo.stale_cleanup_calls == []


def test_pdf_download_failure_is_warning_and_durable_retry() -> None:
    drive_file = replace(
        _file("PDF-FAIL", b"%PDF-fake"),
        name="confidential.pdf",
        mime_type="application/pdf",
    )
    repo = _Repository()
    client = MagicMock()
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.side_effect = TimeoutError("temporary")
    collector = _IngestWarningCollector()

    result = _process_one_gdrive_file(
        drive_file,
        _folder_spec(),
        client=client,
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="req",
        skipped=[],
        warning_collector=collector,
        warning_source_id="FOLDER",
        durable_retry=True,
        retry_lease_owner="claim-token",
    )

    assert result == (0, 0)
    assert collector.snapshot("gdrive", "FOLDER").reasons == {"pdf_download_failed": 1}
    assert repo.retry_records[0]["reason"] == "pdf_download_failed"
    assert repo.retry_records[0]["expected_lease_owner"] == "claim-token"
    assert repo.upserts == []


def test_claimed_file_failure_threads_lease_owner_to_retry_write(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("USE_INCREMENTAL_SYNC", "1")
    drive_file = replace(
        _file("CLAIMED-PDF-FAIL", b"%PDF-fake"),
        name="confidential.pdf",
        mime_type="application/pdf",
    )
    client = MagicMock()
    client.get_start_page_token.return_value = "NEXT"
    client.list_files.return_value = ([drive_file], None)
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.side_effect = TimeoutError("temporary")
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )
    repo = _Repository()
    repo.retry_claims = [
        SourceRetry(
            external_id=drive_file.id,
            md5_checksum=drive_file.md5_checksum,
            size_bytes=drive_file.size,
            mime_type=drive_file.mime_type,
            validator_schema_version="pdf-safe-v1",
            attempt_count=1,
            reason="pdf_download_failed",
            lease_owner="worker-a",
        )
    ]

    _ingest_gdrive_folder(
        _folder_spec(),
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="worker-a",
        warning_collector=_IngestWarningCollector(),
    )

    assert repo.retry_records[-1]["expected_lease_owner"] == "worker-a"


@pytest.mark.parametrize(
    ("failure_mode", "expected_reason"),
    [
        ("extract", "pdf_extract_failed"),
        ("empty", "pdf_empty_text"),
    ],
)
def test_pdf_extract_and_empty_text_are_warnings_and_durable_retries(
    monkeypatch: pytest.MonkeyPatch,
    failure_mode: str,
    expected_reason: str,
) -> None:
    import teamagent.ingest.pdf_extract as pdf_extract

    drive_file = replace(
        _file("PDF-CONTENT-FAIL", b"%PDF-fake"),
        name="confidential.pdf",
        mime_type="application/pdf",
    )
    repo = _Repository()
    client = MagicMock()
    client.list_permissions.return_value = [_owner_permission()]
    client.download_file_bytes.return_value = b"%PDF-fake"
    if failure_mode == "extract":

        def _raise_extract_error(data: bytes, **kwargs: Any) -> list[tuple[int, str]]:
            raise ValueError("malformed PDF")

        monkeypatch.setattr(pdf_extract, "extract_pdf_pages", _raise_extract_error)
    else:
        monkeypatch.setattr(pdf_extract, "extract_pdf_pages", lambda data, **kwargs: [])
    collector = _IngestWarningCollector()

    result = _process_one_gdrive_file(
        drive_file,
        _folder_spec(),
        client=client,
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="req",
        skipped=[],
        warning_collector=collector,
        warning_source_id="FOLDER",
        durable_retry=True,
    )

    assert result == (0, 0)
    assert collector.snapshot("gdrive", "FOLDER").reasons == {expected_reason: 1}
    assert repo.retry_records[0]["reason"] == expected_reason
    assert repo.upserts == []


@pytest.mark.parametrize(
    ("failure_mode", "expected_reason"),
    [
        ("download", "pdf_download_failed"),
        ("extract", "pdf_extract_failed"),
        ("empty", "pdf_empty_text"),
    ],
)
def test_shared_drive_pdf_failures_are_connector_warnings(
    monkeypatch: pytest.MonkeyPatch,
    failure_mode: str,
    expected_reason: str,
) -> None:
    import teamagent.ingest.pdf_extract as pdf_extract

    drive_file = replace(
        _file("SHARED-PDF-FAIL", b"%PDF-fake"),
        name="confidential.pdf",
        mime_type="application/pdf",
    )
    client = MagicMock()
    client.list_shared_drives.return_value = [SharedDrive(id="DRIVE", name="drive")]
    client.walk_files_recursive.return_value = [drive_file]
    client.list_permissions.return_value = [_owner_permission()]
    if failure_mode == "download":
        client.download_file_bytes.side_effect = TimeoutError("temporary")
    else:
        client.download_file_bytes.return_value = b"%PDF-fake"
        if failure_mode == "extract":

            def _raise_extract_error(data: bytes, **kwargs: Any) -> list[tuple[int, str]]:
                raise ValueError("malformed PDF")

            monkeypatch.setattr(pdf_extract, "extract_pdf_pages", _raise_extract_error)
        else:
            monkeypatch.setattr(pdf_extract, "extract_pdf_pages", lambda data, **kwargs: [])
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )
    collector = _IngestWarningCollector()

    result = _ingest_shared_drives_crawl(
        SharedDriveCrawlSpec(enabled=True, sales_relevance_filter=False),
        embedder=_Embedder(),
        repository=_Repository(),  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="req",
        warning_collector=collector,
    )

    assert result == (0, 0)
    assert collector.snapshot("shared_drives", "shared_drives").reasons == {expected_reason: 1}


def test_shared_drive_known_invalid_is_observed_and_suppressed_before_acl(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    broken = b"PK\x03\x04" + (b"x" * 64)
    drive_file = _file("SHARED-KNOWN", broken)
    repo = _Repository(
        known={
            (
                "SHARED-KNOWN",
                str(drive_file.md5_checksum),
                len(broken),
                drive_file.mime_type,
                OFFICE_VALIDATOR_SCHEMA_VERSION,
            )
        }
    )
    client = MagicMock()
    client.list_shared_drives.return_value = [SharedDrive(id="DRIVE", name="drive")]
    client.walk_files_recursive.return_value = [drive_file]
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )
    observed: set[str] = set()

    result = _ingest_shared_drives_crawl(
        SharedDriveCrawlSpec(enabled=True, sales_relevance_filter=False),
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="req",
        observed_gdrive_ids=observed,
        warning_collector=_IngestWarningCollector(),
    )

    assert result == (0, 0)
    assert observed == {"SHARED-KNOWN"}
    client.list_permissions.assert_not_called()
    client.download_file_bytes.assert_not_called()
    assert repo.upserts == []


def test_shared_drive_walk_saturation_fails_closed(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    drive_file = _file("SHARED-SATURATED", _pptx_bytes())
    client = MagicMock()
    client.list_shared_drives.return_value = [SharedDrive(id="DRIVE", name="drive")]
    client.walk_files_recursive.return_value = [drive_file]
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )
    truncated: set[str] = set()

    from teamagent.ingest.pipeline import GDrivePaginationIncompleteError

    with pytest.raises(GDrivePaginationIncompleteError, match="safety limit"):
        _ingest_shared_drives_crawl(
            SharedDriveCrawlSpec(
                enabled=True,
                sales_relevance_filter=False,
                max_files_per_drive=1,
            ),
            embedder=_Embedder(),
            repository=_Repository(),  # type: ignore[arg-type]
            owner_email="bot@example.jp",
            dry_run=False,
            request_id="req",
            truncated_walk_roots=truncated,
            warning_collector=_IngestWarningCollector(),
        )

    assert truncated == {"DRIVE"}
    client.list_permissions.assert_not_called()


def test_incremental_cursor_records_success_with_warnings_metadata(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("USE_INCREMENTAL_SYNC", "1")
    broken = b"PK\x03\x04" + (b"x" * 64)
    drive_file = _file("INCREMENTAL-KNOWN", broken)
    repo = _Repository(
        known={
            (
                "INCREMENTAL-KNOWN",
                str(drive_file.md5_checksum),
                len(broken),
                drive_file.mime_type,
                OFFICE_VALIDATOR_SCHEMA_VERSION,
            )
        }
    )
    client = MagicMock()
    client.list_files.return_value = ([drive_file], None)
    client.get_start_page_token.return_value = "NEXT"
    monkeypatch.setattr(
        "teamagent.adapters.gdrive_client.GDriveClient.from_env",
        classmethod(lambda cls, **kwargs: client),
    )

    result = _ingest_gdrive_folder(
        _folder_spec(),
        embedder=_Embedder(),
        repository=repo,  # type: ignore[arg-type]
        owner_email="bot@example.jp",
        dry_run=False,
        request_id="req",
        warning_collector=_IngestWarningCollector(),
    )

    assert result == (0, 0)
    saved = repo.saved_states[-1]
    assert saved["success"] is True
    assert saved["metadata"]["outcome"] == "success_with_warnings"
    assert saved["metadata"]["warning_reasons"] == {"corrupt_zip": 1}
    assert saved["metadata"]["known_invalid_suppressed"] == 1


def test_runner_records_connector_warning_outcome_and_notifies_ops(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    def warning_handler(spec: Any, **kwargs: Any) -> tuple[int, int]:
        collector = kwargs["warning_collector"]
        collector.add("gdrive", spec.folder_id, "corrupt_zip", suppressed=True)
        return 0, 0

    monkeypatch.setattr("teamagent.ingest.pipeline._ingest_gdrive_folder", warning_handler)
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.IngestRunner._maybe_check_freshness",
        lambda self, *, request_id: None,
    )
    repo = _Repository()
    alerter = IngestOpsAlerter(webhook_url="https://hooks.slack.test/ingest")
    runner = IngestRunner(
        repository=repo,  # type: ignore[arg-type]
        embedder=_Embedder(),
        owner_email="bot@example.jp",
        dry_run=False,
        alerter=alerter,
    )
    sources = IngestSources(
        version=1,
        slack_channels=(),
        gdrive_folders=(_folder_spec(),),
        gsheets=(),
    )

    with patch("httpx.post") as mock_post:
        mock_post.return_value.status_code = 200
        result = runner.run(sources, kinds=["gdrive"])

    stats = result.by_kind["gdrive"]
    assert stats.outcome == "success_with_warnings"
    assert stats.warning_reasons == {"corrupt_zip": 1}
    assert result.outcome == "success_with_warnings"
    assert repo.connector_runs[0]["outcome"] == "success_with_warnings"
    assert repo.connector_runs[0]["warning_reasons"] == {"corrupt_zip": 1}
    mock_post.assert_called_once()
    rendered_notification = str(mock_post.call_args.kwargs["json"])
    assert "success_with_warnings" in rendered_notification
    assert "corrupt_zip" in rendered_notification
    assert "file_id" not in rendered_notification
    assert "file_name" not in rendered_notification


def test_reconciliation_keeps_three_pdfs_and_nine_missing_originals_as_warnings(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.IngestRunner._maybe_check_freshness",
        lambda self, *, request_id: None,
    )
    repo = _Repository()
    repo.reconciliation_counts = {
        "unindexed_pdf": 3,
        "source_original_missing": 9,
    }
    alerter = IngestOpsAlerter(webhook_url="https://hooks.slack.test/ingest")
    runner = IngestRunner(
        repository=repo,  # type: ignore[arg-type]
        embedder=_Embedder(),
        owner_email="bot@example.jp",
        dry_run=False,
        alerter=alerter,
    )
    sources = IngestSources(
        version=1,
        slack_channels=(),
        gdrive_folders=(),
        gsheets=(),
    )

    with patch("httpx.post") as mock_post:
        mock_post.return_value.status_code = 200
        result = runner.run(sources, kinds=["gdrive"])

    assert result.outcome == "success_with_warnings"
    assert result.by_kind["gdrive"].warning_reasons == {
        "unindexed_pdf": 3,
        "source_original_missing": 9,
    }
    assert repo.connector_runs == [
        {
            "request_id": repo.connector_runs[0]["request_id"],
            "source_kind": "gdrive",
            "source_id": "__reconciliation__",
            "outcome": "success_with_warnings",
            "documents_upserted": 0,
            "chunks_inserted": 0,
            "warning_reasons": {
                "unindexed_pdf": 3,
                "source_original_missing": 9,
            },
            "suppressed_retry_count": 0,
            "error": None,
        }
    ]
    notification = str(mock_post.call_args.kwargs["json"])
    assert "unindexed_pdf" in notification
    assert "source_original_missing" in notification
    assert "file_id" not in notification
    assert "title" not in notification.lower()
    assert "customer" not in notification.lower()


def test_normal_no_change_run_remains_success(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    def no_change_handler(spec: Any, **kwargs: Any) -> tuple[int, int]:
        return 0, 0

    monkeypatch.setattr("teamagent.ingest.pipeline._ingest_gdrive_folder", no_change_handler)
    monkeypatch.setattr(
        "teamagent.ingest.pipeline.IngestRunner._maybe_check_freshness",
        lambda self, *, request_id: None,
    )
    repo = _Repository()
    runner = IngestRunner(
        repository=repo,  # type: ignore[arg-type]
        embedder=_Embedder(),
        owner_email="bot@example.jp",
        dry_run=False,
        alerter=IngestOpsAlerter(webhook_url=None),
    )
    sources = IngestSources(
        version=1,
        slack_channels=(),
        gdrive_folders=(_folder_spec(),),
        gsheets=(),
    )

    result = runner.run(sources, kinds=["gdrive"])

    assert result.outcome == "success"
    assert result.by_kind["gdrive"].warning_reasons == {}
    assert repo.connector_runs[0]["outcome"] == "success"
    assert repo.connector_runs[0]["warning_reasons"] == {}


def test_warning_notification_failure_is_fail_open_and_dry_run_is_noop() -> None:
    alerter = IngestOpsAlerter(webhook_url="https://hooks.slack.test/ingest")
    with patch("httpx.post", side_effect=RuntimeError("network unavailable")):
        assert (
            _send_ops_warning_summary(
                alerter,
                kind="gdrive",
                warning_reasons={"corrupt_zip": 1},
                suppressed_retry_count=0,
                request_id="req",
                dry_run=False,
            )
            is False
        )
    with patch("httpx.post") as mock_post:
        assert (
            _send_ops_warning_summary(
                alerter,
                kind="gdrive",
                warning_reasons={"corrupt_zip": 1},
                suppressed_retry_count=0,
                request_id="req",
                dry_run=True,
            )
            is False
        )
    mock_post.assert_not_called()


def test_pipeline_title_only_guard_consults_database_when_process_registry_is_empty() -> None:
    class _GuardRepository:
        def __init__(self) -> None:
            self.title_guard_calls = 0

        def upsert_title_only_if_no_content(
            self,
            doc: Any,
            chunks: list[Any],
            request_id: str,
        ) -> None:
            self.title_guard_calls += 1
            return None

        def upsert_document_with_chunks(self, *args: Any, **kwargs: Any) -> str:
            raise AssertionError("unguarded replacement must not run")

    repo = _GuardRepository()
    doc = DocumentUpsert(
        source_type="gdrive",
        external_id="EXISTING",
        owner_email="bot@example.jp",
    )
    chunks = [
        ChunkUpsert(
            chunk_idx=0,
            content="title",
            embedding=[0.1] * 4,
            metadata={"title_only": True},
        )
    ]
    assert (
        _guarded_upsert(
            repo,  # type: ignore[arg-type]
            doc,
            chunks,
            request_id="req",
            content_registry=set(),
        )
        is False
    )
    assert repo.title_guard_calls == 1
