"""office_extract.py のテスト — docx/pptx/xlsx をテスト内で生成して抽出ロジックを検証。

バイナリ fixture を repo に置かず、各 lib (python-docx / python-pptx / openpyxl) で
生成 → BytesIO に save → 抽出という流れ。lib 自体の動作と統合できる。
"""

from __future__ import annotations

import hashlib
import zipfile
from io import BytesIO
from xml.etree import ElementTree

import pytest

import teamagent.ingest.office_extract as office_extract
from teamagent.ingest.office_extract import (
    DOCX_MIME,
    GDOC_NATIVE_MIME,
    OFFICE_BINARY_MIMES,
    PPTX_MIME,
    XLSX_MIME,
    OfficePayloadError,
    _normalize_text,
    extract_docx_text,
    extract_office_pages,
    extract_pptx_pages,
    extract_xlsx_pages,
)

_PPTX_CONTENT_TYPES = b"""\
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Override PartName="/ppt/presentation.xml"
    ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
</Types>
"""
_PPTX_ROOT = (
    b'<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
)


# -----------------------------------------------------------
# _normalize_text
# -----------------------------------------------------------
def test_normalize_text_removes_nul_and_collapses_whitespace() -> None:
    assert _normalize_text("hello\x00world  \n\nfoo") == "helloworld foo"


def test_normalize_text_empty() -> None:
    assert _normalize_text("") == ""
    assert _normalize_text("   \n\t  ") == ""


# -----------------------------------------------------------
# docx
# -----------------------------------------------------------
def _build_sample_docx() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph("見出しテキスト")
    doc.add_paragraph("本文のサンプルです。")
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "セルA"
    table.rows[0].cells[1].text = "セルB"
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def test_extract_docx_text_includes_paragraphs_and_tables() -> None:
    text = extract_docx_text(_build_sample_docx())
    assert "見出しテキスト" in text
    assert "本文のサンプル" in text
    assert "セルA" in text
    assert "セルB" in text


def test_extract_docx_text_empty() -> None:
    from docx import Document

    buf = BytesIO()
    Document().save(buf)
    # 空 docx は空文字（空白除去後）
    assert extract_docx_text(buf.getvalue()) == ""


# -----------------------------------------------------------
# pptx
# -----------------------------------------------------------
def _build_sample_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # Title Only
    s1 = prs.slides.add_slide(blank_layout)
    s1.shapes.title.text = "スライド1のタイトル"
    # text box を足す
    tb = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    tb.text_frame.text = "ボディテキスト1"
    s2 = prs.slides.add_slide(blank_layout)
    s2.shapes.title.text = "スライド2のタイトル"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_pptx_pages_returns_per_slide_text() -> None:
    pages = extract_pptx_pages(_build_sample_pptx())
    assert len(pages) == 2
    nums = [n for n, _ in pages]
    assert nums == [1, 2]
    s1_text = pages[0][1]
    assert "スライド1のタイトル" in s1_text
    assert "ボディテキスト1" in s1_text
    assert "スライド2のタイトル" in pages[1][1]


def test_extract_pptx_pages_skips_empty_slides() -> None:
    from pptx import Presentation

    prs = Presentation()
    # 空 slide を 1 枚だけ追加（空 layout）
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    buf = BytesIO()
    prs.save(buf)
    assert extract_pptx_pages(buf.getvalue()) == []


# -----------------------------------------------------------
# xlsx
# -----------------------------------------------------------
def _build_sample_xlsx() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1["A1"] = "顧客名"
    ws1["B1"] = "金額"
    ws1["A2"] = "株式会社X"
    ws1["B2"] = 1000
    ws2 = wb.create_sheet("Sheet2")
    ws2["A1"] = "備考"
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_extract_xlsx_pages_returns_per_sheet_text() -> None:
    pages = extract_xlsx_pages(_build_sample_xlsx())
    assert len(pages) == 2
    nums = [n for n, _ in pages]
    assert nums == [1, 2]
    assert "顧客名" in pages[0][1]
    assert "株式会社X" in pages[0][1]
    assert "1000" in pages[0][1]
    assert "備考" in pages[1][1]


def test_extract_xlsx_pages_skips_empty_sheet() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    wb.active.title = "Empty"
    buf = BytesIO()
    wb.save(buf)
    assert extract_xlsx_pages(buf.getvalue()) == []


# -----------------------------------------------------------
# extract_office_pages dispatcher
# -----------------------------------------------------------
def test_extract_office_pages_dispatches_docx() -> None:
    pages = extract_office_pages(_build_sample_docx(), mime_type=DOCX_MIME)
    assert len(pages) == 1
    assert pages[0][0] == 1
    assert "見出し" in pages[0][1]


def test_extract_office_pages_dispatches_pptx() -> None:
    pages = extract_office_pages(_build_sample_pptx(), mime_type=PPTX_MIME)
    assert len(pages) == 2


def test_extract_office_pages_dispatches_xlsx() -> None:
    pages = extract_office_pages(_build_sample_xlsx(), mime_type=XLSX_MIME)
    assert len(pages) == 2


def test_extract_office_pages_unsupported_mime_raises() -> None:
    with pytest.raises(ValueError, match="Unsupported office mime type"):
        extract_office_pages(b"", mime_type="application/octet-stream")


def test_office_payload_classifies_html_response_before_size_mismatch() -> None:
    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(
            b"\xef\xbb\xbf  <!DOCTYPE html><html><body>Drive error</body></html>",
            mime_type=PPTX_MIME,
            expected_size=999,
        )
    assert raised.value.category == "html_response"
    assert raised.value.expected_bytes == 999


def test_office_payload_classifies_truncated_download_from_drive_size() -> None:
    data = _build_sample_pptx()
    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data) + 1)
    assert raised.value.category == "truncated_download"
    assert raised.value.actual_bytes == len(data)


def test_office_payload_classifies_zip_without_central_directory_as_corrupt() -> None:
    data = b"PK\x03\x04" + (b"x" * 64)
    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))
    assert raised.value.category == "corrupt_zip"


def test_office_payload_classifies_ooxml_mime_mismatch() -> None:
    data = _build_sample_docx()
    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))
    assert raised.value.category == "format_mismatch"


def test_office_payload_classifies_drive_md5_mismatch() -> None:
    data = _build_sample_pptx()
    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(
            data,
            mime_type=PPTX_MIME,
            expected_size=len(data),
            expected_md5="0" * 32,
        )
    assert raised.value.category == "checksum_mismatch"


def test_office_payload_runs_zip_crc_integrity_before_extract() -> None:
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_STORED) as package:
        package.writestr(
            "ppt/presentation.xml",
            _PPTX_ROOT,
        )
        package.writestr("payload.bin", b"CRC-CONTENT")
    data = bytearray(buf.getvalue())
    payload_at = data.index(b"CRC-CONTENT")
    data[payload_at] ^= 0x01

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(bytes(data), mime_type=PPTX_MIME, expected_size=len(data))
    assert raised.value.category == "corrupt_zip"


def test_office_payload_parses_required_ooxml_part_before_extract() -> None:
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as package:
        package.writestr("ppt/presentation.xml", b"<p:presentation")
        package.writestr("[Content_Types].xml", _PPTX_CONTENT_TYPES)
    data = buf.getvalue()

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))
    assert raised.value.category == "corrupt_zip"


@pytest.mark.parametrize(
    "declaration",
    [
        '<!DOCTYPE p:presentation [<!ENTITY payload "expanded">]>',
        (
            '<!DOCTYPE p:presentation ['
            '<!ENTITY payload SYSTEM "file:///definitely-not-readable">'
            "]>"
        ),
    ],
)
def test_office_payload_rejects_dtd_and_entities_even_in_utf16(
    declaration: str,
) -> None:
    required_xml = (
        '<?xml version="1.0" encoding="UTF-16"?>'
        f"{declaration}"
        '<p:presentation xmlns:p="http://schemas.openxmlformats.org/'
        'presentationml/2006/main">&payload;</p:presentation>'
    ).encode("utf-16")
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_STORED) as package:
        package.writestr("ppt/presentation.xml", required_xml)
        package.writestr("[Content_Types].xml", _PPTX_CONTENT_TYPES)
    data = buf.getvalue()

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "unsafe_archive"


def test_office_payload_rejects_dtd_in_non_required_slide_xml() -> None:
    original = _build_sample_pptx()
    source = zipfile.ZipFile(BytesIO(original))
    buf = BytesIO()
    malicious_slide = (
        b'<?xml version="1.0" encoding="UTF-8"?>'
        b'<!DOCTYPE p:sld [<!ENTITY payload "expanded">]>'
        b'<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
        b"&payload;</p:sld>"
    )
    with source, zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as package:
        for info in source.infolist():
            payload = source.read(info)
            if info.filename == "ppt/slides/slide1.xml":
                payload = malicious_slide
            package.writestr(info, payload)
    data = buf.getvalue()

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "unsafe_archive"


def test_office_extracted_text_hard_cap_is_classified() -> None:
    data = _build_sample_docx()

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(
            data,
            mime_type=DOCX_MIME,
            max_extracted_chars=8,
        )

    assert raised.value.category == "unsafe_content_volume"


def test_xlsx_cell_limit_fails_instead_of_returning_partial_text(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(office_extract, "_XLSX_MAX_CELLS_PER_SHEET", 2)

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(
            _build_sample_xlsx(),
            mime_type=XLSX_MIME,
        )

    assert raised.value.category == "unsafe_content_volume"


def test_office_progress_callback_runs_and_its_error_is_not_payload_corruption() -> None:
    class LeaseLostError(RuntimeError):
        pass

    calls = 0

    def heartbeat() -> None:
        nonlocal calls
        calls += 1
        if calls == 2:
            raise LeaseLostError("lease lost")

    with pytest.raises(LeaseLostError, match="lease lost"):
        extract_office_pages(
            _build_sample_pptx(),
            mime_type=PPTX_MIME,
            progress_callback=heartbeat,
        )

    assert calls == 2


def test_office_payload_rejects_16kb_to_16mb_zip_bomb_by_ratio() -> None:
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as package:
        package.writestr("ppt/presentation.xml", _PPTX_ROOT)
        package.writestr("[Content_Types].xml", _PPTX_CONTENT_TYPES)
        package.writestr("ppt/media/amplifier.bin", b"A" * (16 * 1024 * 1024))
    data = buf.getvalue()
    assert len(data) < 32 * 1024

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "unsafe_archive"


def test_office_payload_rejects_compressed_input_over_hard_limit(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    data = _build_sample_pptx()
    monkeypatch.setattr(office_extract, "MAX_OFFICE_COMPRESSED_BYTES", len(data) - 1)

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "unsafe_archive"


def test_office_payload_rejects_required_xml_over_parse_limit(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(office_extract, "MAX_OFFICE_REQUIRED_XML_BYTES", 128)
    oversized_xml = (
        b'<p:presentation xmlns:p="http://schemas.openxmlformats.org/'
        b'presentationml/2006/main">' + (b" " * 128) + b"</p:presentation>"
    )
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_STORED) as package:
        package.writestr("ppt/presentation.xml", oversized_xml)
        package.writestr("[Content_Types].xml", _PPTX_CONTENT_TYPES)
    data = buf.getvalue()

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "unsafe_archive"


def test_office_payload_rejects_oversized_non_required_xml(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(office_extract, "MAX_OFFICE_XML_MEMBER_BYTES", 256)
    oversized_xml = b"<root>" + (b" " * 512) + b"</root>"
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_STORED) as package:
        package.writestr("ppt/presentation.xml", _PPTX_ROOT)
        package.writestr("[Content_Types].xml", _PPTX_CONTENT_TYPES)
        package.writestr("ppt/slides/slide1.xml", oversized_xml)
    data = buf.getvalue()

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "unsafe_archive"


def test_office_payload_rejects_excessive_member_count(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(office_extract, "MAX_OFFICE_ZIP_MEMBERS", 2)
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as package:
        package.writestr("ppt/presentation.xml", _PPTX_ROOT)
        package.writestr("[Content_Types].xml", _PPTX_CONTENT_TYPES)
        package.writestr("docProps/core.xml", b"<core/>")
    data = buf.getvalue()

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "unsafe_archive"


def test_office_payload_validates_required_root_and_content_type() -> None:
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as package:
        package.writestr(
            "ppt/presentation.xml",
            b'<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"/>',
        )
        package.writestr("[Content_Types].xml", _PPTX_CONTENT_TYPES)
    data = buf.getvalue()

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "format_mismatch"


def test_office_payload_rejects_wrong_content_types_override() -> None:
    content_types = ElementTree.fromstring(_PPTX_CONTENT_TYPES)
    override = next(iter(content_types))
    override.set("ContentType", "application/octet-stream")
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as package:
        package.writestr("ppt/presentation.xml", _PPTX_ROOT)
        package.writestr("[Content_Types].xml", ElementTree.tostring(content_types))
    data = buf.getvalue()

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "format_mismatch"


def test_office_payload_classifies_encrypted_compound_office() -> None:
    data = (
        b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
        + "EncryptedPackage".encode("utf-16le")
        + "EncryptionInfo".encode("utf-16le")
    )

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(data, mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "encrypted_office"


def test_office_payload_classifies_encrypted_zip_member() -> None:
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_STORED) as package:
        package.writestr("ppt/presentation.xml", _PPTX_ROOT)
        package.writestr("[Content_Types].xml", _PPTX_CONTENT_TYPES)
    data = bytearray(buf.getvalue())
    local_header = data.index(b"PK\x03\x04")
    central_header = data.index(b"PK\x01\x02")
    data[local_header + 6 : local_header + 8] = (1).to_bytes(2, "little")
    data[central_header + 8 : central_header + 10] = (1).to_bytes(2, "little")

    with pytest.raises(OfficePayloadError) as raised:
        extract_office_pages(bytes(data), mime_type=PPTX_MIME, expected_size=len(data))

    assert raised.value.category == "encrypted_office"


def test_office_payload_does_not_use_unbounded_testzip(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    def _forbidden_testzip(self: zipfile.ZipFile) -> str | None:
        raise AssertionError("testzip must not be used")

    monkeypatch.setattr(zipfile.ZipFile, "testzip", _forbidden_testzip)
    data = _build_sample_pptx()
    assert extract_office_pages(data, mime_type=PPTX_MIME)


def test_office_payload_accepts_exact_drive_size() -> None:
    data = _build_sample_pptx()
    pages = extract_office_pages(
        data,
        mime_type=PPTX_MIME,
        expected_size=len(data),
        expected_md5=hashlib.md5(data, usedforsecurity=False).hexdigest(),
    )
    assert len(pages) == 2


def test_office_binary_mimes_constant_matches_individual() -> None:
    assert DOCX_MIME in OFFICE_BINARY_MIMES
    assert PPTX_MIME in OFFICE_BINARY_MIMES
    assert XLSX_MIME in OFFICE_BINARY_MIMES
    assert GDOC_NATIVE_MIME not in OFFICE_BINARY_MIMES


# -----------------------------------------------------------
# pptx: include_notes
# -----------------------------------------------------------
def _build_pptx_with_notes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    s1.shapes.title.text = "本文タイトル"
    # speaker notes
    s1.notes_slide.notes_text_frame.text = "発表者ノートの内容"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_include_notes_default_excludes_notes() -> None:
    # 既定（include_notes=False）ではノートは拾わない＝現行挙動
    pages = extract_pptx_pages(_build_pptx_with_notes())
    assert len(pages) == 1
    assert "本文タイトル" in pages[0][1]
    assert "発表者ノート" not in pages[0][1]


def test_pptx_include_notes_true_picks_up_notes() -> None:
    pages = extract_pptx_pages(_build_pptx_with_notes(), include_notes=True)
    assert len(pages) == 1
    assert "本文タイトル" in pages[0][1]
    assert "発表者ノートの内容" in pages[0][1]


def test_pptx_include_notes_no_notes_slide_is_safe() -> None:
    # ノート未設定の slide で include_notes=True でも落ちず本文のみ
    pages = extract_pptx_pages(_build_sample_pptx(), include_notes=True)
    assert len(pages) == 2
    assert "スライド1のタイトル" in pages[0][1]


# -----------------------------------------------------------
# pptx: include_tables
# -----------------------------------------------------------
def _build_pptx_with_table() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    graphic_frame = s1.shapes.add_table(
        rows=2, cols=2, left=Inches(1), top=Inches(1), width=Inches(4), height=Inches(2)
    )
    tbl = graphic_frame.table
    tbl.cell(0, 0).text = "表ヘッダA"
    tbl.cell(0, 1).text = "表ヘッダB"
    tbl.cell(1, 0).text = "値1"
    tbl.cell(1, 1).text = "値2"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_include_tables_default_excludes_table() -> None:
    # 既定（include_tables=False）では表セルを拾わない＝表だけの slide は空扱い
    pages = extract_pptx_pages(_build_pptx_with_table())
    assert pages == []


def test_pptx_include_tables_true_picks_up_cells() -> None:
    pages = extract_pptx_pages(_build_pptx_with_table(), include_tables=True)
    assert len(pages) == 1
    text = pages[0][1]
    assert "表ヘッダA" in text
    assert "表ヘッダB" in text
    assert "値1" in text
    assert "値2" in text


# -----------------------------------------------------------
# pptx: group shape 再帰
# -----------------------------------------------------------
def _build_pptx_with_group() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    shapes = s1.shapes
    # 2 つの textbox を作ってグループ化
    tb1 = shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb1.text_frame.text = "グループ内テキスト1"
    tb2 = shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    tb2.text_frame.text = "グループ内テキスト2"
    # python-pptx の group_shapes API（GroupShapes.group）
    s1.shapes.add_group_shape([tb1, tb2])
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_group_recursion_default_skips_group_children() -> None:
    # 既定（両 flag False）では group は再帰しない＝現行挙動
    pages = extract_pptx_pages(_build_pptx_with_group())
    joined = " ".join(t for _, t in pages)
    assert "グループ内テキスト1" not in joined
    assert "グループ内テキスト2" not in joined


def test_pptx_group_recursion_with_flag_picks_up_children() -> None:
    # いずれかの拡張 flag が立つと group を再帰して子テキストを拾う
    pages = extract_pptx_pages(_build_pptx_with_group(), include_tables=True)
    assert len(pages) == 1
    text = pages[0][1]
    assert "グループ内テキスト1" in text
    assert "グループ内テキスト2" in text


# -----------------------------------------------------------
# xlsx: formula_fallback
# -----------------------------------------------------------
def _build_xlsx_with_formula_only() -> bytes:
    """openpyxl で式だけ書いた book（キャッシュ値なし＝data_only で None）."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Calc"
    ws["A1"] = "=SUM(1,2)"
    ws["A2"] = "=A1*10"
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def test_xlsx_formula_fallback_default_drops_uncached_formula_sheet() -> None:
    # 既定（formula_fallback=False）= data_only のみ → キャッシュ無し式は None で空
    pages = extract_xlsx_pages(_build_xlsx_with_formula_only())
    assert pages == []


def test_xlsx_formula_fallback_true_picks_up_formula_strings() -> None:
    pages = extract_xlsx_pages(_build_xlsx_with_formula_only(), formula_fallback=True)
    assert len(pages) == 1
    assert pages[0][0] == 1
    text = pages[0][1]
    assert "=SUM(1,2)" in text
    assert "=A1*10" in text


def test_xlsx_formula_fallback_keeps_value_sheets_and_orders() -> None:
    # 値ありsheet(1) + 式のみsheet(2) → fallback で両方拾い index 昇順
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Values"
    ws1["A1"] = "実値データ"
    ws2 = wb.create_sheet("Formula")
    ws2["A1"] = "=NOW()"
    buf = BytesIO()
    wb.save(buf)

    pages = extract_xlsx_pages(buf.getvalue(), formula_fallback=True)
    nums = [n for n, _ in pages]
    assert nums == [1, 2]
    assert "実値データ" in pages[0][1]
    assert "=NOW()" in pages[1][1]


def test_xlsx_formula_fallback_does_not_change_value_only_books() -> None:
    # キャッシュ値のある book では fallback ON/OFF で結果が一致
    data = _build_sample_xlsx()
    off = extract_xlsx_pages(data)
    on = extract_xlsx_pages(data, formula_fallback=True)
    assert off == on


# -----------------------------------------------------------
# min_chars guard（dispatcher）
# -----------------------------------------------------------
def _build_pptx_tiny_and_full() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    # slide1: 極小テキスト
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    tb1 = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    tb1.text_frame.text = "あ"
    # slide2: 十分な長さ
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    tb2 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb2.text_frame.text = "これは十分に長い本文テキストです。"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_min_chars_zero_keeps_all_pages() -> None:
    # min_chars=0（既定）では極小ページも残る＝現行挙動
    pages = extract_office_pages(_build_pptx_tiny_and_full(), mime_type=PPTX_MIME)
    assert len(pages) == 2


def test_min_chars_drops_tiny_pages() -> None:
    pages = extract_office_pages(_build_pptx_tiny_and_full(), mime_type=PPTX_MIME, min_chars=5)
    # 1文字 slide は落ち、長い slide だけ残る
    assert len(pages) == 1
    assert "十分に長い本文" in pages[0][1]


def test_min_chars_applies_to_xlsx_via_dispatcher() -> None:
    # dispatcher 経由で xlsx にも min_chars が効く
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Tiny"
    ws1["A1"] = "x"
    ws2 = wb.create_sheet("Full")
    ws2["A1"] = "十分に長いセルの内容です"
    buf = BytesIO()
    wb.save(buf)

    pages = extract_office_pages(buf.getvalue(), mime_type=XLSX_MIME, min_chars=4)
    assert len(pages) == 1
    assert "十分に長いセル" in pages[0][1]


# -----------------------------------------------------------
# dispatcher: 新 kwarg の後方互換（既定で従来と一致）
# -----------------------------------------------------------
def test_dispatcher_defaults_match_legacy_pptx() -> None:
    data = _build_sample_pptx()
    assert extract_office_pages(data, mime_type=PPTX_MIME) == extract_pptx_pages(data)


def test_dispatcher_defaults_match_legacy_xlsx() -> None:
    data = _build_sample_xlsx()
    assert extract_office_pages(data, mime_type=XLSX_MIME) == extract_xlsx_pages(data)


def test_dispatcher_pptx_passes_through_notes_and_tables() -> None:
    pages = extract_office_pages(
        _build_pptx_with_notes(),
        mime_type=PPTX_MIME,
        include_notes=True,
    )
    assert "発表者ノートの内容" in pages[0][1]
